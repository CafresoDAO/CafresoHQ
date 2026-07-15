#!/usr/bin/env python3
import sys as _sys
if hasattr(_sys.stdout, 'reconfigure'):
    _sys.stdout.reconfigure(encoding='utf-8', errors='replace')
    _sys.stderr.reconfigure(encoding='utf-8', errors='replace')
"""Serve CafresoHQ static files and proxy local LM Studio / Ollama / Brave / Vault.

Same-origin proxy avoids browser CORS. SSE streams pass through unbuffered.
Brave proxy:  /brave/search?q=...   forwards to api.search.brave.com with the
caller's X-Brave-Key header rewritten to X-Subscription-Token.
Vault proxy:  /vault/*              read/write Markdown notes under a configured
root directory. Root set via CAFRESOHQ_VAULT env or POST /vault/configure.
"""
import base64
import hashlib
import http.client
import http.server
import json
import os
# ── env compat shim: mirror legacy OPENCLAW_* vars to CAFRESOHQ_* ───────────────
# Deployed container/entrypoint still export OPENCLAW_* names; mirror them so the
# renamed CAFRESOHQ_* reads keep working until images are rebuilt. Remove later.
for _k, _v in list(os.environ.items()):
    if _k.startswith('OPENCLAW_'):
        os.environ.setdefault('CAFRESOHQ_' + _k[len('OPENCLAW_'):], _v)
import pathlib
import re
import secrets
import select
import socket
import shlex
import shutil
import socketserver
import ssl
import struct
import subprocess
import sys
import threading
import time
import urllib.error
import urllib.parse
import urllib.request
import uuid

# Listen port. Env-configurable (the Dockerfile + entrypoint set PORT) so a
# self-hoster can avoid a clash with another local service; defaults to 8787.
PORT = int(os.environ.get('PORT', '8787') or '8787')
ROUTES = {
    '/lmstudio/': ('10.0.0.100', 1234),
    '/ollama/':   ('localhost', 11434),
}

# Hermes Agent (Nous Research) — the per-container `hermes gateway` exposes an
# OpenAI-compatible API server (enable via ~/.hermes/.env:
# API_SERVER_ENABLED=true, API_SERVER_KEY=…). We proxy /hermes/* → that server,
# injecting the Bearer key server-side so it never reaches the browser, and
# metering the OpenAI `usage` object for per-principal billing. Not in ROUTES
# because it needs the dedicated _hermes_proxy (auth injection + usage tap).
HERMES_HOST = os.environ.get('HERMES_API_HOST', '127.0.0.1')
HERMES_PORT = int(os.environ.get('HERMES_API_PORT', '8642') or '8642')

# ── Idle tracking (powers fleet reap-idle → stop idle containers, free A1 pool) ─
# Single-slot list so the request handler can mutate it without `global`.
# /idle, /health, and /idle's own polls do NOT count as activity.
_LAST_ACTIVITY = [time.time()]
# /market is exempt too: the ambient ticker polls every 60s and must never be
# what keeps an idle fleet container awake/billed.
_IDLE_EXEMPT_PREFIXES = ('/idle', '/health', '/market')

def _touch_activity(path):
    if not any(path == p or path.startswith(p) for p in _IDLE_EXEMPT_PREFIXES):
        _LAST_ACTIVITY[0] = time.time()

HOP_HEADERS = {'host', 'connection', 'keep-alive', 'proxy-authenticate',
               'proxy-authorization', 'te', 'trailers', 'transfer-encoding',
               'upgrade', 'content-length'}
# Claude Code (Pro/Max subscription) — invoked as a subprocess so the user's
# already-authenticated CLI does the auth. Path is overridable; if blank we
# look up `claude` in PATH at request time.
_claudecode_bin = os.environ.get('CAFRESOHQ_CLAUDE_BIN', '').strip()

# Codex CLI — alternative elevated agent backend. Uses `codex exec --json`
# with the same allowed-dirs sandbox. Auth is read from ~/.codex/config.toml
# automatically; no API key is passed by serve.py.
# Override binary path with CAFRESOHQ_CODEX_BIN env var.
_codex_bin = os.environ.get('CAFRESOHQ_CODEX_BIN', '').strip()

# Hermes CLI (Nous Research) — the container's default agent runtime. Normally
# on PATH (pip-installed in the image, or in the user's WSL/unix env). Override
# the binary path with CAFRESOHQ_HERMES_BIN. Hermes is unix-only (its gateway +
# pty_bridge import termios/pty/fcntl) so it never resolves on native Windows —
# the supported Windows path runs the whole stack inside WSL (see Start-CafresoHQ).
_hermes_bin = os.environ.get('CAFRESOHQ_HERMES_BIN', '').strip()

# Gemini CLI (Google) — npm `@google/gemini-cli`, native on every platform.
# Override the binary path with CAFRESOHQ_GEMINI_BIN.
_gemini_bin = os.environ.get('CAFRESOHQ_GEMINI_BIN', '').strip()

# Extra browser origins allowed to open the terminal WebSocket and fetch the PTY
# nonce. Needed when the HQ UI is served cross-origin from an ICP asset canister
# (the frontend/backend split). Comma-separated, e.g.
#   CAFRESOHQ_ALLOWED_WS_ORIGINS=https://<canister>.icp0.io,https://ai.cafreso.com
_extra_app_origins = {o.strip() for o in
                      os.environ.get('CAFRESOHQ_ALLOWED_WS_ORIGINS', '').split(',')
                      if o.strip()}

# API contract version between the (canister-served) UI and this backend. Bump
# only on a BREAKING change to an endpoint the UI depends on; the UI reads it
# from /health and degrades gracefully rather than hard-failing across a
# version-skewed deploy (UI on a canister, API in a slower-to-update container).
API_VERSION = 1

# ── Client path translation (Windows path → WSL mount) ───────────────────────
# The HQ UI may run on Windows and store project paths in Windows form
# (C:\Users\me\proj). The supported Windows deployment runs the whole serve.py
# stack inside WSL, where those live under /mnt/<drive>. Translate any
# client-supplied path so the unix host can resolve it. No-op on native Windows
# (serve.py there already speaks Windows paths) and for paths that are already
# unix-style. Accepts both back- and forward-slash Windows paths.
_WINPATH_RE = re.compile(r'^([A-Za-z]):[\\/](.*)$')
_WINDRIVE_RE = re.compile(r'^([A-Za-z]):[\\/]?$')

def _client_path(p):
    if not p or sys.platform == 'win32':
        return p
    m = _WINPATH_RE.match(p)
    if m:
        return f'/mnt/{m.group(1).lower()}/' + m.group(2).replace('\\', '/')
    m2 = _WINDRIVE_RE.match(p)
    if m2:
        return f'/mnt/{m2.group(1).lower()}'
    return p

# /cafresohq/stream — elevated agent endpoint with computer access. Same
# `claude` CLI as /claudecode/stream but tools are ENABLED and constrained
# by a server-side allowlist (paths and tool names). Both lists are env-
# only — the client can't widen them by sending a bigger spec. If either
# allowlist is empty the endpoint refuses requests, so the unconfigured
# default is safe.
#
#   CAFRESOHQ_ALLOWED_DIRS: os.pathsep-separated absolute paths claude can
#                          read/write. Empty → endpoint disabled.
#   CAFRESOHQ_ALLOWED_TOOLS: comma-separated tool names Claude Code accepts
#                          (Read, Write, Edit, Glob, Grep, Bash, ...).
#                          Default is read-only ("Read,Glob,Grep") so the
#                          first-time experience can't mutate the disk
#                          without an explicit opt-in.
_ALLOWED_DIRS_EXPLICIT = 'CAFRESOHQ_ALLOWED_DIRS' in os.environ
_cafresohq_allowed_dirs = [d.strip() for d in
                          os.environ.get('CAFRESOHQ_ALLOWED_DIRS',
                              os.path.expanduser('~') + os.pathsep +
                              os.path.join(os.path.expanduser('~'), 'Documents')
                          ).split(os.pathsep)
                          if d.strip()]

def _within_allowed_dirs(p):
    """True iff resolved path `p` is inside one of _cafresohq_allowed_dirs.
    Uses Path.relative_to (NOT str.startswith, which lets '/data/proj' authorize
    the sibling '/data/proj-secret'). Enforced in EVERY runtime mode — the fs
    read/browse routes are otherwise an unauthenticated arbitrary-read hole in
    local/BYO deployments. If the allow-list is explicitly emptied, allow (the
    documented opt-out)."""
    if not _cafresohq_allowed_dirs:
        return True
    try:
        rp = pathlib.Path(p).resolve()
    except Exception:
        return False
    for d in _cafresohq_allowed_dirs:
        try:
            rp.relative_to(pathlib.Path(d).resolve())
            return True
        except (ValueError, Exception):
            continue
    return False
_cafresohq_allowed_tools = [t.strip() for t in
                           os.environ.get('CAFRESOHQ_ALLOWED_TOOLS',
                               # Default expanded set — gives elevated agents
                               # a real workshop instead of read-only browsing.
                               # Bash + Edit + Write enable code edits; Glob +
                               # Grep cover discovery; WebFetch + WebSearch
                               # cover external research; the ROOT TodoWrite
                               # / Notebook* tools are deliberately excluded
                               # because they're meta-tools that confuse the
                               # sub-agent's own task list.
                               'Read,Glob,Grep,Bash,Edit,Write,WebFetch,WebSearch'
                           ).split(',')
                           if t.strip()]

# ── Runtime environment detection ────────────────────────────────────────────
# 'container' → OCI/Docker container (no display; terminal spawn disabled)
# 'local'     → developer machine (all features enabled)
def _detect_runtime():
    if sys.platform == 'win32':
        return 'local'
    # Linux/macOS: check standard container markers
    if os.path.exists('/.dockerenv') or os.path.exists('/run/.containerenv'):
        return 'container'
    try:
        with open('/proc/1/cgroup', 'r') as _f:
            _cg = _f.read()
        if any(kw in _cg for kw in ('docker', 'kubepods', 'lxc', 'containerd')):
            return 'container'
    except Exception:
        pass
    # Our OCI fleet env vars are definitive
    if os.environ.get('OCI_TENANCY_ID') or os.environ.get('FLEET_REGISTRY'):
        return 'container'
    return 'local'

_RUNTIME_ENV = _detect_runtime()

# ── Bearer-key auth for local / BYO deployments ──────────────────────────────
# When CAFRESOHQ_API_KEY is set, the dangerous routes (terminal, agents, vault,
# code-agent streams, hq-state) require it via an X-API-Key header or a ?k=
# query param (WebSocket handshakes can't set custom headers). In OCI-fleet mode
# the Caddy gateway + verifier already gate access, so the key stays unset/optional
# there — this is the security floor for Local-native/WSL and any non-loopback
# exposure (the terminal PTY is effectively RCE). Static UI + /health stay open
# so the app shell can bootstrap.
CAFRESOHQ_API_KEY = os.environ.get('CAFRESOHQ_API_KEY', '').strip()
_KEY_PROTECTED_PREFIXES = (
    '/vault', '/hermes', '/terminal', '/cafresohq', '/codex', '/claudecode',
    '/agents', '/hq/', '/hq-state', '/spawn', '/graph/publish', '/approvals',
    '/browser', '/missions',
    # Code-execution + filesystem routes: RCE-/write-equivalent to the routes
    # above and MUST sit behind the same key. /tools/exec runs a shell;
    # /projects,/export,/generate touch the host; the /fs mutation routes
    # write/delete. NOTE: only the /fs *mutation* prefixes are protected — the
    # preview iframe fetches /fs/site/<root>/<asset> read-only via the browser
    # with no key, so blanket-protecting bare '/fs' would break multi-file
    # previews. The allowed-dirs boundary (enforced in every mode below) caps
    # the read routes instead.
    '/tools', '/projects', '/export', '/generate',
    '/fs/upload', '/fs/mkdir', '/fs/rename', '/fs/delete',
)

# Background CLI-install jobs (POST /agents/install returns 202 immediately;
# the UI polls GET /agents/install/status?agent=…). One job per agent id.
_INSTALL_JOBS = {}
_INSTALL_JOBS_LOCK = threading.Lock()

# ── Market quotes (Trading Floor theme ticker) ──────────────────────────────
# GET /market/quotes proxies Yahoo Finance's public chart endpoint: stock
# indices have no CORS-open free API the browser could hit directly, and
# stooq's CSV endpoints now sit behind a JS anti-bot wall. Cached 60s so a
# whole floor of open tabs costs one upstream sweep per minute; on total
# upstream failure the last good payload is served stale instead of erroring.
_MARKET_SYMBOLS = (
    ('NAS100', '^NDX'),
    ('US30',   '^DJI'),
    ('SPX',    '^GSPC'),
    ('GOLD',   'GC=F'),
)
_market_cache = {'ts': 0.0, 'quotes': []}
_market_lock = threading.Lock()

# Per-process nonce for /terminal/pty WebSocket auth.
# Generated once at startup; never logged.  Frontend fetches it from
# /terminal/nonce before opening the WebSocket.  Any connection that doesn't
# supply it with secrets.compare_digest() equality is rejected with 403.
_PTY_NONCE = secrets.token_hex(32)   # 256-bit random — unguessable

# ── Persistent PTY session registry ────────────────────────────────────────
# Keeps PTY processes alive after the WebSocket drops so mobile clients can
# reconnect and see buffered output from work that ran while they were away.
_PTY_SESSIONS    = {}               # session_id → session-state dict
_PTY_SESSIONS_LK = threading.Lock()
_PTY_SESSION_TTL = 300              # seconds to keep a disconnected PTY alive
_PTY_BUF_CAP     = 524_288          # max bytes buffered per session (512 KB)

def _pty_reaper():
    while True:
        time.sleep(30)
        now = time.time()
        to_kill = []
        with _PTY_SESSIONS_LK:
            for sid, s in list(_PTY_SESSIONS.items()):
                if s['sock'] is None and s.get('expires', 0) < now:
                    to_kill.append((sid, s))
            for sid, _ in to_kill:
                _PTY_SESSIONS.pop(sid, None)
        for _, s in to_kill:
            s['stop'].set()
            for k in ('pty_proc', 'proc'):
                try: s.get(k) and s[k].terminate()
                except Exception: pass

threading.Thread(target=_pty_reaper, daemon=True).start()

# HQ state persistence
_hq_state_dir   = pathlib.Path(os.environ.get('CAFRESOHQ_HQ_STATE_DIR',
                    os.path.join(os.path.dirname(__file__), 'hq-state')))
_hq_memory_dir  = pathlib.Path(os.environ.get('CAFRESOHQ_MEMORY_DIR',
                    os.path.join(os.path.dirname(__file__), 'hq-state', 'memory')))


# ---- Night Shift (Sprint 4 MVP-1) ------------------------------------------
# "Close the laptop, work continues": a 30s scheduler daemon runs missions
# SERVER-side via night_runner.py — a restricted read/search/fetch + vault-write
# tool subset (no WALLET/PUBLISH/BASH/FILE_WRITE; the money/publish seam stays
# in the authenticated browser shell by construction). Files:
#   hq-state/scheduled-missions.json  — the schedules (UI + hqsh manage these)
#   hq-state/mission-runs.json        — ring-capped run log (Gazette lead story)
# One mission runs at a time; shared-activity writes are deferred while a real
# browser session is live (last-writer-wins clobber guard).
_night_lock = threading.Lock()
_night_running = {}          # scheduleId -> True while a run is in flight
_night_base_url = ['']       # set in __main__ once scheme + port are known
_LAST_UI_ACTIVITY = [0.0]    # last request that came from a real browser
MAX_NIGHT_RUNS_KEPT = 100
MAX_NIGHT_SCHEDULES = 20


def _night_path(name):
    return _hq_state_dir / name


def _night_load(name, default):
    try:
        with open(_night_path(name), 'r', encoding='utf-8') as f:
            data = json.load(f)
        return data if isinstance(data, type(default)) else default
    except Exception:
        return default


def _night_save(name, data):
    _hq_state_dir.mkdir(parents=True, exist_ok=True)
    tmp = _night_path(name + '.tmp')
    with open(tmp, 'w', encoding='utf-8') as f:
        json.dump(data, f)
    os.replace(tmp, _night_path(name))


# ── Shared "trial brain" metering ──────────────────────────────────────────
# hermes-bootstrap.py writes HERMES_HOME/trial.json when a fresh HQ falls back to
# the operator's shared trial key (so a new user's agents work with no signup).
# That key draws on ONE shared upstream account, so serve.py caps each principal
# to a small number of trial completions/day and clears the trial the moment the
# user brings their own key. A user's OWN key is never metered here.
_TRIAL_DAILY_CAP = max(1, int(os.environ.get('CAFRESOHQ_TRIAL_DAILY_CAP', '25') or '25'))
_trial_usage_lock = threading.Lock()

# ── Operator control plane (network-wide switches, read from the canister) ──
# The planAdmin sets ONE JSON blob on cafresohq_state; every client + container
# reads it from the public /operator/config.json route. serve.py caches it and
# honors: trialBrain.enabled (kill the shared trial network-wide),
# trialBrain.dailyCap (live cap override), gpuNode.enabled (a worker self-pauses
# when the operator takes the node down). Never blocks on the network — a failed
# fetch keeps the last good config (or {}), so a canister blip can't break chat.
_OPERATOR_BASE = os.environ.get(
    'SEARCH_STATE_URL', 'https://ydacz-riaaa-aaaal-qxeja-cai.icp0.io').rstrip('/')
_operator_cfg = {'ts': 0.0, 'data': {}}
_operator_cfg_lock = threading.Lock()


def _operator_config():
    with _operator_cfg_lock:
        if time.time() - _operator_cfg['ts'] < 60:
            return _operator_cfg['data']
    data = None
    try:
        with urllib.request.urlopen(_OPERATOR_BASE + '/operator/config.json', timeout=6) as r:
            parsed = json.loads(r.read().decode('utf-8'))
        if isinstance(parsed, dict):
            data = parsed
    except Exception:
        pass
    with _operator_cfg_lock:
        if data is not None:
            _operator_cfg['data'] = data
        _operator_cfg['ts'] = time.time()
        return _operator_cfg['data']


def _trial_cap():
    """Operator override (trialBrain.dailyCap) wins over the env default."""
    op = _operator_config().get('trialBrain') or {}
    c = op.get('dailyCap')
    if isinstance(c, (int, float)) and c >= 1:
        return int(c)
    return _TRIAL_DAILY_CAP


def _hermes_home():
    return os.environ.get('HERMES_HOME', '').strip() or os.path.expanduser('~/.hermes')


def _validate_local_base_url(u):
    """(ok, error) for a LOCAL model backend URL the operator typed.

    ALLOWLIST, not blocklist — whatever this accepts, the container will later
    POST to, and /hermes is only key-protected when CAFRESOHQ_API_KEY is set. So
    this is the whole defence, not a layer of it.

    IP-literal/localhost only is deliberate and is what closes DNS rebinding: we
    validate at write time but hermes resolves at call time, so a hostname that
    passes here can point anywhere seconds later. An IP literal has no such gap.
    Operators with a real internal DNS name can set HQ_ALLOW_REMOTE_BACKEND=1
    and accept that risk knowingly.
    """
    import ipaddress
    u = (u or '').strip()
    if not u:
        return False, 'base_url is required'
    if len(u) > 300:
        return False, 'base_url too long'
    try:
        p = urllib.parse.urlsplit(u)
    except Exception:
        return False, 'unparseable base_url'
    if p.scheme not in ('http', 'https'):
        return False, 'base_url must be http or https'
    if p.username or p.password:
        return False, 'base_url must not contain credentials'
    if p.query or p.fragment:
        return False, 'base_url must not contain a query or fragment'
    if '..' in p.path:
        return False, 'base_url path must not contain ".."'
    if not re.fullmatch(r'(/[\w\-.]+)*/?', p.path or ''):
        return False, 'base_url has an unexpected path'
    host = p.hostname
    if not host:
        return False, 'base_url has no host'
    if os.environ.get('HQ_ALLOW_REMOTE_BACKEND', '').strip() == '1':
        return True, ''
    try:
        infos = socket.getaddrinfo(host, p.port or (443 if p.scheme == 'https' else 80),
                                   proto=socket.IPPROTO_TCP)
    except Exception:
        return False, 'could not resolve %s' % host
    addrs = {i[4][0] for i in infos}
    if not addrs:
        return False, 'could not resolve %s' % host
    for a in addrs:
        try:
            ip = ipaddress.ip_address(a)
        except ValueError:
            return False, 'unexpected address for %s' % host
        # link-local (169.254/16, fe80::/10) covers the cloud metadata endpoint,
        # which some definitions call "private" — it must never be reachable.
        if ip.is_link_local:
            return False, 'link-local addresses are not allowed'
        if not (ip.is_private or ip.is_loopback):
            return False, ('%s resolves to the public address %s — a local backend must be '
                           'on private or loopback space (set HQ_ALLOW_REMOTE_BACKEND=1 to '
                           'override)' % (host, a))
    return True, ''


def _trial_state():
    """{'active': bool, 'provider': str} — read fresh so a BYOK clear takes effect
    without a restart. The operator can also kill the shared trial network-wide
    (trialBrain.enabled=false) — that override wins over the local marker."""
    try:
        with open(os.path.join(_hermes_home(), 'trial.json'), 'r', encoding='utf-8') as f:
            d = json.load(f)
        active = bool(d.get('active'))
        provider = str(d.get('provider') or '')
    except Exception:
        active, provider = False, ''
    if active:
        op = _operator_config().get('trialBrain') or {}
        if op.get('enabled') is False:
            active = False   # operator kill switch
    return {'active': active, 'provider': provider}


def _trial_deactivate():
    """Called when the user brings their own key — the trial (and its cap) end."""
    try:
        with open(os.path.join(_hermes_home(), 'trial.json'), 'w', encoding='utf-8') as f:
            json.dump({'active': False, 'provider': ''}, f)
    except Exception:
        pass


def _trial_usage(principal):
    """Return (used_today, cap) for a principal without mutating."""
    cap = _trial_cap()
    today = time.strftime('%Y-%m-%d')
    u = _night_load('trial-usage.json', {})
    if u.get('date') != today:
        return 0, cap
    return int((u.get('counts') or {}).get(principal, 0)), cap


def _trial_check_and_bump(principal):
    """Atomically count one trial completion. Returns (allowed, used, cap):
    allowed False (over cap) means DON'T forward — surface the upsell."""
    cap = _trial_cap()
    today = time.strftime('%Y-%m-%d')
    with _trial_usage_lock:
        u = _night_load('trial-usage.json', {})
        if u.get('date') != today:
            u = {'date': today, 'counts': {}}
        counts = u.setdefault('counts', {})
        used = int(counts.get(principal, 0))
        if used >= cap:
            return False, used, cap
        counts[principal] = used + 1
        _night_save('trial-usage.json', u)
        return True, used + 1, cap


def _night_log_run(run):
    """Upsert a run record (called progressively so a crash keeps partial log)."""
    with _night_lock:
        runs = [r for r in _night_load('mission-runs.json', [])
                if r.get('id') != run.get('id')]
        runs.append(run)
        _night_save('mission-runs.json', runs[-MAX_NIGHT_RUNS_KEPT:])


def _night_browser_active():
    return (time.time() - _LAST_UI_ACTIVITY[0]) < 180


def _night_ctx():
    import night_runner as _nr
    return _nr.NightContext(
        _night_base_url[0] or ('http://127.0.0.1:%d' % PORT),
        api_key=CAFRESOHQ_API_KEY)


def _night_post_activity(run):
    """Surface a finished run on the shared activity ticker. Deferred while a
    browser is live — the app's merge-by-id load makes concurrent writes mostly
    safe, but last-writer-wins can still clobber, so we only write when nobody
    is watching. The Gazette ingests mission-runs.json regardless."""
    if _night_browser_active():
        return
    try:
        import night_runner as _nr
        ctx = _night_ctx()
        s, raw = _nr._self_call(ctx, 'GET', '/hq/state/activity')
        cur = json.loads(raw.decode('utf-8', 'replace')) if s == 200 else []
        if not isinstance(cur, list):
            cur = []
        entry = {
            'id': 'act_night_%s' % run.get('id', ''),
            'ts': int(time.time() * 1000),
            'agentId': run.get('agentId', ''), 'agentName': run.get('agentName', ''),
            'action': 'night', 'priority': 'attention', 'unread': True,
            'text': 'night shift: %d note(s) on %s' % (
                len(run.get('writes', [])), (run.get('topic') or '')[:60]),
        }
        _nr._self_call(ctx, 'PUT', '/hq/state/activity',
                       body=json.dumps([entry] + cur[:199]))
    except Exception:
        pass


def _night_run_one(sched):
    try:
        import night_runner as _nr
        run = _nr.run_mission(_night_ctx(), sched, on_progress=_night_log_run)
        _night_log_run(run)
        _night_post_activity(run)
    except Exception as e:
        now_ms = int(time.time() * 1000)
        _night_log_run({
            'id': 'run_%d' % now_ms, 'scheduleId': sched.get('id', ''),
            'agentId': sched.get('agentId', ''), 'agentName': sched.get('agentName', ''),
            'topic': sched.get('topic', ''), 'vaultFolder': sched.get('vaultFolder', ''),
            'startedAt': now_ms, 'finishedAt': now_ms, 'iterations': 0, 'writes': [],
            'tokensUsed': 0, 'errors': 1, 'lastError': str(e)[:300], 'summary': ''})
    finally:
        _night_running.pop(sched.get('id', ''), None)


def _night_scan():
    now_ms = int(time.time() * 1000)
    with _night_lock:
        scheds = _night_load('scheduled-missions.json', [])
        changed = False
        for s in scheds:
            if not s.get('enabled'):
                continue
            if int(s.get('nextRunAt', 0) or 0) > now_ms:
                continue
            if _night_running:
                break   # one mission at a time — keeps provider usage sane
            sid = str(s.get('id', ''))
            _night_running[sid] = True
            s['lastRunAt'] = now_ms
            if s.get('recurrence') == 'daily':
                nxt = int(s.get('nextRunAt', now_ms) or now_ms)
                while nxt <= now_ms:
                    nxt += 86_400_000
                s['nextRunAt'] = nxt
            else:
                s['enabled'] = False
            changed = True
            threading.Thread(target=_night_run_one, args=(dict(s),),
                             daemon=True, name='night-%s' % sid).start()
        if changed:
            _night_save('scheduled-missions.json', scheds)


def _night_loop():
    while True:
        time.sleep(30)
        try:
            _night_scan()
        except Exception as e:
            print('[night] scan error:', e)


threading.Thread(target=_night_loop, daemon=True, name='night-shift').start()


# ---- Ai Cafreso Search — network worker ------------------------------------
# Opt-in: this container joins the on-chain search network, claiming queued
# queries from cafresohq_state over HMAC-signed plain HTTPS (Python has no IC
# agent — the canister verifies HMAC-SHA256 over the raw body against the
# secret registered from the browser). Pipeline per job: Brave (this
# container's own key) → local LLM (hermes gateway first, night-runner
# provider fallback) → graph snapshot → fulfill. Fulfilled answers become
# public on-chain library entries attributed (and paid) to WORKER_PRINCIPAL.
#
#   SEARCH_WORKER=1
#   WORKER_PRINCIPAL=<principal registered via ai.cafreso.com settings>
#   WORKER_SECRET=<64-hex secret shown once at registration>
#   BRAVE_API_KEY=<your free Brave key>
#   SEARCH_STATE_URL=<override for local-replica testing; defaults to mainnet>
#   WORKER_JOB_BUDGET=<seconds for the whole claim→fulfill window; default 200>
#   WORKER_IDLE_TIMEOUT=<seconds of LLM silence before giving up; default 25>
#   WORKER_MAX_TOKENS=<completion cap for the answer; default 700>
#   WORKER_MODEL=<OPTIONAL override; by default search follows the operator's
#                 HQ brain picker — see _sw_model()>
#
# The model comes from the brain picker by default, deliberately: an operator
# who switches their brain in HQ expects search to follow, and a stale
# WORKER_MODEL pin pointing at a model the backend no longer has loaded fails
# EVERY job (silently, sources-only) with no hint as to why.
#
# The loop is SINGLE-THREADED by design: the canister's replay guard requires
# each worker's signed timestamps to strictly increase, which serialized calls
# guarantee for free. Workers run standalone — this loop makes only outbound
# calls and never touches the idle tracker, so it will not keep a fleet
# container awake by itself; a worker box must have idle-stop disabled.
#
# EVERY job runs against one wall-clock deadline. The canister's claim lease is
# 240s (CLAIM_LEASE_NS) and is NOT renewable — heartbeat does not touch
# claimedAt and there is no renew op. Overrunning it is silent and lossy: the
# fulfill comes back not-your-claim, the work is discarded, attempts increments,
# and three strikes fail the job for good. So Brave + LLM + graph + upload all
# have to fit inside _SW_JOB_BUDGET, which sits 40s inside the lease.
_SW_ENABLED = os.environ.get('SEARCH_WORKER', '').strip() == '1'
_SW_PRINCIPAL = os.environ.get('WORKER_PRINCIPAL', '').strip()
_SW_SECRET_HEX = os.environ.get('WORKER_SECRET', '').strip()
_SW_BASE = os.environ.get(
    'SEARCH_STATE_URL', 'https://ydacz-riaaa-aaaal-qxeja-cai.icp0.io').rstrip('/')
_SW_JOB_BUDGET = float(os.environ.get('WORKER_JOB_BUDGET', '') or 200)
_SW_TTFT_TIMEOUT = 60.0    # first token: prefill on a saturated box is slow
_SW_IDLE_TIMEOUT = float(os.environ.get('WORKER_IDLE_TIMEOUT', '') or 25)
_SW_MAX_TOKENS = int(os.environ.get('WORKER_MAX_TOKENS', '') or 700)
_SW_ANSWER_CAP = 4000      # LIB_MAX_ANSWER — the canister rejects anything longer
_sw_last_ts = 0


def _sw_left(deadline, floor=1.0):
    """Seconds until `deadline`, floored. The floor is load-bearing: urlopen
    reads timeout=0 as 'block forever', so passing a non-positive remaining
    straight through would turn a blown deadline into a hang."""
    if deadline is None:
        return None
    return max(floor, deadline - time.monotonic())




def _sw_call(op, extra_lines):
    """Signed worker call. Returns (status, dict). Envelope + HMAC per docs/SEARCH_NETWORK.md."""
    global _sw_last_ts
    import hmac as _hmac
    ts = max(int(time.time() * 1000), _sw_last_ts + 1)
    _sw_last_ts = ts
    body = '\n'.join(['v1', _SW_PRINCIPAL, str(ts), secrets.token_hex(8), op]
                     + list(extra_lines)).encode('utf-8')
    sig = _hmac.new(bytes.fromhex(_SW_SECRET_HEX), body, hashlib.sha256).hexdigest()
    req = urllib.request.Request(_SW_BASE + '/worker/' + op, data=body, headers={
        'x-worker-signature': sig, 'Content-Type': 'text/plain'})
    try:
        with urllib.request.urlopen(req, timeout=30) as r:
            return r.status, json.loads(r.read().decode('utf-8'))
    except urllib.error.HTTPError as e:
        try:
            return e.code, json.loads(e.read().decode('utf-8'))
        except Exception:
            return e.code, {}
    except Exception as e:
        return 0, {'error': str(e)[:200]}


def _sw_brave(q, deadline=None):
    """This worker's own Brave key. Descriptions feed the LLM prompt ONLY —
    stored sources are title+url (deliberate Brave-ToS posture)."""
    key = os.environ.get('BRAVE_API_KEY', '').strip()
    if not key:
        return None
    url = ('https://api.search.brave.com/res/v1/web/search?q='
           + urllib.parse.quote(q) + '&count=8')
    req = urllib.request.Request(url, headers={
        'Accept': 'application/json', 'Accept-Encoding': 'identity',
        'X-Subscription-Token': key, 'User-Agent': 'CafresoHQ/1.0'})
    timeout = 30 if deadline is None else min(30, _sw_left(deadline))
    with urllib.request.urlopen(req, timeout=timeout) as r:
        data = json.loads(r.read().decode('utf-8'))
    out = []
    for x in (data.get('web') or {}).get('results', [])[:8]:
        if x.get('url'):
            out.append({'title': (x.get('title') or x['url'])[:600],
                        'url': x['url'][:600],
                        'description': re.sub(r'<[^>]+>', '', x.get('description') or '')})
    return out


def _sw_hermes_key():
    """The gateway key. MUST resolve via _hermes_home(): containers set
    HERMES_HOME=/data/hermes and have no ~/.hermes at all, so hardcoding the
    home here silently returned '' — which made _sw_llm skip the gateway
    entirely and fulfill every job sources-only, with any model."""
    key = os.environ.get('API_SERVER_KEY', '').strip()
    if key:
        return key
    try:
        with open(os.path.join(_hermes_home(), '.env'), 'r', encoding='utf-8') as f:
            m = re.search(r'(?m)^API_SERVER_KEY\s*=\s*(\S+)', f.read())
        if m:
            return m.group(1).strip().strip('"\'')
    except Exception:
        pass
    return ''


def _sw_model():
    """Which model answers searches.

    Follows the operator's HQ brain picker by default: the picker writes
    `model.default` into the hermes config, and reading it here means flipping
    the picker moves the search worker too — resolved per job, so no restart.
    WORKER_MODEL stays as an explicit override for operators who deliberately
    want search on a different (e.g. cheaper/faster) model than their agents.
    """
    hint = os.environ.get('WORKER_MODEL', '').strip()
    if hint:
        return hint
    import night_runner as _nr
    return _nr.read_model_config(_hermes_home())['model'] or 'default'


def _sw_backend():
    """The operator's brain, callable directly — or None when only the agent
    gateway can reach it. The single seam the tests monkeypatch."""
    import night_runner as _nr
    return _nr.resolve_backend(_hermes_home(),
                               model_override=os.environ.get('WORKER_MODEL', '').strip())


def _sw_trial_notice():
    """Say it out loud when search is about to spend the SHARED trial key.

    Deliberately a log line, not a meter: search is uncapped by decision. The
    backstops are the canister's 500/day answer budget and the fact that becoming
    a worker requires manual planAdmin approval. _trial_usage is non-mutating, so
    reading it here doesn't consume anyone's chat allowance."""
    try:
        st = _trial_state()
        if not st.get('active'):
            return
        used, cap = _trial_usage(_SW_PRINCIPAL or 'local')
        print('[search-worker] spending the SHARED TRIAL key (provider=%s) — search is NOT '
              'metered by the trial cap; %s is at %d/%d today'
              % (st.get('provider') or '?', (_SW_PRINCIPAL or 'local')[:12], used, cap))
    except Exception:
        pass


_SW_ESCAPES = {'"': '"', '\\': '\\', '/': '/', 'b': '\b',
               'f': '\f', 'n': '\n', 'r': '\r', 't': '\t'}
_SW_MIN_SALVAGE_WORDS = 3  # a salvaged sentence shorter than this isn't worth publishing


def _sw_json_str_at(t, i):
    """Scan the JSON string literal starting at t[i] == '"'.
    Returns (value, end_index, complete). complete=False means the text ran out
    mid-literal (a truncated stream) — value is then everything decoded so far,
    which is exactly what we want to keep."""
    if i >= len(t) or t[i] != '"':
        return '', i, False
    out = []
    i += 1
    while i < len(t):
        c = t[i]
        if c == '"':
            return ''.join(out), i + 1, True
        if c != '\\':
            out.append(c)
            i += 1
            continue
        # Escape: needs at least one more char, \u needs four more.
        if i + 1 >= len(t):
            break
        e = t[i + 1]
        if e == 'u':
            if i + 5 >= len(t):
                break                       # cut inside \uXXXX — keep what we have
            try:
                out.append(chr(int(t[i + 2:i + 6], 16)))
            except ValueError:
                out.append(t[i + 2:i + 6])  # malformed \u — pass through, never raise
            i += 6
        else:
            out.append(_SW_ESCAPES.get(e, e))
            i += 2
    return ''.join(out), len(t), False


def _sw_salvage_json(t):
    """Best-effort (summary, notes) from possibly-truncated JSON.

    Targeted scanning, NOT general repair — brace-balancing mis-repairs
    silently. This works because the prompt pins "summary" FIRST: a stream cut
    anywhere inside "notes" still yields a complete summary, which is the
    artifact that actually matters. Returns (None, None) when there's no
    "summary" key at all, so the caller can fall through to plain text."""
    ks = t.find('"summary"')
    if ks < 0:
        return None, None
    vs = t.find('"', ks + len('"summary"'))
    if vs < 0:
        return None, None
    summary, end, complete = _sw_json_str_at(t, vs)
    summary = summary.strip()
    if not complete:
        # Library entries are permanent and public, so only publish a truncated
        # summary if a WHOLE sentence survived: trim back to the last boundary,
        # and if there wasn't one, a half-thought like "ICP i" is worse than no
        # answer — degrade to source-only and let the entry be honestly
        # summary-less. Applies to truncated text ONLY; a model that
        # deliberately returns one terse complete sentence keeps it.
        cut = max(summary.rfind('.'), summary.rfind('!'), summary.rfind('?'))
        if cut <= 0:
            return None, None
        summary = summary[:cut + 1]
        if len(summary.split()) < _SW_MIN_SALVAGE_WORDS:
            return None, None
        return summary, {}
    notes = {}
    kn = t.find('"notes"', end)
    if kn < 0:
        return (summary or None), notes
    i = t.find('{', kn)
    if i < 0:
        return (summary or None), notes
    i += 1
    while i < len(t):
        ks = t.find('"', i)
        if ks < 0:
            break
        # A '}' before the next key means the notes object closed — stop here
        # rather than reading keys from whatever follows it.
        close = t.find('}', i)
        if close >= 0 and close < ks:
            break
        key, i, ok = _sw_json_str_at(t, ks)
        if not ok:
            break                            # truncated key — nothing usable past here
        vs = t.find('"', i)
        if vs < 0:
            break
        val, i, ok = _sw_json_str_at(t, vs)
        if not ok:
            break                            # truncated value — drop this partial note
        try:
            notes[int(key) - 1] = val.strip()[:200]
        except (ValueError, TypeError):
            pass
    return (summary or None), notes


def _sw_parse_analysis(text):
    """Parse the model's {"summary", "notes"} JSON; degrade to plain text.
    Returns (summary, notes) where notes maps 0-based source index → blurb.

    Three tiers: strict parse (well-formed output) → salvage (stream truncated
    by the deadline) → raw text."""
    t = (text or '').strip()
    t = re.sub(r'^```(?:json)?\s*|\s*```$', '', t)
    m = re.search(r'\{.*\}', t, re.S)
    if m:
        try:
            d = json.loads(m.group(0))
            summary = str(d.get('summary') or '').strip()
            notes = {}
            for k, v in (d.get('notes') or {}).items():
                try:
                    notes[int(k) - 1] = str(v).strip()[:200]
                except (ValueError, TypeError):
                    pass
            if summary:
                return summary, notes
        except Exception:
            pass
    summary, notes = _sw_salvage_json(t)
    if summary:
        return summary, notes or {}
    if t.startswith('{') or '"summary"' in t:
        # It was trying to be our JSON and nothing usable survived. Returning the
        # raw text here would publish a broken `{"summary": "ICP i` fragment as
        # the permanent public answer — an empty answer degrades to source-only,
        # which is honest.
        return '', {}
    return t, {}


def _sw_fmt_tokens(n):
    return '%.1fk' % (n / 1000.0) if n >= 1000 else str(n)


def _sw_settimeout(resp, secs):
    """Re-bound an open response's socket. Reaches through private API
    (fp.raw._sock), so it's best-effort by design: on failure the socket keeps
    the coarser timeout from urlopen, which still bounds us — just less tightly.
    Never raises."""
    try:
        fp = getattr(resp, 'fp', None)
        s = getattr(getattr(fp, 'raw', None), '_sock', None) or getattr(fp, '_sock', None)
        if s is not None:
            s.settimeout(secs)
            return True
    except Exception:
        pass
    return False


def _sw_chat(url, headers, payload, deadline, stream=True, usage_opt=True):
    """One OpenAI-compatible chat call → (text, model, tokens, truncated).

    Streams by default so a slow backend yields a PARTIAL answer at the deadline
    instead of nothing. Raises on hard failure (connection refused, HTTP error,
    zero content) so the caller can try the next backend; returns whatever
    arrived when the wall clock or the idle timer fires."""
    body = dict(payload)
    if stream:
        body['stream'] = True
        if usage_opt:
            # usage is omitted from streamed responses unless asked for explicitly.
            body['stream_options'] = {'include_usage': True}
    req = urllib.request.Request(url, data=json.dumps(body).encode('utf-8'), headers=headers)
    # Open with the WHOLE remaining budget, not the TTFT bound. A blocking
    # gateway withholds headers until generation finishes, so a TTFT-bounded open
    # would kill a healthy 100s blocking generation with 140s still on the clock.
    # Streaming costs nothing here — SSE headers land immediately — and its
    # tighter TTFT/idle bounds go on the socket below, where they belong. If the
    # socket isn't reachable those bounds degrade to this one: coarser, but the
    # deadline still holds and we still salvage.
    r = urllib.request.urlopen(req, timeout=_sw_left(deadline) or (_SW_TTFT_TIMEOUT * 3))
    try:
        # Backward-compat gate: a gateway that ignored `stream` just sends JSON.
        # Read it as a blocking response — no exception, no retry, no wasted decode.
        if 'text/event-stream' not in (r.headers.get('Content-Type') or ''):
            data = json.loads(r.read().decode('utf-8'))
            text = data['choices'][0]['message']['content']
            tokens = int((data.get('usage') or {}).get('total_tokens') or 0)
            return text, data.get('model') or '', tokens, False

        # Streaming: bound the FIRST token by TTFT (prefill on a saturated box is
        # slow but not unbounded), then tighten to the idle timeout once tokens
        # are flowing. Best-effort — see _sw_sock.
        _sw_settimeout(r, _SW_TTFT_TIMEOUT)
        chunks, model, tokens, deltas, truncated = [], '', 0, 0, False
        while True:
            if deadline is not None and time.monotonic() > deadline:
                truncated = True
                break
            try:
                line = r.readline()
            except (socket.timeout, TimeoutError):
                truncated = True          # backend went quiet — keep what we have
                break
            if not line:
                break
            line = line.decode('utf-8', 'replace').strip()
            if not line.startswith('data:'):
                continue
            data = line[5:].strip()
            if data == '[DONE]':
                break
            try:
                obj = json.loads(data)
            except ValueError:
                continue
            model = obj.get('model') or model
            usage = obj.get('usage')
            if usage:
                tokens = int(usage.get('total_tokens') or 0)
            for ch in (obj.get('choices') or []):
                piece = (ch.get('delta') or {}).get('content')
                if piece:
                    chunks.append(piece)
                    deltas += 1
                    if deltas == 1:
                        # First token is in: tighten from the prefill bound to
                        # the idle bound so a mid-stream stall aborts fast.
                        _sw_settimeout(r, _SW_IDLE_TIMEOUT)
        text = ''.join(chunks)
        if not text:
            raise RuntimeError('empty stream')
        # A truncated stream never carries the final usage chunk, so fall back to
        # the delta count. That's completion-only (excludes prefill) and so
        # UNDER-reports — deliberate: the chip says "~", and guessing a prompt
        # estimate to close the gap would be inventing numbers.
        return text, model, (tokens or deltas), truncated
    finally:
        try:
            r.close()          # drop a slow socket rather than leak it into the next job
        except Exception:
            pass


# (stream, usage_opt, schema) — richest first. response_format is dropped LAST
# because it's the rung that carries real value: it's what makes a note per
# source structurally required rather than a polite request.
_SW_ATTEMPTS = ((True, True, True), (True, False, True),
                (True, False, False), (False, False, False))
_SW_CAPS = {}          # url → index of the last rung that worked


def _sw_stream_or_block(url, headers, payload, deadline):
    """_sw_chat with a graceful climb-down for backends that reject a param.

    A 4xx means the backend dislikes something we sent, so retry smaller. A 5xx
    means the backend itself is unhappy and re-sending the same work won't help
    — raise so the caller moves on.

    The per-url memo is what makes a 4-rung ladder affordable: a backend that
    hates response_format costs ONE wasted round-trip per process, not one per
    job. A 4xx carries no decode, so even uncached the cost is milliseconds."""
    start = _SW_CAPS.get(url, 0)
    for i in range(start, len(_SW_ATTEMPTS)):
        stream, usage_opt, schema = _SW_ATTEMPTS[i]
        body = dict(payload)
        if not schema:
            body.pop('response_format', None)
        try:
            out = _sw_chat(url, headers, body, deadline, stream=stream, usage_opt=usage_opt)
            _SW_CAPS[url] = i
            return out
        except urllib.error.HTTPError as e:
            if not (400 <= e.code < 500) or i == len(_SW_ATTEMPTS) - 1:
                raise      # 5xx, or the last rung 4xx'd — genuinely broken
            dropped = ('stream_options' if usage_opt else
                       'response_format' if schema else 'stream')
            print('[search-worker] backend rejected %s (%s) — climbing down' % (dropped, e.code))
    raise RuntimeError('unreachable')


def _sw_schema(n):
    """Strict json_schema for the {summary, notes} answer.

    Two deliberate choices:
    - `summary` is FIRST in properties. Constrained decoders emit in schema
      order, and _sw_salvage_json depends on summary-before-notes to rescue a
      deadline-truncated answer.
    - notes enumerates keys "1".."N" with `required` rather than a typed
      additionalProperties. OpenAI-flavoured strict mode REJECTS a typed
      additionalProperties, and `required` makes a note per source structurally
      mandatory instead of a polite request the model may ignore (observed:
      gpt-oss-20b dropping notes and blowing the word budget on some prompts).
    """
    keys = [str(i + 1) for i in range(n)]
    return {'type': 'json_schema', 'json_schema': {
        'name': 'search_analysis', 'strict': True,
        'schema': {'type': 'object', 'properties': {
            'summary': {'type': 'string'},
            'notes': {'type': 'object',
                      'properties': {k: {'type': 'string'} for k in keys},
                      'required': keys, 'additionalProperties': False}},
            'required': ['summary', 'notes'], 'additionalProperties': False}}}


_SW_STOP = frozenset(('the', 'and', 'for', 'are', 'was', 'were', 'what', 'which',
                      'who', 'how', 'why', 'when', 'where', 'this', 'that', 'with',
                      'from', 'into', 'does', 'did', 'has', 'have', 'can', 'you'))


def _sw_focus(desc, q, max_chars=220):
    """Keep only the query-relevant sentence(s) of a Brave description.

    Decode dominates the job budget, but prefill isn't free either, and a
    saturated box feels every token. This only ever NARROWS text that was
    already prompt-only, so the Brave-ToS posture is unchanged: descriptions
    still never leave this function's caller, and graph notes stay the LLM's
    own words."""
    d = (desc or '').strip()
    if len(d) <= max_chars:
        return d
    terms = set(w for w in re.findall(r'[a-z0-9]+', q.lower())
                if len(w) > 2 and w not in _SW_STOP)
    sents = [s.strip() for s in re.split(r'(?<=[.!?])\s+', d) if s.strip()]
    if not terms or not sents:
        return d[:max_chars]
    scored = sorted(((sum(1 for w in terms if w in s.lower()), -i, s)
                     for i, s in enumerate(sents)), reverse=True)
    if scored[0][0] == 0:
        return d[:max_chars]          # nothing matched (synonyms only) — take the head
    keep, out = [], 0
    for score, negi, s in scored[:2]:  # 2 not 1: cheap hedge against a mis-scored pick
        if score == 0 or out + len(s) > max_chars:
            break
        keep.append((-negi, s))
        out += len(s)
    if not keep:
        return d[:max_chars]
    return ' '.join(s for _, s in sorted(keep))


def _sw_llm(q, results, deadline=None):
    """(answer, model, notes, tokens) via the local hermes gateway, falling
    back to the night-runner provider config. Empty values when no model is
    reachable — sources + graph are still worth fulfilling."""
    src = '\n\n'.join('[%d] %s\n%s\n%s' % (i + 1, r['title'], r['url'],
                                           _sw_focus(r['description'], q))
                      for i, r in enumerate(results))
    # "summary" MUST come first and the prompt must say so: the salvage parser
    # in _sw_parse_analysis relies on it, so a stream the deadline cuts short
    # still yields a complete summary and merely loses some notes. Field order
    # here is load-bearing, not cosmetic.
    prompt = (
        'You are a research summarizer. Use ONLY the numbered sources below.\n'
        '"summary": 2-3 sentences, 60 words MAX total, on what the sources '
        'collectively say about the query, citing like [1][3]. If they do not '
        'answer it, say what they DO cover.\n'
        '"notes": for each source number, at most 12 words in your own words on '
        'what that page contributes. Do not copy phrases from the source text.\n'
        'Output STRICT JSON only — no prose, no code fence, "summary" first:\n'
        '{"summary": "...", "notes": {"1": "...", "2": "..."}}\n\n'
        'Query: %s\n\nSources:\n%s' % (q, src))
    want = _sw_model()
    payload = {'model': want,
               'messages': [{'role': 'user', 'content': prompt}],
               'max_tokens': _SW_MAX_TOKENS,
               'response_format': _sw_schema(len(results))}
    _sw_trial_notice()

    def _finish(text, rmodel, tokens, truncated, via, trust_reported=True):
        """trust_reported=False for the gateway: it ECHOES the model you asked
        for while actually running config.yaml's. That echo can never disagree
        with our request, so believing it would write a confident lie onto a
        permanent public entry (ask for nemotron, get gpt-oss, entry says
        nemotron). On that path the config IS the truth."""
        summary, notes = _sw_parse_analysis(text)
        if truncated:
            print('[search-worker] %s truncated at deadline — salvaged %d chars, %d notes'
                  % (via, len(summary or ''), len(notes)))
        if trust_reported:
            actual = rmodel or want
        else:
            import night_runner as _nr
            actual = _nr.read_model_config(_hermes_home())['model'] or via
        if actual and want and actual.split('/')[-1] != want.split('/')[-1]:
            print('[search-worker] asked for %s but %s answered as %s — the chip reports what '
                  'actually ran' % (want, via, actual))
        # A salvaged partial IS success: falling through would trade a usable
        # summary for a second full decode we can't afford inside the lease.
        return summary, (actual or via)[:80], notes, tokens

    # 1. Direct to the operator's brain. Preferred: the agent gateway layers a
    #    ~13k-token system prompt on every call and ignores max_tokens/model.
    b = _sw_backend()
    if b is not None:
        try:
            text, rmodel, tokens, truncated = _sw_stream_or_block(
                b.url, b.headers, dict(payload, model=b.model or want), deadline)
            return _finish(text, rmodel, tokens, truncated, b.provider or 'direct')
        except Exception as e:
            print('[search-worker] direct backend failed (provider=%s url=%s model=%s): %s'
                  % (b.provider, b.url, b.model or want, str(e)[:120]))

    # 2. Fall back to the agent gateway. Costly and it ignores our params, but
    #    it's the only path for operators we can't resolve directly (e.g.
    #    anthropic, which has no _PROVIDER_ENDPOINTS entry). Do NOT delete.
    key = _sw_hermes_key()
    if key:
        try:
            text, rmodel, tokens, truncated = _sw_stream_or_block(
                'http://%s:%d/v1/chat/completions' % (HERMES_HOST, HERMES_PORT),
                {'Content-Type': 'application/json', 'Authorization': 'Bearer ' + key},
                payload, deadline)
            return _finish(text, rmodel, tokens, truncated, 'hermes-local', trust_reported=False)
        except Exception as e:
            print('[search-worker] hermes llm failed (model=%s): %s' % (want, str(e)[:120]))
    elif b is None:
        print('[search-worker] no direct backend and no gateway key (%s/.env, '
              'API_SERVER_KEY) — answers will be sources-only' % _hermes_home())

    # 3. Sources + graph are still worth fulfilling.
    return '', '', {}, 0


_SW_PALETTE = ['#7DC9B0', '#C9B8E0', '#E8A9A9', '#F0C987',
               '#9BC0E8', '#B8E09A', '#E0A47C', '#D89BE0']


def _sw_domain(u):
    try:
        return (urllib.parse.urlparse(u).hostname or u).replace('www.', '', 1)
    except Exception:
        return u


def _sw_color(name):
    h = 0
    for ch in name:
        h = (h * 31 + ord(ch)) & 0xFFFFFFFF
    return _SW_PALETTE[h % len(_SW_PALETTE)]


def _sw_graph(q, results, notes=None):
    """Query hub → result ring → domain ring, wrapped exactly as
    graph-viewer.html expects. Result nodes carry url/domain/note attributes
    so the viewer can render hover cards and click-through — the note is the
    worker LLM's own one-liner (never Brave's description text)."""
    import math
    notes = notes or {}
    nodes = [{'key': 'q', 'attributes': {'label': q, 'size': 18, 'x': 0, 'y': 0,
                                         'color': '#F5D25D', 'kind': 'query'}}]
    edges = []
    domains = {}
    for i, r in enumerate(results):
        a = (i / max(1, len(results))) * math.pi * 2
        d = _sw_domain(r['url'])
        nodes.append({'key': 'r%d' % i, 'attributes': {
            'label': r['title'][:60], 'size': 8,
            'x': math.cos(a) * 10, 'y': math.sin(a) * 10,
            'color': _sw_color(d), 'kind': 'source',
            'url': r['url'][:600], 'domain': d,
            'note': (notes.get(i) or '')[:200]}})
        edges.append({'key': 'eq%d' % i, 'source': 'q', 'target': 'r%d' % i, 'attributes': {}})
        domains.setdefault(d, []).append(i)
    for di, (d, ixs) in enumerate(domains.items()):
        a = (di / max(1, len(domains))) * math.pi * 2 + 0.35
        nodes.append({'key': 'd:' + d, 'attributes': {
            'label': d, 'size': 5 + len(ixs),
            'x': math.cos(a) * 17, 'y': math.sin(a) * 17, 'color': _sw_color(d),
            'kind': 'domain', 'domain': d}})
        for i in ixs:
            edges.append({'key': 'ed%d_%d' % (di, i), 'source': 'r%d' % i,
                          'target': 'd:' + d, 'attributes': {}})
    return json.dumps({'graph': {'options': {'type': 'mixed', 'multi': False,
                                             'allowSelfLoops': True},
                                 'attributes': {}, 'nodes': nodes, 'edges': edges},
                       'title': 'Search: ' + q, 'ts': int(time.time() * 1000)})


def _sw_process(job, deadline=None):
    jid, q = job['id'], job['q']
    started = time.monotonic()
    P = lambda s: urllib.parse.quote(s, safe='')
    try:
        print('[search-worker] claimed %s: %s' % (jid, q[:80]))
        results = _sw_brave(q, deadline)
        if not results:
            print('[search-worker] %s: no results (key set: %s)'
                  % (jid, bool(os.environ.get('BRAVE_API_KEY', '').strip())))
            _sw_call('fail', [jid, P('no brave key or no results')])
            return
        answer, model, notes, tokens = _sw_llm(q, results, deadline)
        if tokens:
            # Token burn rides in the model field ("llama-3.1-8b · ~1.2k tok")
            # — the canister entry has no dedicated field and this shows
            # everywhere the model chip renders, with zero schema changes.
            model = ('%s · ~%s tok' % (model or 'local', _sw_fmt_tokens(tokens)))[:80]
        # libValidate hard-rejects an over-cap answer, and a rejected fulfill
        # leaves the job claimed until the lease expires — burning an attempt
        # each time. The degrade path can hand us raw model prose, so clamp.
        if len(answer) > _SW_ANSWER_CAP:
            print('[search-worker] answer %d chars — clamping to %d' % (len(answer), _SW_ANSWER_CAP))
            answer = answer[:_SW_ANSWER_CAP]
        graph = _sw_graph(q, results, notes)
        lines = [jid, P(model), P('brave'), P(answer), str(len(results))]
        for r in results:
            lines.append('%s %s' % (P(r['title']), P(r['url'])))
        lines.append(P(graph))
        # Always attempt fulfill, even past our budget: the budget is soft and
        # sits 40s inside the lease, so the common case still lands. When we're
        # genuinely late the reap fires on the next worker's claim regardless,
        # so `attempts` increments whether we call or stay quiet — bailing would
        # save nothing and guarantee the loss.
        status, resp = _sw_call('fulfill', lines)
        if status == 200 and resp.get('ok'):
            print('[search-worker] fulfilled %s -> %s (%.0fs)'
                  % (jid, resp.get('libraryId'), time.monotonic() - started))
        else:
            print('[search-worker] fulfill rejected %s: %s %s' % (jid, status, resp))
    except Exception as e:
        print('[search-worker] job %s error: %s' % (jid, str(e)[:200]))
        _sw_call('fail', [jid, P(str(e)[:180])])


def _sw_loop():
    print('[search-worker] joining the search network as %s -> %s'
          % (_SW_PRINCIPAL[:12] + '…', _SW_BASE))
    paused = False
    while True:
        try:
            # Operator can take this node down from the admin panel
            # (gpuNode.enabled=false): stop claiming AND stop heartbeating so it
            # ages out of the network and clients see it offline — no SSH needed.
            gpu = _operator_config().get('gpuNode') or {}
            if gpu.get('enabled') is False:
                if not paused:
                    print('[search-worker] paused by operator (gpuNode.enabled=false)')
                    paused = True
                time.sleep(20)
                continue
            if paused:
                print('[search-worker] resumed by operator')
                paused = False
            # Clock starts BEFORE the claim: the canister stamps claimedAt during
            # that request, so a pre-call reading is guaranteed conservative.
            # monotonic, never time.time() — an NTP step mid-job would silently
            # blow (or extend) the budget.
            t0 = time.monotonic()
            status, resp = _sw_call('claim', [])
            if status == 403:
                # Not approved yet (or suspended): heartbeat keeps the
                # registration's lastSeen fresh so the operator sees "connected".
                _sw_call('heartbeat', [])
            elif status == 200 and resp.get('job'):
                _sw_process(resp['job'], t0 + _SW_JOB_BUDGET)
                continue          # drain the queue without sleeping between jobs
        except Exception as e:
            print('[search-worker] loop error:', str(e)[:200])
        time.sleep(20)


if _SW_ENABLED and _SW_PRINCIPAL and re.fullmatch(r'[0-9a-fA-F]{64}', _SW_SECRET_HEX or ''):
    threading.Thread(target=_sw_loop, daemon=True, name='search-worker').start()
elif _SW_ENABLED:
    print('[search-worker] SEARCH_WORKER=1 but WORKER_PRINCIPAL/WORKER_SECRET missing — not starting')


# ---- External approvals (Claude Code PreToolUse hook bridge) ---------------
# When the local `claude` CLI is wired to call our approval hook, each tool
# invocation it wants to make blocks on a long-poll against this server. We
# park the request here, surface it in the HQ UI's ApprovalTray, and signal
# the waiting hook process when the boss clicks APPROVE / REJECT.
#
# Single-process, thread-safe. Entries auto-expire so a closed terminal
# doesn't leak forever in the tray.
_approvals_lock = threading.Lock()
_approvals_pending = {}   # id -> dict (see _new_approval)
_APPROVAL_TTL_SEC = 30 * 60

def _new_approval(payload: dict) -> dict:
    aid = 'ce_' + uuid.uuid4().hex[:10]
    entry = {
        'id': aid,
        'ts': time.time(),
        'tool': str(payload.get('tool') or '')[:80],
        'input': payload.get('input') or {},
        'cwd': str(payload.get('cwd') or '')[:400],
        'agent': str(payload.get('agent') or 'claude-code')[:60],
        'sessionId': str(payload.get('sessionId') or '')[:80],
        'summary': str(payload.get('summary') or '')[:400],
        'decision': None,           # 'allow' | 'deny' | None
        'reason': '',
        'event': threading.Event(),
    }
    with _approvals_lock:
        _approvals_pending[aid] = entry
    return entry

def _gc_approvals():
    """Drop entries older than TTL with no decision — terminal probably went away."""
    now = time.time()
    with _approvals_lock:
        stale = [k for k, e in _approvals_pending.items()
                 if e['decision'] is None and now - e['ts'] > _APPROVAL_TTL_SEC]
        for k in stale:
            e = _approvals_pending.pop(k, None)
            if e:
                e['decision'] = 'deny'
                e['reason'] = 'expired (no human present)'
                e['event'].set()


# Vault config: env vars at startup, can be re-set at runtime via POST /vault/configure.
# By default CafresoHQ owns a plain Markdown vault under hq-state, so Obsidian is
# optional instead of required for notes to work.
_default_vault_root = _hq_state_dir / 'vault'
_vault_root     = os.environ.get('CAFRESOHQ_VAULT', str(_default_vault_root)).strip()
_vault_backend  = os.environ.get('CAFRESOHQ_VAULT_BACKEND', 'fs').strip() or 'fs'
_vault_rest_url = os.environ.get('CAFRESOHQ_OBSIDIAN_URL', 'https://127.0.0.1:27124').strip()
_vault_rest_key = os.environ.get('CAFRESOHQ_OBSIDIAN_KEY', '').strip()

# ── OCI Object Storage vault config (CAFRESOHQ_VAULT_BACKEND=oci) ─────────────
# Used by OCI Fleet containers.  serve.py never reads ~/.oci/config on its own;
# the Container Instance's instance-principal auth or a mounted config file
# provides credentials.  All three vars must be set for OCI vault to work.
_oci_vault_namespace = os.environ.get('OCI_VAULT_NAMESPACE', '').strip()
_oci_vault_bucket    = os.environ.get('OCI_VAULT_BUCKET', 'cafresohq-fleet-vault').strip()
# Prefix isolates this user's objects: e.g. "2vxsx-principal-hash/"
_oci_vault_prefix    = os.environ.get('OCI_VAULT_PREFIX', '').strip()
_oci_client_lock     = threading.Lock()
_oci_client_obj: object = None   # lazy singleton — oci.object_storage.ObjectStorageClient

def _oci_object_client():
    """Lazy-load OCI Object Storage client. Thread-safe. Raises RuntimeError on failure."""
    global _oci_client_obj
    if _oci_client_obj is None:
        with _oci_client_lock:
            if _oci_client_obj is None:
                try:
                    import oci as _oci_sdk  # pip install oci

                    # ── 1. Env-var credentials (preferred inside Container Instances) ──
                    # Fleet-manager passes these at provision time so we don't rely on
                    # instance-principal IMDS, which can hang if IAM policies aren't ready.
                    _ev_tenancy     = os.environ.get('OCI_TENANCY_OCID', '').strip()
                    _ev_user        = os.environ.get('OCI_USER_OCID', '').strip()
                    _ev_fingerprint = os.environ.get('OCI_FINGERPRINT', '').strip()
                    _ev_region      = os.environ.get('OCI_REGION', 'us-ashburn-1').strip()
                    _ev_key_b64     = os.environ.get('OCI_KEY_B64', '').strip()
                    if all([_ev_tenancy, _ev_user, _ev_fingerprint, _ev_key_b64]):
                        import base64 as _b64
                        _cfg = {
                            'tenancy':     _ev_tenancy,
                            'user':        _ev_user,
                            'fingerprint': _ev_fingerprint,
                            'region':      _ev_region,
                            'key_content': _b64.b64decode(_ev_key_b64).decode(),
                        }
                        _oci_sdk.config.validate_config(_cfg)
                        _oci_client_obj = _oci_sdk.object_storage.ObjectStorageClient(_cfg)
                        return _oci_client_obj

                    # ── 2. Config file (~/.oci/config) ───────────────────────────────
                    try:
                        _cfg = _oci_sdk.config.from_file()
                        _oci_client_obj = _oci_sdk.object_storage.ObjectStorageClient(_cfg)
                        return _oci_client_obj
                    except Exception:
                        pass

                    # ── 3. Instance principal (OCI Compute; may be slow to init) ────
                    _cfg = {}
                    signer = _oci_sdk.auth.signers.InstancePrincipalsSecurityTokenSigner()
                    _oci_client_obj = _oci_sdk.object_storage.ObjectStorageClient(
                        config=_cfg, signer=signer)
                except ImportError:
                    raise RuntimeError(
                        'OCI SDK not installed — run: pip install oci  '
                        '(or pip install -r oci-fleet/requirements-serve.txt)')
    return _oci_client_obj

def _oci_obj_key(rel: str) -> str:
    """Build the full OCI Object key for a vault-relative path."""
    prefix = (_oci_vault_prefix.rstrip('/') + '/') if _oci_vault_prefix else ''
    return prefix + rel.lstrip('/')

# Fleet identity (set by fleet-manager when provisioning the container)
_fleet_mode      = os.environ.get('CAFRESOHQ_FLEET_MODE', 'local').strip()
_fleet_user_principal = os.environ.get('USER_PRINCIPAL', '').strip()

# Uptime tracking for /health
_server_start_time = time.time()

if not os.environ.get('CAFRESOHQ_VAULT'):
    try:
        pathlib.Path(_vault_root).mkdir(parents=True, exist_ok=True)
    except OSError:
        pass


def _clean_vault_root(raw) -> str:
    """Normalize common pasted path forms before validating a local vault."""
    root = str(raw or '').strip()
    while len(root) >= 2 and root[0] == root[-1] and root[0] in ('"', "'"):
        root = root[1:-1].strip()
    if root.lower().startswith('file:'):
        parsed = urllib.parse.urlparse(root)
        if parsed.scheme == 'file':
            path = urllib.parse.unquote(parsed.path or '')
            if parsed.netloc:
                path = f'//{parsed.netloc}{path}'
            if os.name == 'nt' and re.match(r'^/[A-Za-z]:[\\/]', path):
                path = path[1:]
            root = path
    root = os.path.expandvars(os.path.expanduser(root))
    if '%' in root:
        root = urllib.parse.unquote(root)
    return root.strip()


def _obsidian_config_dirs() -> list:
    """Return likely Obsidian desktop config directories for this platform."""
    out = []
    appdata = os.environ.get('APPDATA')
    if appdata:
        out.append(pathlib.Path(appdata) / 'obsidian')
    home = pathlib.Path.home()
    if sys.platform == 'darwin':
        out.append(home / 'Library' / 'Application Support' / 'obsidian')
    else:
        xdg = os.environ.get('XDG_CONFIG_HOME')
        out.append((pathlib.Path(xdg) if xdg else home / '.config') / 'obsidian')
    return out


def _discover_obsidian_vaults() -> list:
    """Read Obsidian's desktop config and return known local vault paths."""
    seen = set()
    vaults = []
    for cfg_dir in _obsidian_config_dirs():
        cfg = cfg_dir / 'obsidian.json'
        if not cfg.is_file():
            continue
        try:
            data = json.loads(cfg.read_text(encoding='utf-8'))
        except Exception:
            continue
        raw_vaults = data.get('vaults') or {}
        if not isinstance(raw_vaults, dict):
            continue
        for vault_id, info in raw_vaults.items():
            if not isinstance(info, dict):
                continue
            raw_path = str(info.get('path') or '').strip()
            if not raw_path:
                continue
            key = os.path.normcase(os.path.abspath(raw_path))
            if key in seen:
                continue
            seen.add(key)
            p = pathlib.Path(raw_path)
            vaults.append({
                'id': str(vault_id),
                'path': raw_path,
                'name': p.name or raw_path,
                'open': bool(info.get('open')),
                'ts': int(info.get('ts') or 0),
                'exists': p.is_dir(),
                'source': str(cfg),
            })
    vaults.sort(key=lambda v: (not v['open'], not v['exists'], -v['ts'], v['name'].lower()))
    return vaults


def _vault_resolve(rel: str) -> pathlib.Path:
    """Resolve `rel` (e.g. "Daily/2026-04-25.md") under the vault directory with
    traversal protection. Raises ValueError if escape is attempted or unset."""
    if not _vault_root:
        raise ValueError('vault directory not configured')
    root = pathlib.Path(_vault_root).resolve()
    if not root.is_dir():
        raise ValueError(f'vault directory does not exist: {root}')
    rel = rel.lstrip('/').replace('\\', '/')
    if not rel.endswith('.md'):
        rel += '.md'
    candidate = (root / rel).resolve()
    # Reject anything outside the vault directory.
    try:
        candidate.relative_to(root)
    except ValueError:
        raise ValueError('path escapes vault directory')
    return candidate


# ---- Obsidian Local REST API client -----------------------------------------
def _obsidian_request(method: str, upstream_path: str,
                      body: bytes = None, extra_headers: dict = None,
                      content_type: str = 'application/json'):
    """Forward a request to the Obsidian Local REST API plugin. Returns
    (status, headers_list, body_bytes). Raises ValueError if not configured.
    Tolerates self-signed HTTPS certs (the plugin's default)."""
    if not _vault_rest_url:
        raise ValueError('Obsidian REST URL not configured')
    if not _vault_rest_key:
        raise ValueError('Obsidian REST key not configured')
    parsed = urllib.parse.urlparse(_vault_rest_url)
    is_https = parsed.scheme == 'https'
    host = parsed.hostname or '127.0.0.1'
    port = parsed.port or (27124 if is_https else 27123)
    headers = {
        'Authorization': f'Bearer {_vault_rest_key}',
        'Accept': 'application/json',
    }
    if body is not None:
        headers['Content-Type'] = content_type
    if extra_headers:
        headers.update(extra_headers)
    if is_https:
        ctx = ssl._create_unverified_context()  # plugin uses self-signed cert
        conn = http.client.HTTPSConnection(host, port, timeout=30, context=ctx)
    else:
        conn = http.client.HTTPConnection(host, port, timeout=30)
    try:
        conn.request(method, upstream_path, body=body, headers=headers)
        resp = conn.getresponse()
        return resp.status, resp.getheaders(), resp.read()
    finally:
        conn.close()


# ---- Obsidian REST: vault adapter -------------------------------------------
def _rest_list_all() -> list:
    """Walk the vault via the REST API. Returns the same shape as the FS
    listing: [{path, title, mtime, size}, ...]."""
    out = []
    def walk(folder: str):
        s, _, body = _obsidian_request('GET', '/vault/' + urllib.parse.quote(folder))
        if s != 200:
            return
        try:
            data = json.loads(body.decode('utf-8'))
        except Exception:
            return
        for entry in data.get('files', []):
            full = (folder + entry).lstrip('/')
            if entry.endswith('/'):
                # Subfolder; skip Obsidian internals
                if entry.lstrip('/').startswith('.'):
                    continue
                walk(full)
            elif entry.endswith('.md'):
                # Plugin's listing doesn't include mtime/size in the simple list;
                # fetch a stat-like via the file endpoint headers if needed. For
                # speed we just stub them. The graph + browse views still work.
                title = pathlib.PurePosixPath(full).stem
                out.append({'path': full, 'title': title, 'mtime': 0, 'size': 0})
    walk('')
    out.sort(key=lambda f: f['path'])
    return out


def _rest_search(query: str, limit: int = 10) -> list:
    qs = urllib.parse.urlencode({'query': query, 'contextLength': 120})
    s, _, body = _obsidian_request('GET', '/search/simple/?' + qs)
    if s != 200:
        return []
    try:
        data = json.loads(body.decode('utf-8'))
    except Exception:
        return []
    hits = []
    for r in data[:limit]:
        filename = r.get('filename') or r.get('path') or ''
        matches = r.get('matches', [])
        snippet = matches[0].get('context', '') if matches else ''
        hits.append({
            'path': filename,
            'title': pathlib.PurePosixPath(filename).stem,
            'score': r.get('score', len(matches)),
            'snippet': snippet.replace('\n', ' ').strip(),
        })
    return hits


# ---- Graph extraction (works for both backends) -----------------------------
import re as _re
_WIKILINK_RE = _re.compile(r'\[\[([^\]\|]+?)(?:\|[^\]]*)?\]\]')
_TAG_RE = _re.compile(r'(?:^|\s)#([A-Za-z0-9_/\-]+)')


def _graph_node_type(path: str, tags=None) -> str:
    """Lightweight product-facing graph taxonomy. Keep this intentionally small.
    The UI can reason with these without requiring a canonical backend schema yet."""
    p = (path or '').replace('\\', '/').lower()
    tagset = {str(t).lstrip('#').lower() for t in (tags or [])}
    name = pathlib.PurePosixPath(p).stem
    if 'research/' in p or 'research' in tagset:
        return 'research'
    if 'project' in tagset or p.startswith('projects/') or p.startswith('05-projects/') or '/projects/' in p:
        return 'project'
    if 'task' in tagset or 'tasks/' in p or 'todo' in name:
        return 'task'
    if 'agent' in tagset or 'agents/' in p or name.startswith('agent-'):
        return 'agent'
    if 'decision' in tagset or 'decisions/' in p or 'decision' in name:
        return 'decision'
    if 'proposal' in tagset or 'proposal' in name or '/proposals/' in p:
        return 'proposal'
    if 'memory' in tagset or 'memory/' in p or 'memories/' in p:
        return 'memory'
    if 'risk' in tagset or 'risk' in name:
        return 'risk'
    if p.endswith(('.py.md', '.jsx.md', '.js.md', '.ts.md', '.tsx.md')) or 'code' in tagset:
        return 'code_file'
    return 'note'


def _graph_edge_payload(source: str, target: str, edge_type='links_to', status='canonical', confidence=1.0, source_text='') -> dict:
    return {
        'source': source,
        'target': target,
        'type': edge_type,
        'status': status,
        'confidence': confidence,
        'sourcePath': source,
        'sourceText': source_text,
    }


# ──────────────────────────────────────────────────────────────────────────
# Edge-type extraction (knowledge-graph style relationships)
#
# Every wikilink in a markdown file becomes an edge. By default that edge is
# `links_to` (the historical behavior). But we look at the SURROUNDING context
# of each link — section heading, callout type, line prefix, frontmatter key —
# to upgrade generic links into typed knowledge-graph relationships
# (cites, supports, contradicts, child_of, decided, blocked_by, etc).
#
# This makes the graph reason-able instead of just browse-able. Confidence
# stays high (>=0.7) for explicit cues; pure body links remain confidence 1.0
# typed links_to so we don't degrade the existing graph.
# ──────────────────────────────────────────────────────────────────────────

# Section heading → edge type for any links found INSIDE that section.
# Match is case-insensitive on the heading text; trailing punctuation stripped.
_SECTION_EDGE_TYPES = {
    'sources': 'cites',
    'source': 'cites',
    'references': 'cites',
    'reference': 'cites',
    'citations': 'cites',
    'see also': 'related_to',
    'related': 'related_to',
    'related notes': 'related_to',
    'further reading': 'related_to',
    'children': 'parent_of',
    'sub-pages': 'parent_of',
    'parent': 'child_of',
    'parents': 'child_of',
    'implementation': 'implements',
    'implemented by': 'implemented_by',
    'implements': 'implements',
    'risks': 'has_risk',
    'risk': 'has_risk',
    'decisions': 'decided',
    'decided': 'decided',
    'decision': 'decided',
    'evidence': 'supports',
    'supports': 'supports',
    'contradicts': 'contradicts',
    'objections': 'contradicts',
    'counter-evidence': 'contradicts',
    'blocks': 'blocks',
    'blocked by': 'blocked_by',
    'depends on': 'depends_on',
    'dependencies': 'depends_on',
    'supersedes': 'supersedes',
    'superseded by': 'superseded_by',
    'agents': 'created_by',
    'assigned to': 'assigned_to',
    'tasks': 'has_task',
    'proposals': 'has_proposal',
}

# A line that BEGINS with one of these prefixes (case-insensitive) and contains
# a wikilink is treated as a typed edge to that link. e.g. `Parent: [[Foo]]`.
_LINE_PREFIX_EDGE_TYPES = {
    'source': 'cites',
    'sources': 'cites',
    'cite': 'cites',
    'cites': 'cites',
    'reference': 'cites',
    'references': 'cites',
    'parent': 'child_of',
    'parents': 'child_of',
    'child': 'parent_of',
    'children': 'parent_of',
    'see also': 'related_to',
    'related': 'related_to',
    'implements': 'implements',
    'implemented by': 'implemented_by',
    'blocks': 'blocks',
    'blocked by': 'blocked_by',
    'depends on': 'depends_on',
    'depends': 'depends_on',
    'supersedes': 'supersedes',
    'superseded by': 'superseded_by',
    'replaces': 'supersedes',
    'replaced by': 'superseded_by',
    'assigned to': 'assigned_to',
    'assigned': 'assigned_to',
    'owner': 'assigned_to',
    'created by': 'created_by',
    'author': 'created_by',
    'authors': 'created_by',
    'edited by': 'edited_by',
    'reviewer': 'reviewed_by',
    'reviewed by': 'reviewed_by',
    'decision': 'decided',
    'decided': 'decided',
    'decides': 'decided',
    'note': 'notes',
    'tldr': 'notes',
    'summary': 'notes',
    'mention': 'mentions',
    'mentions': 'mentions',
}

# Callout block types (> [!type]) → edge type for links inside the callout.
_CALLOUT_EDGE_TYPES = {
    'supports': 'supports',
    'support': 'supports',
    'evidence': 'supports',
    'contradicts': 'contradicts',
    'contradict': 'contradicts',
    'refutes': 'contradicts',
    'objection': 'contradicts',
    'derived': 'derived_from',
    'derived-from': 'derived_from',
    'derives': 'derived_from',
    'cite': 'cites',
    'cites': 'cites',
    'source': 'cites',
    'note': 'notes',
    'info': 'notes',
    'abstract': 'notes',
    'summary': 'notes',
    'warning': 'has_risk',
    'caution': 'has_risk',
    'danger': 'has_risk',
    'risk': 'has_risk',
    'tip': 'related_to',
    'hint': 'related_to',
    'example': 'exemplifies',
    'decision': 'decided',
    'decided': 'decided',
    'todo': 'has_task',
    'task': 'has_task',
    'question': 'questions',
    'q': 'questions',
}

# Frontmatter key → (edge_type, is_inverse). is_inverse means the edge points
# FROM the linked target back TO this note (e.g. `children: [[X]]` means
# this note is parent_of X, but `parent: [[X]]` means this note is child_of X).
_FRONTMATTER_EDGE_TYPES = {
    'parent': ('child_of', False),
    'parents': ('child_of', False),
    'children': ('parent_of', False),
    'child': ('parent_of', False),
    'source': ('cites', False),
    'sources': ('cites', False),
    'cites': ('cites', False),
    'references': ('cites', False),
    'related': ('related_to', False),
    'see-also': ('related_to', False),
    'see_also': ('related_to', False),
    'supersedes': ('supersedes', False),
    'superseded-by': ('superseded_by', False),
    'replaces': ('supersedes', False),
    'replaced-by': ('superseded_by', False),
    'implements': ('implements', False),
    'implemented-by': ('implemented_by', False),
    'blocks': ('blocks', False),
    'blocked-by': ('blocked_by', False),
    'depends-on': ('depends_on', False),
    'depends': ('depends_on', False),
    'agent': ('created_by', False),
    'author': ('created_by', False),
    'authors': ('created_by', False),
    'created-by': ('created_by', False),
    'assigned-to': ('assigned_to', False),
    'assignedto': ('assigned_to', False),
    'assignee': ('assigned_to', False),
    'owner': ('assigned_to', False),
    'reviewer': ('reviewed_by', False),
    'reviewed-by': ('reviewed_by', False),
}

# Default confidence per type. Explicit frontmatter / line-prefix is 1.0;
# section context is 0.85 (fairly likely); callouts are 0.95 (explicit author
# intent). Generic body links_to stays at 1.0 (the link itself is a fact).
_EDGE_TYPE_CONFIDENCE = {
    'links_to': 1.0,
    'cites': 0.95,
    'related_to': 0.85,
    'child_of': 1.0,
    'parent_of': 1.0,
    'implements': 0.95,
    'implemented_by': 0.95,
    'blocks': 0.95,
    'blocked_by': 0.95,
    'supersedes': 1.0,
    'superseded_by': 1.0,
    'depends_on': 0.95,
    'supports': 0.95,
    'contradicts': 0.95,
    'derived_from': 0.95,
    'has_risk': 0.9,
    'decided': 0.9,
    'created_by': 1.0,
    'assigned_to': 1.0,
    'edited_by': 1.0,
    'reviewed_by': 1.0,
    'notes': 0.7,
    'mentions': 0.85,
    'has_task': 0.9,
    'has_proposal': 0.9,
    'exemplifies': 0.85,
    'questions': 0.85,
    # HQ-state ingestion edge types (tasks/missions/receipts/etc → vault)
    'references': 0.9,   # task/mission body mentions a note
    'produces':   0.95,  # mission wrote a note
    'modified':   0.85,  # receipt records a write to a note
    'targets':    0.85,  # mission scoped to a folder/note
    'runs_as':    1.0,   # mission runs as an agent
    # Message-registry edges
    'sent_to':    1.0,   # agent originated a message-thread
    'received':   1.0,   # agent received a message-thread
    # Org-chart edges
    'reports_to': 1.0,   # assistant agent → senior agent
}


_FRONTMATTER_RE = _re.compile(r'\A---\s*\n(.*?\n)---\s*\n', _re.DOTALL)
_HEADING_RE = _re.compile(r'^(#{1,6})\s+(.+?)\s*$')
_CALLOUT_OPEN_RE = _re.compile(r'^\s*>\s*\[!([A-Za-z][\w\-]*)\]')
# Grabs any line of the form "Word(s):" at the very start, captures the label.
_LINE_PREFIX_RE = _re.compile(r'^\s*(?:[-*]\s+)?\*{0,2}([A-Za-z][\w\s\-]{0,30}?)\*{0,2}\s*:\s*(.*)$')


def _norm_label(s: str) -> str:
    """Lowercase + collapse whitespace + strip trailing punctuation. Used to
    match section headings / line prefixes against the type lookup tables."""
    s = (s or '').strip().lower().rstrip('.:!?')
    return _re.sub(r'\s+', ' ', s)


def _classify_edge(context: dict) -> tuple:
    """Given the context surrounding a wikilink, return (edge_type, confidence,
    source_text). Priority: callout > line prefix > section heading > default.
    `context` keys: callout, prefix, section, line_text."""
    callout = context.get('callout')
    prefix  = context.get('prefix')
    section = context.get('section')
    text    = (context.get('line_text') or '').strip()[:240]
    if callout:
        t = _CALLOUT_EDGE_TYPES.get(callout.lower())
        if t:
            return t, _EDGE_TYPE_CONFIDENCE.get(t, 0.9), text
    if prefix:
        t = _LINE_PREFIX_EDGE_TYPES.get(_norm_label(prefix))
        if t:
            return t, _EDGE_TYPE_CONFIDENCE.get(t, 1.0), text
    if section:
        t = _SECTION_EDGE_TYPES.get(_norm_label(section))
        if t:
            return t, _EDGE_TYPE_CONFIDENCE.get(t, 0.85), text
    return 'links_to', 1.0, text


def _extract_typed_edges(rel: str, text: str, all_paths: dict):
    """Walk `text` line-by-line and yield (target_path, edge_type, confidence,
    source_text) for every resolvable wikilink, classifying each by its
    surrounding context. Also processes YAML frontmatter for typed metadata.

    Yields tuples; the caller dedups by (rel, target, type). We dedup on type
    too so a note can have BOTH a `cites` AND a `related_to` edge to the same
    target if the note links to it from two different contexts — useful for
    showing the strongest relationship while keeping evidence."""
    # ── Frontmatter pass ───────────────────────────────────────────────
    body = text
    fm_match = _FRONTMATTER_RE.match(text)
    if fm_match:
        fm_text = fm_match.group(1)
        body = text[fm_match.end():]
        for fm_line in fm_text.splitlines():
            kv = fm_line.split(':', 1)
            if len(kv) != 2: continue
            key = _norm_label(kv[0]).replace(' ', '-')
            val = kv[1].strip()
            mapping = _FRONTMATTER_EDGE_TYPES.get(key)
            if not mapping: continue
            edge_type, _ = mapping
            conf = _EDGE_TYPE_CONFIDENCE.get(edge_type, 0.9)
            for m in _WIKILINK_RE.finditer(val):
                tgt = _norm_link(m.group(1), all_paths)
                if tgt and tgt != rel:
                    yield (tgt, edge_type, conf, f'{key}: {val}'[:240])

    # ── Body pass: line-by-line, tracking section + callout state ──────
    current_section = None
    callout_type = None
    callout_lines_left = 0  # how many continuation `>` lines still belong to the open callout
    for line in body.splitlines():
        # Heading? Update current_section, no links here.
        h = _HEADING_RE.match(line)
        if h:
            current_section = h.group(2).strip()
            callout_type = None
            callout_lines_left = 0
            continue
        # Callout open?
        co = _CALLOUT_OPEN_RE.match(line)
        if co:
            callout_type = co.group(1)
            callout_lines_left = 12  # generous: callouts up to ~12 lines
        elif callout_type:
            # Callout continues only while line begins with `>`
            if line.lstrip().startswith('>') and callout_lines_left > 0:
                callout_lines_left -= 1
            else:
                callout_type = None
                callout_lines_left = 0
        # Line prefix?  (Only outside callouts; callouts already imply type.)
        prefix = None
        if not callout_type:
            lp = _LINE_PREFIX_RE.match(line)
            if lp:
                cand = lp.group(1)
                if _norm_label(cand) in _LINE_PREFIX_EDGE_TYPES:
                    prefix = cand
        # Now scan wikilinks on this line.
        for m in _WIKILINK_RE.finditer(line):
            tgt = _norm_link(m.group(1), all_paths)
            if not tgt or tgt == rel:
                continue
            ctx = {
                'callout': callout_type,
                'prefix': prefix,
                'section': current_section,
                'line_text': line.strip(),
            }
            etype, conf, text_snip = _classify_edge(ctx)
            yield (tgt, etype, conf, text_snip)

def _norm_link(target: str, all_paths: dict) -> str:
    """Resolve a wikilink target string to a known path key, or '' if unknown.
    Wikilinks can be by stem ("Foo") or path ("dir/Foo"); the .md is implicit."""
    target = target.strip()
    if not target:
        return ''
    # Strip section anchors (#heading) and block refs (^id)
    for sep in ('#', '^'):
        i = target.find(sep)
        if i >= 0: target = target[:i]
    if not target:
        return ''
    candidates = [
        target,
        target + '.md',
        target.rstrip('/') + '.md',
    ]
    # Lowercase fallback for case-insensitive matching
    by_stem = all_paths['by_stem']
    by_path = all_paths['by_path']
    for c in candidates:
        if c in by_path:
            return c
    stem = pathlib.PurePosixPath(target).stem.lower()
    return by_stem.get(stem, '')


def _build_graph_fs() -> dict:
    """Build a graph from raw .md files under the configured vault directory."""
    if not _vault_root:
        raise ValueError('vault not configured')
    root = pathlib.Path(_vault_root).resolve()
    nodes = []
    raw = []
    by_path = {}
    by_stem = {}
    for p in root.rglob('*.md'):
        try:
            rel = str(p.relative_to(root)).replace('\\', '/')
        except ValueError:
            continue
        if any(part.startswith('.') for part in p.relative_to(root).parts):
            continue
        try:
            text = p.read_text(encoding='utf-8', errors='replace')
        except OSError:
            continue
        by_path[rel] = True
        by_stem[p.stem.lower()] = rel
        raw.append((rel, p.stem, text, p.stat().st_mtime if p.exists() else 0, p.stat().st_size if p.exists() else 0))
    all_paths = {'by_path': by_path, 'by_stem': by_stem}
    edges = []
    # Dedupe by (source, target, type) so a note can have multiple edge types
    # to the same target (e.g. cited AND related_to) but never duplicates
    # within a single type. We also track per-pair best-confidence to upgrade
    # repeat plain links_to into typed edges if a typed mention shows up later.
    seen_typed_edge = set()
    for rel, title, text, mtime, size in raw:
        tags = sorted(set(m.group(1) for m in _TAG_RE.finditer(text)))
        out_targets = set()
        # Walk every (target, type, conf, snippet) the extractor finds.
        for tgt, etype, conf, snip in _extract_typed_edges(rel, text, all_paths):
            out_targets.add(tgt)
            key = (rel, tgt, etype)
            if key in seen_typed_edge: continue
            seen_typed_edge.add(key)
            edges.append(_graph_edge_payload(
                rel, tgt, edge_type=etype, confidence=conf, source_text=snip))
        nodes.append({
            'id': rel, 'title': title, 'path': rel,
            'mtime': int(mtime * 1000), 'size': size,
            'tags': tags[:8], 'outlinks': len(out_targets),
            'type': _graph_node_type(rel, tags),
            'source': 'markdownvault',
        })

    # ── HQ-state ingestion: tasks/projects/missions/receipts/agents ────
    # The graph stops being a vault-only viewer and starts representing
    # everything CafresoHQ knows. New nodes use prefixed IDs (`task:`,
    # `agent:` etc) so they never collide with vault paths. Edges link
    # them to each other AND to vault notes when references resolve.
    hq_nodes, hq_edges = _build_hq_state_graph(all_paths, seen_typed_edge)
    nodes.extend(hq_nodes)
    edges.extend(hq_edges)

    # Compute inlinks for sizing — across BOTH vault and HQ-state nodes.
    indeg = {n['id']: 0 for n in nodes}
    for e in edges:
        if e['target'] in indeg:
            indeg[e['target']] += 1
    for n in nodes:
        n['inlinks'] = indeg[n['id']]
    return {'nodes': nodes, 'edges': edges}


# ──────────────────────────────────────────────────────────────────────────
# Graph build cache (FS backend)
#
# /vault/graph rebuilds by reading + parsing every .md file and the hq-state
# JSON on each call — O(total text). On a cold container or a big vault that's
# a slow first paint of the graph. This cache keys the built graph on a CHEAP
# stat-only signature (each file's mtime+size, no reads); a repeat call with no
# file changes returns the cached graph instantly. Any edit/add/delete changes
# the signature and triggers exactly one rebuild. Single-process, so a plain
# module global under the server's thread is sufficient.
# ──────────────────────────────────────────────────────────────────────────

_graph_cache = {'sig': None, 'graph': None}
_graph_cache_lock = threading.Lock()


def _vault_graph_signature() -> str:
    """Cheap fingerprint of everything _build_graph_fs reads: each vault .md
    file's (rel, mtime, size) + the top-level hq-state JSON files. Stat-only,
    so it's far cheaper than the read+parse a full rebuild does."""
    if not _vault_root:
        return 'unconfigured'
    h = hashlib.sha1()
    root = pathlib.Path(_vault_root).resolve()
    try:
        for p in sorted(root.rglob('*.md')):
            try:
                rel = p.relative_to(root)
            except ValueError:
                continue
            if any(part.startswith('.') for part in rel.parts):
                continue
            try:
                st = p.stat()
            except OSError:
                continue
            h.update(('%s|%d|%d\n' % (str(rel).replace('\\', '/'), int(st.st_mtime), st.st_size)).encode('utf-8'))
    except OSError:
        return 'walk-error'
    try:
        for jp in sorted(_hq_state_dir.glob('*.json')):
            try:
                st = jp.stat()
                h.update(('S:%s|%d|%d\n' % (jp.name, int(st.st_mtime), st.st_size)).encode('utf-8'))
            except OSError:
                pass
    except OSError:
        pass
    return h.hexdigest()


def _build_graph_fs_cached() -> dict:
    """_build_graph_fs() with an mtime-signature cache (see above)."""
    sig = _vault_graph_signature()
    with _graph_cache_lock:
        if _graph_cache['sig'] == sig and _graph_cache['graph'] is not None:
            return _graph_cache['graph']
    graph = _build_graph_fs()
    with _graph_cache_lock:
        _graph_cache['sig'] = sig
        _graph_cache['graph'] = graph
    return graph


# ──────────────────────────────────────────────────────────────────────────
# HQ-state → graph ingestion
#
# Reads the JSON files in hq-state/ and turns them into typed graph nodes,
# with edges to each other AND to vault notes when references resolve. This
# is what makes GraphView the unified "everything CafresoHQ knows" view
# rather than just a markdown vault visualizer.
# ──────────────────────────────────────────────────────────────────────────

def _hq_state_path(name: str) -> pathlib.Path:
    """Resolve `hq-state/<name>.json` relative to the configured state dir."""
    return _hq_state_dir / f'{name}.json'


def _load_hq_state(name: str):
    """Safely load `hq-state/<name>.json`. Returns [] / {} on missing/invalid;
    never raises so a missing file doesn't break the whole graph build."""
    p = _hq_state_path(name)
    if not p.exists():
        return None
    try:
        with open(p, 'r', encoding='utf-8') as f:
            return json.load(f)
    except (OSError, json.JSONDecodeError) as e:
        sys.stderr.write(f'[graph] could not load {name}.json: {e}\n')
        return None


def _load_hq_memory(name: str):
    """Safely load `<hq-memory>/<name>.json`. Same shape contract as
    _load_hq_state but reads from _hq_memory_dir (where the CafresoHQ
    frontend persists agents, journals, etc via useFileStored('memory', ...))."""
    p = _hq_memory_dir / f'{name}.json'
    if not p.exists():
        return None
    try:
        with open(p, 'r', encoding='utf-8') as f:
            return json.load(f)
    except (OSError, json.JSONDecodeError) as e:
        sys.stderr.write(f'[graph] could not load memory/{name}.json: {e}\n')
        return None


def _load_authoritative_agents() -> dict:
    """Load the canonical agents list from hq-memory/agents.json (where the
    CafresoHQ frontend persists hired agents). Returns a dict keyed by both
    the agent id (e.g. 'a_nykw53') AND its name (e.g. 'Selvin'), with the
    agent dict as value. This lets the graph builder upgrade synthesised
    agent nodes — which would otherwise show as 'a_nykw53' — to friendly
    names ('Selvin · Code Gremlin') with proper sprite color metadata."""
    raw = _load_hq_memory('agents') or []
    if not isinstance(raw, list):
        return {}
    by_key = {}
    for a in raw:
        if not isinstance(a, dict): continue
        aid = a.get('id') or ''
        name = a.get('name') or ''
        if aid:  by_key[aid.lower()] = a
        if name: by_key[name.lower()] = a
    return by_key


def _slug_agent(s: str) -> str:
    """Normalize an agent name/id into a stable graph node id suffix."""
    return _re.sub(r'[^A-Za-z0-9_]+', '_', (s or 'unknown').strip()).strip('_').lower() or 'unknown'


def _build_hq_state_graph(all_paths: dict, seen_typed_edge: set) -> tuple:
    """Return (nodes, edges) ingested from hq-state/. Safe to call even if
    no JSON files exist — returns ([], []) in that case.

    `all_paths` is the same {by_path, by_stem} the wikilink resolver uses;
    we reuse it so when a task.detail mentions [[NoteName]] we can build a
    proper edge to the actual vault note.

    `seen_typed_edge` is the dedup set from the vault pass; we extend it
    so HQ-state edges don't duplicate vault edges by accident."""
    nodes_out = []
    edges_out = []

    def add_edge(src, tgt, etype, conf, snippet=''):
        key = (src, tgt, etype)
        if key in seen_typed_edge: return
        seen_typed_edge.add(key)
        edges_out.append(_graph_edge_payload(
            src, tgt, edge_type=etype, confidence=conf, source_text=snippet))

    # Authoritative agent registry (hq-memory/agents.json). Lookup by either
    # id or name. When a task/mission/receipt references an agent, we use
    # this to get the friendly name + role + color so the graph node is
    # readable ('Selvin · Code Gremlin') instead of opaque ('a_nykw53').
    auth_agents = _load_authoritative_agents()

    def _resolve_agent_meta(agent_ref: str):
        """Return the canonical agent dict for a given id-or-name reference,
        or None if the registry doesn't know about it. Case-insensitive."""
        if not agent_ref: return None
        return auth_agents.get(str(agent_ref).lower())

    # Track agent identities we've seen — keyed by canonical slug. We unify
    # references so 'Selvin', 'a_nykw53', and 'selvin' all collapse to one
    # node. Slug derives from the canonical id when known, else from the
    # raw reference (so unknown agents still get a node, just unbranded).
    agents_by_slug = {}  # slug → {'id', 'title', 'role', 'color', 'sourceIds', 'roles', 'meta'}

    def note_agent(agent_id_or_name: str, role: str = ''):
        """Register an agent reference; returns the slug node id. If the
        registry knows this agent, all later refs (by id or by name) collapse
        to the same slug — which is critical because tasks reference agents
        by id (a_xxx) but receipts reference them by name (Selvin)."""
        if not agent_id_or_name:
            return None
        meta = _resolve_agent_meta(agent_id_or_name)
        # Canonical slug: from the agent's own id if registry knows them,
        # else from the literal reference. This is what unifies cross-refs.
        canonical = (meta.get('id') if meta else str(agent_id_or_name))
        slug = _slug_agent(canonical)
        rec = agents_by_slug.get(slug)
        if not rec:
            rec = {
                'id': f'agent:{slug}',
                'title': str(agent_id_or_name),
                'role': '',
                'color': None,
                'meta': None,
                'sourceIds': set(),
                'roles': set(),
            }
            agents_by_slug[slug] = rec
        # Apply registry data when available — friendly name wins over a_xxx.
        if meta:
            rec['title'] = meta.get('name') or rec['title']
            rec['role']  = meta.get('role') or rec['role']
            rec['color'] = meta.get('color') or rec['color']
            rec['meta']  = meta
        else:
            # Fallback heuristic: prefer non-`a_xxx` displays when no registry hit.
            new = str(agent_id_or_name)
            if rec['title'].startswith('a_') and not new.startswith('a_'):
                rec['title'] = new
        rec['sourceIds'].add(str(agent_id_or_name))
        if role: rec['roles'].add(role)
        return rec['id']

    def resolve_wikilinks_in(text: str):
        """Yield resolved vault paths for every wikilink in `text`."""
        if not text: return
        for m in _WIKILINK_RE.finditer(text):
            tgt = _norm_link(m.group(1), all_paths)
            if tgt: yield tgt

    def resolve_path_to_note(p: str) -> str:
        """If `p` looks like a path that ends in a markdown filename we know,
        return the resolved vault path. Used to wire receipts whose title is
        e.g. 'FILE_WRITE: graphview_design_doc.md' to the actual note."""
        if not p: return ''
        # Try the literal stem.
        stem = pathlib.PurePosixPath(p.replace('\\', '/')).stem
        return all_paths['by_stem'].get(stem.lower(), '')

    # ── Tasks ──────────────────────────────────────────────────────────
    tasks = _load_hq_state('tasks') or []
    if isinstance(tasks, list):
        for t in tasks:
            if not isinstance(t, dict): continue
            tid = t.get('id') or ''
            if not tid: continue
            node_id = f'task:{tid}'
            title = (t.get('title') or '(untitled task)')[:120]
            detail = t.get('detail') or ''
            status = t.get('status') or ''
            priority = t.get('priority') or ''
            assigned = t.get('assignedTo')
            created = t.get('createdAt') or 0
            tags = [f'status/{status}'] if status else []
            if priority: tags.append(f'priority/{priority}')
            nodes_out.append({
                'id': node_id, 'title': title, 'path': '',
                'mtime': int(created) if created else 0,
                'size': len(detail) + len(title),
                'tags': tags, 'outlinks': 0,
                'type': 'task',
                'source': 'hq-state',
            })
            # Edges
            if assigned:
                aid = note_agent(assigned, 'assignee')
                if aid: add_edge(node_id, aid, 'assigned_to', 1.0,
                                  f'task assigned to {assigned}')
            for ref in resolve_wikilinks_in(title + '\n' + detail):
                add_edge(node_id, ref, 'references', 0.9,
                          f'task body mentions [[{pathlib.PurePosixPath(ref).stem}]]')

    # ── Projects ───────────────────────────────────────────────────────
    projects = _load_hq_state('projects') or []
    if isinstance(projects, list):
        for p in projects:
            if not isinstance(p, dict): continue
            pid = p.get('id') or ''
            if not pid: continue
            node_id = f'project:{pid}'
            name = p.get('name') or '(unnamed project)'
            ptype = p.get('type') or 'local'
            agentIds = p.get('agentIds') or []
            created = p.get('createdAt') or 0
            tags = [f'project-type/{ptype}']
            nodes_out.append({
                'id': node_id, 'title': name, 'path': p.get('path') or '',
                'mtime': int(created) if created else 0,
                'size': 0, 'tags': tags, 'outlinks': 0,
                'type': 'project',
                'source': 'hq-state',
            })
            for aid_raw in agentIds:
                aid = note_agent(aid_raw, 'project-member')
                if aid: add_edge(node_id, aid, 'assigned_to', 1.0,
                                  f'project includes agent {aid_raw}')

    # ── Missions ───────────────────────────────────────────────────────
    missions = _load_hq_state('missions') or []
    if isinstance(missions, list):
        for m in missions:
            if not isinstance(m, dict): continue
            mid = m.get('id') or ''
            if not mid: continue
            node_id = f'mission:{mid}'
            mtype = m.get('type') or 'mission'
            agent = m.get('agentId')
            topic = m.get('topic') or ''
            folder = m.get('vaultFolder') or ''
            started = m.get('startedAt') or 0
            status = m.get('status') or ''
            iters = m.get('iterations') or 0
            written = m.get('notesWritten') or []
            tags = [f'mission-type/{mtype}']
            if status: tags.append(f'status/{status}')
            short_topic = topic.split('.')[0][:80] if topic else ''
            title = f'{mtype.title()}: {short_topic}' if short_topic else f'{mtype.title()} mission'
            nodes_out.append({
                'id': node_id, 'title': title, 'path': '',
                'mtime': int(started) if started else 0,
                'size': len(topic), 'tags': tags, 'outlinks': 0,
                'type': 'mission',
                'source': 'hq-state',
            })
            if agent:
                aid = note_agent(agent, 'mission-runner')
                if aid: add_edge(node_id, aid, 'runs_as', 1.0,
                                  f'mission runs as {agent}')
            # Mission topic mentioning [[wikilinks]] → references
            for ref in resolve_wikilinks_in(topic):
                add_edge(node_id, ref, 'references', 0.9, 'mission topic')
            # `notesWritten` array → produces edges to those vault notes.
            if isinstance(written, list):
                for w in written[:50]:
                    rel = ''
                    if isinstance(w, str):
                        rel = w.replace('\\','/').lstrip('/')
                        if rel and rel not in all_paths['by_path']:
                            # Try stem fallback.
                            rel = resolve_path_to_note(rel)
                    elif isinstance(w, dict):
                        rel = (w.get('path') or '').replace('\\','/').lstrip('/')
                        if rel and rel not in all_paths['by_path']:
                            rel = resolve_path_to_note(rel)
                    if rel:
                        add_edge(node_id, rel, 'produces', 0.95,
                                  f'mission wrote {rel}')
            # If vaultFolder maps to a known folder root, edge to the index/MOC note in it.
            if folder:
                fp = folder.replace('\\','/').strip('/')
                # Look for an index-ish note inside that folder.
                for cand in (f'{fp}/index.md', f'{fp}/README.md', f'{fp}/{pathlib.PurePosixPath(fp).name}.md'):
                    if cand in all_paths['by_path']:
                        add_edge(node_id, cand, 'targets', 0.85, f'mission folder: {fp}')
                        break

    # ── Receipts ───────────────────────────────────────────────────────
    receipts = _load_hq_state('receipts') or []
    if isinstance(receipts, list):
        for r in receipts:
            if not isinstance(r, dict): continue
            rid = r.get('id') or ''
            if not rid: continue
            node_id = f'receipt:{rid}'
            title = (r.get('title') or '(untitled receipt)')[:120]
            kind = r.get('kind') or ''
            decision = r.get('decision') or ''
            decided_at = r.get('decidedAt') or 0
            elevated = r.get('elevated', False)
            by_agent = r.get('by') or ''
            tags = []
            if kind: tags.append(f'receipt-kind/{kind}')
            if decision: tags.append(f'decision/{decision}')
            if elevated: tags.append('elevated')
            # Use type 'decision' for explicit decisions, 'receipt' for tool-execution audit.
            ntype = 'decision' if decision in ('approved', 'rejected') else 'receipt'
            nodes_out.append({
                'id': node_id, 'title': title, 'path': '',
                'mtime': int(decided_at) if decided_at else 0,
                'size': 0, 'tags': tags, 'outlinks': 0,
                'type': ntype,
                'source': 'hq-state',
            })
            if by_agent:
                aid = note_agent(by_agent, 'receipt-actor')
                if aid:
                    etype = 'decided' if ntype == 'decision' else 'created_by'
                    add_edge(node_id, aid, etype, 1.0, f'receipt by {by_agent}')
            # If the receipt title looks like a file action, link to the affected note.
            # e.g. "FILE_WRITE: foo.md" or "EDIT: bar/baz.md".
            m = _re.search(r':\s*(\S+\.md)\b', title)
            if m:
                cand = resolve_path_to_note(m.group(1))
                if cand:
                    affect_type = {
                        'tool-execution': 'modified',
                    }.get(kind, 'references')
                    add_edge(node_id, cand, affect_type, 0.85, title)

    # Add every registry agent up front so even agents who haven't been
    # referenced by any task/mission/receipt still show up in the graph
    # (otherwise newly-hired agents are invisible until they do work).
    # We also seed `reports_to` edges here so the org-chart structure is
    # visible even for assistants who haven't done work yet.
    pending_reports_to = []  # [(assistant_aid_ref, senior_aid_ref)]
    for key, meta in auth_agents.items():
        # Skip duplicate entries — registry indexes by both id and name.
        aid = meta.get('id')
        if not aid: continue
        if key.lower() != aid.lower(): continue
        note_agent(aid)
        # If this agent has a senior (reportsTo set in agents.json), record
        # it for an edge — we resolve both ends to canonical agent slugs
        # AFTER all agents are noted (so the senior exists in the graph
        # regardless of registration order).
        senior_id = meta.get('reportsTo') or meta.get('parentAgentId')
        if senior_id:
            pending_reports_to.append((aid, senior_id, meta.get('name') or aid))
    # Resolve pending reports_to edges now that all agents are noted.
    for assistant_id, senior_id, assistant_name in pending_reports_to:
        a_node = note_agent(assistant_id, 'assistant')
        s_node = note_agent(senior_id, 'senior')
        if a_node and s_node and a_node != s_node:
            add_edge(a_node, s_node, 'reports_to', 1.0,
                     f'{assistant_name} reports to senior')

    # ── Messages: agent-comms registry (Phase 1 of comms refactor) ─────
    # The frontend's MessageRegistry persists every agent↔agent handoff to
    # hq-state/messages.json. Each thread becomes a `message-thread` node
    # in the graph with edges to the participating agents — so the boss can
    # see live conversations as visual structure, not just inbox rows.
    # Per-message nodes would explode the graph on busy days; per-thread
    # gives one node per coherent conversation which is far more readable.
    messages = _load_hq_state('messages') or []
    if isinstance(messages, list) and messages:
        # Group messages into threads (preserving creation order).
        threads = {}
        for m in messages:
            if not isinstance(m, dict): continue
            tid = m.get('threadId') or ''
            if not tid: continue
            threads.setdefault(tid, []).append(m)
        for tid, thread_msgs in threads.items():
            thread_msgs.sort(key=lambda m: m.get('createdAt') or 0)
            head = thread_msgs[0]
            tail = thread_msgs[-1]
            node_id = f'thread:{tid}'
            # Display: "Selvin ↔ Plato (3)" or "boss → Selvin"
            from_n = head.get('fromAgentName') or head.get('fromAgentId') or '?'
            to_n   = head.get('toAgentName')   or head.get('toAgentId')   or '?'
            count = len(thread_msgs)
            title = f'{from_n} → {to_n}'
            if count > 1: title += f' ({count})'
            # Body preview = first message body, truncated.
            body_preview = (head.get('body') or '').split('\n')[0][:90]
            # State: take the latest non-terminal state if any (so an
            # otherwise-completed thread with a new in-progress reply
            # surfaces as in_progress); else take the tail's terminal state.
            terminal_states = {'completed', 'failed', 'cancelled'}
            current_state = tail.get('state') or 'queued'
            for m in reversed(thread_msgs):
                if m.get('state') and m['state'] not in terminal_states:
                    current_state = m['state']; break
            tags = [f'state/{current_state}']
            # Promote priority/taskType to tags for filtering.
            if head.get('priority') and head['priority'] != 'med':
                tags.append(f'priority/{head["priority"]}')
            if head.get('taskType'):
                tags.append(f'task-type/{head["taskType"]}')
            # Aggregate failure cause across thread (any failed msg => tagged).
            had_failure = any(m.get('state') == 'failed' for m in thread_msgs)
            if had_failure: tags.append('had-failure')
            updated = max((m.get('updatedAt') or m.get('createdAt') or 0) for m in thread_msgs)
            nodes_out.append({
                'id': node_id,
                'title': title,
                'path': '',
                'mtime': int(updated) if updated else 0,
                'size': sum(len(m.get('body') or '') for m in thread_msgs),
                'tags': tags,
                'outlinks': 0,
                'type': 'message-thread',
                'source': 'hq-state',
                'messageCount': count,
                'threadState': current_state,
                'preview': body_preview,
            })
            # Edges: thread → each unique participant agent. Use 'sent_to'
            # for the from→thread direction and 'received' for to→thread,
            # so the graph reads as "agent originates / agent receives this
            # conversation". Boss participation collapses into a 'boss'
            # synthetic agent so it's visible in the graph too.
            participants = set()
            for m in thread_msgs:
                fid = m.get('fromAgentId') or ''
                tid_ag = m.get('toAgentId') or ''
                if fid and fid != 'boss': participants.add((fid, m.get('fromAgentName') or fid, 'sender'))
                elif fid == 'boss': participants.add(('boss', 'You (boss)', 'sender'))
                if tid_ag: participants.add((tid_ag, m.get('toAgentName') or tid_ag, 'receiver'))
            for ag_id, ag_name, role in participants:
                aid = note_agent(ag_id, role)
                if aid:
                    etype = 'sent_to' if role == 'sender' else 'received'
                    add_edge(aid, node_id, etype, 1.0,
                             f'{ag_name} {role} on this thread')
            # Edge to the artifact notes the thread produced (vault paths
            # in any message's `artifacts` array).
            for m in thread_msgs:
                for art in (m.get('artifacts') or []):
                    if not isinstance(art, dict): continue
                    p = (art.get('path') or '').replace('\\', '/').lstrip('/')
                    if not p: continue
                    if p in all_paths['by_path']:
                        kind = art.get('kind') or 'produced'
                        etype = 'produces' if kind in ('wrote', 'produced') else 'modified'
                        add_edge(node_id, p, etype, 0.95,
                                 f'thread artifact ({kind})')

    # ── Synthesise agent nodes from everything that referenced them ────
    for slug, rec in agents_by_slug.items():
        meta = rec.get('meta') or {}
        roles = sorted(rec['roles'])
        # Display title: 'Selvin · Code Gremlin' when role is known, else just name.
        display = rec['title']
        if rec.get('role'):
            display = f'{display} · {rec["role"]}'
        # Tags surface role + activity context for filter syntax.
        tags = []
        if rec.get('role'): tags.append(f'role/{_re.sub(r"[^A-Za-z0-9]+", "-", rec["role"]).strip("-").lower()}')
        for r in roles: tags.append(f'activity/{r}')
        if meta.get('elevated'): tags.append('elevated')
        if meta.get('model'): tags.append(f'model/{meta["model"].split(":")[0]}')
        # mtime: agent's lastRun if available, else hiredAt — for time-scrub
        # correctness. Both fields may legitimately be strings like
        # "just hired" or "never" (they're free-form display labels in some
        # CafresoHQ states), so int() must be defensive — fall back to 0
        # whenever the field isn't a numeric timestamp.
        def _safe_ts(v):
            try: return int(v)
            except (TypeError, ValueError): return 0
        mtime = _safe_ts(meta.get('lastRun')) or _safe_ts(meta.get('hiredAt')) or 0
        node = {
            'id': rec['id'],
            'title': display,
            'path': '',
            'mtime': mtime,
            'size': 0,
            'tags': tags,
            'outlinks': 0,
            'type': 'agent',
            'source': 'hq-state',
        }
        # Attach color/meta hints the frontend can use for sprite rendering
        # without breaking the node schema (extra fields are passed through).
        if rec.get('color'): node['agentColor'] = rec['color']
        if meta.get('id'):   node['agentId']    = meta['id']
        if meta.get('model'): node['agentModel'] = meta['model']
        if meta.get('status'): node['agentStatus'] = meta['status']
        nodes_out.append(node)

    return nodes_out, edges_out


def _build_graph_rest() -> dict:
    """Use Dataview via the REST plugin to extract links + tags in one query."""
    query = (
        'TABLE WITHOUT ID '
        'file.path AS path, file.name AS title, file.size AS size, '
        'file.mtime AS mtime, file.outlinks AS outlinks, file.tags AS tags '
        'FROM "" SORT file.mtime DESC'
    )
    body = json.dumps({'query': query, 'queryType': 'dataview'}).encode('utf-8')
    s, _, raw = _obsidian_request('POST', '/search/', body=body,
                                  content_type='application/json')
    if s != 200:
        # Fall back to FS extraction if Dataview isn't installed.
        if _vault_root:
            return _build_graph_fs()
        raise ValueError(f'dataview query failed (status {s}); install Dataview plugin or set vault directory for FS fallback')
    try:
        data = json.loads(raw.decode('utf-8'))
    except Exception as e:
        raise ValueError(f'invalid graph response: {e}')
    rows = data if isinstance(data, list) else data.get('values', []) or data.get('rows', [])
    by_path = {}
    by_stem = {}
    raw_rows = []
    for r in rows:
        # Dataview returns objects keyed by query column names, or arrays.
        if isinstance(r, dict):
            path = r.get('path') or ''
            title = r.get('title') or pathlib.PurePosixPath(path).stem
            size = r.get('size') or 0
            mtime = r.get('mtime') or 0
            outs = r.get('outlinks') or []
            tags = r.get('tags') or []
        elif isinstance(r, list) and len(r) >= 6:
            path, title, size, mtime, outs, tags = r[:6]
        else:
            continue
        if not path: continue
        by_path[path] = True
        stem = pathlib.PurePosixPath(path).stem.lower()
        by_stem[stem] = path
        raw_rows.append((path, title, size, mtime, outs or [], tags or []))
    all_paths = {'by_path': by_path, 'by_stem': by_stem}
    nodes, edges, seen_edge = [], [], set()
    for path, title, size, mtime, outs, tags in raw_rows:
        out_paths = []
        for o in outs:
            tgt = ''
            if isinstance(o, dict):
                tgt = o.get('path') or o.get('display') or ''
            elif isinstance(o, str):
                tgt = o
            tgt = _norm_link(tgt, all_paths)
            if tgt and tgt != path:
                out_paths.append(tgt)
        clean_tags = []
        for t in tags:
            if isinstance(t, str):
                clean_tags.append(t.lstrip('#'))
        for tgt in out_paths:
            key = (path, tgt)
            if key in seen_edge: continue
            seen_edge.add(key)
            edges.append(_graph_edge_payload(path, tgt))
        # mtime from Dataview is an ISO string or a luxon DateTime serialized object.
        if isinstance(mtime, str):
            try:
                from datetime import datetime
                mtime_ms = int(datetime.fromisoformat(mtime.replace('Z','+00:00')).timestamp() * 1000)
            except Exception:
                mtime_ms = 0
        elif isinstance(mtime, (int, float)):
            mtime_ms = int(mtime)
        else:
            mtime_ms = 0
        nodes.append({
            'id': path, 'title': title or pathlib.PurePosixPath(path).stem, 'path': path,
            'mtime': mtime_ms, 'size': size or 0,
            'tags': clean_tags[:8], 'outlinks': len(out_paths),
            'type': _graph_node_type(path, clean_tags),
            'source': 'markdownvault',
        })
    indeg = {n['id']: 0 for n in nodes}
    for e in edges:
        if e['target'] in indeg:
            indeg[e['target']] += 1
    for n in nodes:
        n['inlinks'] = indeg[n['id']]
    return {'nodes': nodes, 'edges': edges}


# ────────────────────────────────────────────────────────────────────────────
# Minimal Chrome DevTools Protocol WebSocket client.
#
# Just enough to drive a fresh tab through navigate + screenshot using only
# stdlib (`socket` + `struct`). Replaces a full `websockets` dependency for
# what's a single short-lived RPC dance per screenshot.
#
# Spec subset:
#   - Client always masks frames (RFC 6455).
#   - Server frames may be unmasked (and usually are for browsers).
#   - Only TEXT (opcode 0x1) frames are handled — CDP uses JSON text frames.
#   - Single-frame messages only (we never see fragmented CDP messages
#     in practice for the tiny payloads we exchange).
# ────────────────────────────────────────────────────────────────────────────
_CDP_GUID = '258EAFA5-E914-47DA-95CA-C5AB0DC85B11'

class _CDPWebSocket:
    def __init__(self, ws_url, timeout=20):
        u = urllib.parse.urlparse(ws_url)
        self.host = u.hostname
        self.port = u.port or (443 if u.scheme == 'wss' else 80)
        self.path = u.path + (('?' + u.query) if u.query else '')
        self.use_tls = (u.scheme == 'wss')
        self.timeout = timeout
        self.sock = None
        self._next_id = 1
        self._lock = threading.Lock()

    def connect(self):
        raw = socket.create_connection((self.host, self.port), timeout=self.timeout)
        if self.use_tls:
            ctx = ssl.create_default_context()
            raw = ctx.wrap_socket(raw, server_hostname=self.host)
        self.sock = raw
        # Send the upgrade request.
        key = base64.b64encode(os.urandom(16)).decode()
        req = (
            f'GET {self.path} HTTP/1.1\r\n'
            f'Host: {self.host}:{self.port}\r\n'
            f'Upgrade: websocket\r\n'
            f'Connection: Upgrade\r\n'
            f'Sec-WebSocket-Key: {key}\r\n'
            f'Sec-WebSocket-Version: 13\r\n\r\n'
        )
        self.sock.sendall(req.encode())
        # Read the handshake response.
        buf = b''
        while b'\r\n\r\n' not in buf:
            chunk = self.sock.recv(4096)
            if not chunk:
                raise RuntimeError('CDP handshake closed')
            buf += chunk
            if len(buf) > 8192:
                raise RuntimeError('CDP handshake oversized')
        first_line = buf.split(b'\r\n', 1)[0].decode('latin-1', 'replace')
        if '101' not in first_line:
            raise RuntimeError(f'CDP upgrade rejected: {first_line}')
        # Discard any payload bytes after the handshake (shouldn't happen
        # since we haven't sent any RPC yet).
        return self

    def _send_frame(self, payload_bytes):
        """Send a single TEXT frame, masked (per RFC 6455 client rules)."""
        opcode = 0x1                  # text
        first = 0x80 | opcode          # FIN=1, opcode=text
        length = len(payload_bytes)
        mask_key = os.urandom(4)
        if length < 126:
            header = struct.pack('!BB', first, 0x80 | length)
        elif length < 65536:
            header = struct.pack('!BBH', first, 0x80 | 126, length)
        else:
            header = struct.pack('!BBQ', first, 0x80 | 127, length)
        masked = bytearray(payload_bytes)
        for i in range(length):
            masked[i] ^= mask_key[i % 4]
        with self._lock:
            self.sock.sendall(header + mask_key + bytes(masked))

    def _read_exact(self, n):
        buf = b''
        while len(buf) < n:
            chunk = self.sock.recv(n - len(buf))
            if not chunk:
                raise RuntimeError('CDP socket closed mid-frame')
            buf += chunk
        return buf

    def _recv_frame(self):
        b1, b2 = struct.unpack('!BB', self._read_exact(2))
        # We only handle single-frame TEXT messages (most CDP responses are).
        opcode = b1 & 0x0F
        masked = (b2 & 0x80) != 0
        length = b2 & 0x7F
        if length == 126:
            (length,) = struct.unpack('!H', self._read_exact(2))
        elif length == 127:
            (length,) = struct.unpack('!Q', self._read_exact(8))
        mask_key = self._read_exact(4) if masked else None
        payload = self._read_exact(length)
        if masked:
            payload = bytes(b ^ mask_key[i % 4] for i, b in enumerate(payload))
        if opcode == 0x8:    # close
            raise RuntimeError('CDP server closed connection')
        if opcode in (0x9, 0xA):  # ping/pong — skip
            return self._recv_frame()
        return payload.decode('utf-8', 'replace')

    def call(self, method, params, timeout=None):
        """Send an RPC and wait for the matching response."""
        rid = self._next_id
        self._next_id += 1
        msg = json.dumps({'id': rid, 'method': method, 'params': params})
        self._send_frame(msg.encode('utf-8'))
        deadline = time.time() + (timeout or self.timeout)
        while time.time() < deadline:
            text = self._recv_frame()
            try: ev = json.loads(text)
            except Exception: continue
            if ev.get('id') == rid:
                if 'error' in ev:
                    raise RuntimeError(f'CDP error in {method}: {ev["error"]}')
                return ev
            # Ignore async events (Page.loadEventFired etc.) until our id arrives
        raise TimeoutError(f'CDP {method} timed out after {timeout or self.timeout}s')

    def wait_for_event(self, method, timeout=10):
        deadline = time.time() + timeout
        while time.time() < deadline:
            text = self._recv_frame()
            try: ev = json.loads(text)
            except Exception: continue
            if ev.get('method') == method:
                return ev
        raise TimeoutError(f'CDP event {method} did not fire within {timeout}s')

    def close(self):
        try:
            # Send close frame (opcode 0x8) — empty payload, masked.
            mask = os.urandom(4)
            self.sock.sendall(struct.pack('!BB', 0x88, 0x80) + mask)
        except Exception:
            pass
        try: self.sock.close()
        except Exception: pass


class Handler(http.server.SimpleHTTPRequestHandler):
    # HTTP/1.0 + Connection: close gives us read-until-close streaming,
    # which is exactly what SSE needs without fighting chunked encoding.
    protocol_version = 'HTTP/1.0'

    # RFC 6455 WebSocket handshake magic GUID — concatenated with the client's
    # Sec-WebSocket-Key and SHA-1'd to form Sec-WebSocket-Accept. Used by the
    # /terminal/pty upgrade handler.
    _WS_GUID = '258EAFA5-E914-47DA-95CA-C5AB0DC85B11'

    def end_headers(self):
        # Disable caching for local-dev iteration so HTML/JSX edits land on
        # next reload without browser caching tripping us up.
        self.send_header('Cache-Control', 'no-store, must-revalidate')
        self.send_header('Pragma', 'no-cache')
        self.send_header('Expires', '0')
        # CORS — only ALLOWLISTED app origins (the SvelteKit frontend + canister
        # UI shell, see _app_origins) get credentialed cross-origin access. Any
        # other origin gets a credential-less '*' (a browser cannot send cookies
        # with ACAO:'*'), so a malicious site can't read this container's
        # cookie-authenticated responses cross-origin (the hq_session is
        # SameSite=None, so it WOULD otherwise ride along). Public probes like
        # /health still work for anyone.
        try:
            origin  = self.headers.get('Origin', '') if hasattr(self, 'headers') else ''
            allowed = bool(origin) and origin in self._app_origins()
        except Exception:
            origin, allowed = '', False
        if allowed:
            self.send_header('Access-Control-Allow-Origin', origin)
            self.send_header('Access-Control-Allow-Credentials', 'true')
        else:
            self.send_header('Access-Control-Allow-Origin', '*')
        self.send_header('Access-Control-Allow-Methods', 'GET, POST, PUT, DELETE, OPTIONS')
        self.send_header('Access-Control-Allow-Headers',
                         'Content-Type, Authorization, X-User-Principal, X-API-Key, '
                         'anthropic-version, x-api-key, X-Vault-Format')
        self.send_header('Access-Control-Max-Age', '86400')
        self.send_header('Vary', 'Origin')
        super().end_headers()

    def _route(self):
        for prefix, target in ROUTES.items():
            if self.path.startswith(prefix):
                return prefix, target
        return None, None

    def _api_key_ok(self):
        """Bearer-key gate (see CAFRESOHQ_API_KEY). No key configured → open.
        Otherwise a request to a protected prefix must present the key via the
        X-API-Key header or a ?k= query param (for WebSocket handshakes). Public
        paths (static UI, /health, /idle) are exempt."""
        if not CAFRESOHQ_API_KEY:
            return True
        path = self.path.split('?', 1)[0]
        if not path.startswith(_KEY_PROTECTED_PREFIXES):
            return True
        supplied = self.headers.get('X-API-Key', '') or ''
        if not supplied:
            try:
                supplied = urllib.parse.parse_qs(
                    urllib.parse.urlparse(self.path).query).get('k', [''])[0]
            except Exception:
                supplied = ''
        import hmac as _hmac
        return bool(supplied) and _hmac.compare_digest(str(supplied), CAFRESOHQ_API_KEY)

    def do_GET(self):
        if not self._api_key_ok():
            return self._send_json(401, {'error': 'API key required'})
        # /idle is the signal the fleet's reap-idle uses to stop idle containers
        # (free the A1 pool + pause billing). Report seconds since the last
        # *user-facing* request. Health/idle pings themselves don't count as
        # activity (they'd keep a container "busy" forever).
        if self.path == '/idle':
            import time as _t
            return self._send_json(200, {'idle_seconds': int(_t.time() - _LAST_ACTIVITY[0])})
        _touch_activity(self.path)
        # Track REAL browser traffic separately (the night runner self-calls
        # tag themselves) — _night_post_activity defers shared-file writes
        # while a session is live.
        # Only genuine browser traffic marks a live UI session. /health is
        # excluded: load-balancer probes (docker-compose curls it every 30s)
        # would otherwise keep _night_browser_active() true forever and
        # permanently suppress the night-shift ticker entry.
        if 'X-Night-Runner' not in self.headers and self.path.startswith(('/hq/', '/vault/')):
            _LAST_UI_ACTIVITY[0] = time.time()
        if self.path == '/health':
            return self._health()
        if self.path == '/market/quotes':
            return self._market_quotes()
        if self.path.startswith('/hq/'):
            return self._hq_handler('GET')
        if self.path == '/missions/scheduled':
            return self._missions_scheduled_get()
        if self.path == '/missions/runs':
            return self._missions_runs()
        if self.path.startswith('/brave/'):
            return self._brave_search()
        if self.path == '/hermes/capability':
            return self._hermes_get_capability()
        if self.path == '/hermes/model':
            return self._hermes_get_model()
        if self.path == '/hermes/provider':
            return self._hermes_get_provider()
        if self.path.split('?')[0] == '/hermes/local-models':
            return self._hermes_local_models(
                self.path.split('?', 1)[1] if '?' in self.path else '')
        if self.path == '/hermes/config/export':
            return self._hermes_export_config()
        if self.path == '/hermes/trial-status':
            return self._hermes_trial_status()
        if self.path.startswith('/hermes/'):
            return self._hermes_proxy('GET')
        if self.path.startswith('/vault/'):
            return self._vault('GET')
        if self.path == '/claudecode/status':
            return self._claudecode_status()
        if self.path == '/cafresohq/status':
            return self._cafresohq_status()
        if self.path == '/codex/status':
            return self._codex_status()
        if self.path.startswith('/agents/install/status'):
            return self._agents_install_status()
        if self.path == '/agents':
            return self._agents_status()
        if self.path == '/terminal/status':
            return self._terminal_status()
        if self.path == '/terminal/nonce':
            return self._terminal_nonce()
        if self.path.startswith('/terminal/spawn'):
            return self._terminal_spawn()
        if self.path.startswith('/terminal/pty'):
            return self._terminal_pty_ws()
        if self.path.startswith('/fs/browse'):
            return self._fs_browse()
        if self.path.startswith('/fs/collect'):
            return self._fs_collect()
        if self.path.startswith('/fs/site/'):
            return self._fs_site()
        if self.path.startswith('/fs/stat'):
            return self._fs_stat()
        if self.path.startswith('/fs/file'):
            return self._fs_file()
        if self.path == '/browser/status':
            return self._browser_status()
        if self.path.startswith('/browser/fetch'):
            return self._browser_fetch()
        if self.path.startswith('/browser/screenshot'):
            return self._browser_screenshot()
        if self.path.startswith('/approvals/external/wait'):
            return self._approval_wait()
        if self.path == '/approvals/external/list':
            return self._approval_list()
        if self.path.startswith('/bundle/'):
            return self._serve_bundle()
        _p0 = self.path.split('?', 1)[0]
        if _p0 in ('/', '/hq.html', '/index.html'):
            return self._serve_hq_html()
        if _p0 == '/graph-viewer.html':
            return self._serve_cwd_file('graph-viewer.html', 'text/html; charset=utf-8')
        if _p0 == '/graph-viewer.js':
            return self._serve_dist_file('graph-viewer.js', 'application/javascript; charset=utf-8')
        if _p0.startswith('/graph/snapshot/'):
            return self._graph_snapshot(_p0.rsplit('/', 1)[-1])
        prefix, target = self._route()
        if target:
            return self._proxy('GET', prefix, target)
        # Static fallthrough serves legit UI assets (styles.css, sw.js,
        # manifest, assets/*, dist-ui) from cwd — but SimpleHTTPRequestHandler
        # would otherwise hand out ANYTHING under the repo root, incl. the TLS
        # private key + vault under hq-state/ and the .py source. Gate it.
        if not self._static_path_allowed(_p0):
            return self.send_error(404, 'Not found')
        return super().do_GET()

    # Static assets that legitimately live at the web root; everything else
    # under cwd (state dir, TLS key, vault, source) is off-limits to the
    # SimpleHTTPRequestHandler fallthrough regardless of the API key.
    def _static_path_allowed(self, url_path):
        try:
            rel = urllib.parse.unquote(url_path.split('?', 1)[0]).lstrip('/')
            root = os.path.realpath(os.getcwd())
            resolved = os.path.realpath(os.path.join(root, rel))
            # (a) must stay inside the web root
            if resolved != root and not resolved.startswith(root + os.sep):
                return False
            # (b) never serve Python source
            if resolved.endswith('.py'):
                return False
            # (c) never serve the state dir (tls/, vault/, memory/, *.json)
            state_root = os.path.realpath(str(_hq_state_dir))
            if resolved == state_root or resolved.startswith(state_root + os.sep):
                return False
            return True
        except Exception:
            return False

    # No directory listings — hq-state/ etc. must not be enumerable.
    def list_directory(self, path):
        self.send_error(404, 'Not found')
        return None

    # --- HQ UI build (no in-browser Babel) ------------------------------------
    # The HQ app is built into dist-ui/ by scripts/build_ui_bundle.mjs: vendor
    # globals (React/ReactDOM/xterm) + content-hashed, pre-transformed JSX. hq.html
    # carries an <!--HQ_SCRIPTS--> placeholder; we substitute it from the manifest.
    # (Mirrors scripts/ui_manifest.py — inlined so the packaged/frozen build needs
    # no scripts/ dir alongside the exe.)
    def _hq_manifest_tags(self):
        import json as _json
        with open(os.path.join(os.getcwd(), 'dist-ui', 'manifest.json'), encoding='utf-8') as fh:
            m = _json.load(fh)
        parts = []
        for css in m.get('vendorCss', []):
            parts.append('<link rel="stylesheet" href="%s"/>' % css)
        for js in m.get('vendor', []):
            parts.append('<script src="%s"></script>' % js)
        if m.get('analyticsWorker'):
            parts.append('<script>window.__CAFRESO_BUNDLE__=%s;</script>'
                         % _json.dumps({'analyticsWorker': m['analyticsWorker']}))
        if m.get('graphEngine'):
            parts.append('<script src="%s"></script>' % m['graphEngine'])
        for js in m.get('app', []):
            parts.append('<script src="%s"></script>' % js)
        return '\n'.join(parts)

    def _serve_hq_html(self):
        try:
            with open(os.path.join(os.getcwd(), 'hq.html'), encoding='utf-8') as fh:
                html = fh.read()
            html = html.replace('<!--HQ_SCRIPTS-->', self._hq_manifest_tags())
        except Exception as e:
            return self.send_error(500, 'HQ UI not built: %s (run `npm run build`)' % e)
        body = html.encode('utf-8')
        self.send_response(200)
        self.send_header('Content-Type', 'text/html; charset=utf-8')
        self.send_header('Content-Length', str(len(body)))
        self.end_headers()
        try:
            self.wfile.write(body)
        except Exception:
            pass

    def _serve_bundle(self):
        rel = self.path.split('?', 1)[0].lstrip('/')          # 'bundle/app-<hash>.js'
        base = os.path.normpath(os.path.join(os.getcwd(), 'dist-ui'))
        full = os.path.normpath(os.path.join(base, rel))
        if not full.startswith(base) or not os.path.isfile(full):
            return self.send_error(404)
        ctype = ('application/javascript' if full.endswith('.js')
                 else 'text/css' if full.endswith('.css')
                 else 'application/octet-stream')
        with open(full, 'rb') as fh:
            data = fh.read()
        self.send_response(200)
        self.send_header('Content-Type', ctype + '; charset=utf-8')
        self.send_header('Content-Length', str(len(data)))
        self.end_headers()
        try:
            self.wfile.write(data)
        except Exception:
            pass

    # --- Public shareable graphs (Phase 2) ------------------------------------
    # Publish exports a laid-out snapshot (positions + community + betweenness) and
    # stores it under hq-state/public-graphs/<slug>.json; the static graph-viewer
    # renders it read-only via ?g=/graph/snapshot/<slug>. (Production target is an
    # ICP asset/Motoko canister — same URL contract, swappable storage.)
    def _public_graph_dir(self):
        base = os.environ.get('CAFRESOHQ_HQ_STATE_DIR') or os.path.join(os.getcwd(), 'hq-state')
        d = os.path.join(base, 'public-graphs')
        os.makedirs(d, exist_ok=True)
        return d

    def _serve_cwd_file(self, name, ctype):
        try:
            with open(os.path.join(os.getcwd(), name), 'rb') as fh:
                data = fh.read()
        except Exception:
            return self.send_error(404)
        self.send_response(200)
        self.send_header('Content-Type', ctype)
        self.send_header('Content-Length', str(len(data)))
        self.end_headers()
        try: self.wfile.write(data)
        except Exception: pass

    def _serve_dist_file(self, name, ctype):
        full = os.path.normpath(os.path.join(os.getcwd(), 'dist-ui', name))
        if not os.path.isfile(full):
            return self.send_error(404)
        with open(full, 'rb') as fh:
            data = fh.read()
        self.send_response(200)
        self.send_header('Content-Type', ctype)
        self.send_header('Content-Length', str(len(data)))
        self.end_headers()
        try: self.wfile.write(data)
        except Exception: pass

    def _graph_snapshot(self, slug):
        import re as _re
        if not _re.fullmatch(r'[A-Za-z0-9_-]{1,64}', slug or ''):
            return self.send_error(404)
        full = os.path.join(self._public_graph_dir(), slug + '.json')
        if not os.path.isfile(full):
            return self.send_error(404)
        with open(full, 'rb') as fh:
            data = fh.read()
        self.send_response(200)
        self.send_header('Content-Type', 'application/json; charset=utf-8')
        self.send_header('Content-Length', str(len(data)))
        self.end_headers()
        try: self.wfile.write(data)
        except Exception: pass

    def _graph_publish(self):
        import json as _json, secrets as _secrets
        try:
            length = int(self.headers.get('Content-Length', 0))
            body = self.rfile.read(length) if length else b'{}'
            payload = _json.loads(body.decode('utf-8'))
        except Exception:
            return self._send_json(400, {'error': 'bad payload'})
        if not isinstance(payload, dict) or 'graph' not in payload:
            return self._send_json(400, {'error': 'missing graph'})
        slug = _secrets.token_hex(6)
        try:
            with open(os.path.join(self._public_graph_dir(), slug + '.json'), 'w', encoding='utf-8') as fh:
                _json.dump(payload, fh)
        except Exception as e:
            return self._send_json(500, {'error': str(e)})
        snap_url = '/graph/snapshot/' + slug
        view_url = '/graph-viewer.html?g=' + snap_url + '&background=dark&most_influential=bc&maxnodes=150&show_analytics=1&selected=highlight&demo=1'
        return self._send_json(200, {'slug': slug, 'snapshotUrl': snap_url, 'viewerUrl': view_url})

    def do_POST(self):
        if not self._api_key_ok():
            return self._send_json(401, {'error': 'API key required'})
        _touch_activity(self.path)
        if self.path == '/hermes/capability':
            return self._hermes_set_capability()
        if self.path == '/hermes/model':
            return self._hermes_set_model()
        if self.path == '/hermes/provider':
            return self._hermes_set_provider()
        if self.path == '/hermes/config/import':
            return self._hermes_import_config()
        if self.path == '/hermes/openrouter-key':
            return self._hermes_set_provider('openrouter')  # back-compat alias
        if self.path.startswith('/hermes/'):
            return self._hermes_proxy('POST')
        if self.path.startswith('/vault/'):
            return self._vault('POST')
        if self.path == '/graph/publish':
            return self._graph_publish()
        if self.path == '/agents/install':
            return self._agents_install()
        if self.path == '/claudecode/configure':
            return self._claudecode_configure()
        if self.path == '/codex/configure':
            return self._codex_configure()
        if self.path == '/claudecode/stream':
            return self._claudecode_stream()
        if self.path == '/cafresohq/stream':
            return self._cafresohq_stream()
        if self.path == '/codex/stream':
            return self._codex_stream()
        if self.path == '/terminal/stream':
            return self._terminal_stream()
        if self.path == '/projects/clone':
            return self._projects_clone()
        if self.path == '/missions/schedule':
            return self._missions_schedule()
        if self.path == '/tools/exec':
            return self._tool_exec()
        if self.path.startswith('/fs/upload'):
            return self._fs_upload()
        if self.path.startswith('/fs/mkdir'):
            return self._fs_mkdir()
        if self.path.startswith('/fs/rename'):
            return self._fs_rename()
        if self.path.startswith('/fs/delete'):
            return self._fs_delete()
        if self.path == '/approvals/external':
            return self._approval_submit()
        if self.path == '/approvals/external/decide':
            return self._approval_decide()
        if self.path == '/export/pptx':
            return self._export_pptx()
        if self.path == '/export/docx':
            return self._export_docx()
        if self.path == '/export/pdf':
            return self._export_pdf()
        if self.path == '/generate/image':
            return self._generate_image()
        if self.path == '/generate/video':
            return self._generate_video()
        prefix, target = self._route()
        if target:
            return self._proxy('POST', prefix, target)
        self.send_error(405)

    def do_PUT(self):
        if not self._api_key_ok():
            return self._send_json(401, {'error': 'API key required'})
        if self.path.startswith('/hq/'):
            return self._hq_handler('PUT')
        if self.path.startswith('/vault/'):
            return self._vault('PUT')
        self.send_error(405)

    def do_DELETE(self):
        if not self._api_key_ok():
            return self._send_json(401, {'error': 'API key required'})
        if self.path.startswith('/vault/'):
            return self._vault('DELETE')
        if self.path.startswith('/missions/scheduled/'):
            return self._missions_delete()
        self.send_error(405)

    def do_OPTIONS(self):
        prefix, target = self._route()
        if target:
            return self._proxy('OPTIONS', prefix, target)
        # Default: respond 204 No Content with CORS headers (added by end_headers).
        # This handles preflight requests from the SvelteKit frontend for /health,
        # /vault/*, /api/anthropic/*, etc. without needing per-route OPTIONS.
        self.send_response(204)
        self.send_header('content-length', '0')
        self.end_headers()

    # ---- Claude Code (Pro/Max subscription via local CLI) -----------------
    def _claudecode_resolve(self):
        """Find the claude binary. Returns absolute path or None."""
        if _claudecode_bin and pathlib.Path(_claudecode_bin).is_file():
            return _claudecode_bin
        # shutil.which respects PATH and Windows .cmd/.exe extensions.
        return shutil.which(_claudecode_bin or 'claude')

    def _market_quotes(self):
        """GET /market/quotes — cached index/gold quotes for the office ticker."""
        global _market_cache
        with _market_lock:
            if time.time() - _market_cache['ts'] < 60 and _market_cache['quotes']:
                return self._send_json(200, {'quotes': _market_cache['quotes'], 'ts': _market_cache['ts']})
        quotes = []
        for sym, yq in _MARKET_SYMBOLS:
            try:
                url = ('https://query1.finance.yahoo.com/v8/finance/chart/'
                       + urllib.parse.quote(yq) + '?range=1d&interval=1d')
                req = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0'})
                with urllib.request.urlopen(req, timeout=6) as r:
                    meta = json.load(r)['chart']['result'][0]['meta']
                last = float(meta['regularMarketPrice'])
                prev = float(meta.get('chartPreviousClose') or 0)
                quotes.append({'sym': sym, 'last': last,
                               'pct': ((last - prev) / prev * 100.0) if prev > 0 else None})
            except Exception:
                continue   # partial results are fine; total failure serves stale below
        with _market_lock:
            if quotes:
                _market_cache = {'ts': time.time(), 'quotes': quotes}
            return self._send_json(200, {'quotes': _market_cache['quotes'], 'ts': _market_cache['ts']})

    def _health(self):
        """GET /health — lightweight liveness probe for load balancers and companion detection."""
        import platform
        oci_ready = (
            bool(_oci_vault_namespace and _oci_vault_bucket)
            if _vault_backend == 'oci' else None
        )
        return self._send_json(200, {
            'status':           'ok',
            'version':          '1.0.0',
            # apiVersion: integer contract between the (canister-served) UI and
            # this API. Bump on a BREAKING endpoint change; the UI min-checks it
            # so a freshly-deployed UI degrades gracefully against an older
            # container (and vice-versa) instead of hard-failing. See API_VERSION.
            'apiVersion':       API_VERSION,
            'mode':             _fleet_mode,
            'vault_backend':    _vault_backend,
            'uptime_seconds':   int(time.time() - _server_start_time),
            'platform':         platform.system(),
            'user_principal':   _fleet_user_principal,
            'claude_code':      bool(_claudecode_bin or shutil.which('claude')),
            'codex':            bool(_codex_bin or shutil.which('codex')),
            'hermes':           bool(_hermes_bin or shutil.which('hermes')),
            'gemini':           bool(_gemini_bin or shutil.which('gemini')),
            'runtime_env':      _RUNTIME_ENV,
            'auth_required':    bool(CAFRESOHQ_API_KEY),
            'oci_vault_ready':  oci_ready,
        })

    def _claudecode_status(self):
        bin_ = self._claudecode_resolve()
        return self._send_json(200, {
            'configured': bool(bin_),
            'binary': bin_ or '',
            'override': _claudecode_bin or '',
        })

    def _claudecode_configure(self):
        global _claudecode_bin
        length = int(self.headers.get('content-length', 0) or 0)
        try:
            body = json.loads(self.rfile.read(length).decode('utf-8') or '{}')
        except json.JSONDecodeError:
            return self._send_json(400, {'error': 'bad json'})
        path = (body.get('binary') or '').strip()
        if path and not pathlib.Path(path).is_file():
            return self._send_json(400, {'error': f'not a file: {path}'})
        _claudecode_bin = path
        return self._claudecode_status()

    def _claudecode_stream(self):
        """Stream a chat completion through the local `claude` CLI.

        Body: {messages, system, model, maxTokens?}
        We compose a single prompt from system + history + last user turn,
        spawn `claude --print --output-format=stream-json --model <m> ...`,
        pipe the prompt to stdin, and parse the line-delimited JSON events
        on stdout. Each text delta is forwarded as an SSE-shaped frame
        (matching what claude-client.jsx expects from the OpenAI-compat
        backends, so we can reuse parseSSE on the client).
        """
        bin_ = self._claudecode_resolve()
        if not bin_:
            return self._send_json(503, {
                'error': 'claude CLI not found — install Claude Code or set CAFRESOHQ_CLAUDE_BIN'
            })
        length = int(self.headers.get('content-length', 0) or 0)
        try:
            body = json.loads(self.rfile.read(length).decode('utf-8') or '{}')
        except json.JSONDecodeError:
            return self._send_json(400, {'error': 'bad json'})

        messages = body.get('messages') or []
        system   = (body.get('system') or '').strip()
        model    = (body.get('model') or '').strip()

        # Compose a single prompt for non-interactive --print mode.
        # Roles get prefixed so the model sees the conversation shape; we
        # extract the last user turn to be the "current" message.
        history_lines = []
        last_user = ''
        for m in messages:
            role = (m.get('role') or 'user').upper()
            content = (m.get('content') or '').strip()
            if not content: continue
            if role == 'USER': last_user = content
            history_lines.append(f'{role}: {content}')
        prompt = '\n\n'.join(history_lines) if history_lines else last_user
        if not prompt:
            return self._send_json(400, {'error': 'no prompt'})

        # Use disallowed-tools to suppress all tool use in non-elevated mode.
        # --allowed-tools '' (empty string) is rejected by Claude CLI; using
        # --disallowed-tools with a wildcard is the correct idiom.
        cmd = [bin_, '--print',
               '--output-format', 'stream-json',
               '--verbose',
               '--disallowed-tools', 'Bash,Edit,Write,MultiEdit,NotebookEdit,WebFetch,WebSearch,TodoWrite,Task',
               '--input-format', 'text']
        if model:
            cmd += ['--model', model]
        if system:
            cmd += ['--append-system-prompt', system]

        # Use project cwd if provided and valid, otherwise first allowed dir.
        _cc_cwd = _cafresohq_allowed_dirs[0] if _cafresohq_allowed_dirs else None
        req_cwd = (body.get('cwd') or '').strip()
        if req_cwd:
            try:
                cwd_p = pathlib.Path(_client_path(req_cwd)).resolve()
                for d in _cafresohq_allowed_dirs:
                    try:
                        cwd_p.relative_to(pathlib.Path(d).resolve())
                        if cwd_p.is_dir():
                            _cc_cwd = str(cwd_p)
                        break
                    except ValueError:
                        continue
            except OSError:
                pass
        import copy as _copy
        _cc_env = _copy.deepcopy(os.environ)
        for _d in [r'C:\Program Files\Git\usr\bin',
                   r'C:\Program Files\Git\bin',
                   r'C:\Program Files\Git\mingw64\bin']:
            if os.path.isdir(_d) and _d not in _cc_env.get('PATH', ''):
                _cc_env['PATH'] = _d + os.pathsep + _cc_env.get('PATH', '')

        try:
            proc = subprocess.Popen(
                cmd,
                stdin=subprocess.PIPE,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True, bufsize=1, encoding='utf-8',
                cwd=_cc_cwd,
                env=_cc_env,
            )
        except FileNotFoundError as e:
            return self._send_json(500, {'error': f'spawn: {e}'})

        # Send the prompt and close stdin so the CLI knows we're done.
        try:
            proc.stdin.write(prompt)
            proc.stdin.close()
        except Exception as e:
            try: proc.kill()
            except Exception: pass
            return self._send_json(500, {'error': f'stdin: {e}'})

        # Stream stdout to the client as SSE-style "data: <json>\n\n" frames
        # in the OpenAI-compatible chat-completion delta shape, so the
        # existing client parser works without changes.
        self.send_response(200)
        self.send_header('content-type', 'text/event-stream')
        self.send_header('cache-control', 'no-store')
        self.end_headers()

        def write_sse(obj):
            try:
                self.wfile.write(b'data: ' + json.dumps(obj).encode('utf-8') + b'\n\n')
                self.wfile.flush()
            except (BrokenPipeError, ConnectionResetError):
                return False
            return True

        # Drain stderr in a thread so the pipe doesn't fill up and stall.
        stderr_buf = []
        def _drain_err():
            try:
                for line in proc.stderr:
                    stderr_buf.append(line)
            except Exception:
                pass
        threading.Thread(target=_drain_err, daemon=True).start()

        in_tokens = 0
        out_tokens = 0
        try:
            for line in proc.stdout:
                line = line.strip()
                if not line: continue
                try:
                    ev = json.loads(line)
                except json.JSONDecodeError:
                    continue
                # Claude Code's stream-json wraps Anthropic message events.
                # Common shapes:
                #   {"type":"system","subtype":"init", ...}
                #   {"type":"assistant","message":{"content":[{"type":"text","text":"…"}], "usage":{...}}}
                #   {"type":"result","result":"…","usage":{...}}
                t = ev.get('type')
                if t == 'assistant':
                    msg = ev.get('message') or {}
                    for block in (msg.get('content') or []):
                        if block.get('type') == 'text':
                            text = block.get('text') or ''
                            if text:
                                ok = write_sse({
                                    'choices': [{
                                        'index': 0,
                                        'delta': {'content': text},
                                    }],
                                })
                                if not ok: break
                    u = msg.get('usage')
                    if u:
                        in_tokens  = u.get('input_tokens', in_tokens)
                        out_tokens = u.get('output_tokens', out_tokens)
                elif t == 'result':
                    u = ev.get('usage') or {}
                    in_tokens  = u.get('input_tokens',  in_tokens)
                    out_tokens = u.get('output_tokens', out_tokens)
                elif t == 'error':
                    write_sse({'choices': [{'index': 0, 'delta': {
                        'content': f"\n\n⚠ Claude Code error: {ev.get('message') or ev}"}}]})
            # Final usage frame so the client can attribute tokens.
            if in_tokens or out_tokens:
                write_sse({
                    'choices': [],
                    'usage': {
                        'prompt_tokens': in_tokens,
                        'completion_tokens': out_tokens,
                        'total_tokens': in_tokens + out_tokens,
                    },
                })
        finally:
            try: proc.wait(timeout=2)
            except Exception:
                try: proc.kill()
                except Exception: pass
            # If exit was non-zero and we never wrote any text, surface stderr.
            if proc.returncode and proc.returncode != 0 and not (in_tokens or out_tokens):
                err_text = ('\n'.join(stderr_buf))[:600] or f'exit {proc.returncode}'
                write_sse({'choices': [{'index': 0, 'delta': {
                    'content': f"\n\n⚠ Claude Code exited {proc.returncode}: {err_text}"}}]})

    # ---- Tool execution proxy (bracket-format tools for any provider) ------
    def _validate_path(self, path):
        """Resolve path and verify it falls within CAFRESOHQ_ALLOWED_DIRS.
        In local mode with no explicit CAFRESOHQ_ALLOWED_DIRS env var the check
        is skipped — the user is developing locally and can access their own files.
        In container mode or when the admin explicitly set CAFRESOHQ_ALLOWED_DIRS
        the strict whitelist is enforced.
        """
        p = pathlib.Path(_client_path(path)).resolve()
        if not _ALLOWED_DIRS_EXPLICIT and _RUNTIME_ENV == 'local':
            return p  # local default: no restriction, user accesses own files
        for d in _cafresohq_allowed_dirs:
            try:
                p.relative_to(pathlib.Path(d).resolve())
                return p
            except ValueError:
                continue
        raise PermissionError(f'Path outside allowed directories: {path!r}')

    def _projects_clone(self):
        """Clone a GitHub repo into the first allowed dir.
        Body: {url, name?, depth?}
          url   — github URL or owner/repo shorthand
          name  — optional override for the local folder name
          depth — 1 (default, shallow) or 0 (full history)
        Returns: {ok, path, name, stderr?}
        """
        if not _cafresohq_allowed_dirs:
            return self._send_json(503, {'ok': False,
                'error': 'CAFRESOHQ_ALLOWED_DIRS not set — clone disabled'})
        length = int(self.headers.get('content-length', 0) or 0)
        try:
            req = json.loads(self.rfile.read(length).decode('utf-8') or '{}')
        except json.JSONDecodeError:
            return self._send_json(400, {'ok': False, 'error': 'bad json'})

        raw_url = (req.get('url') or '').strip()
        if not raw_url:
            return self._send_json(400, {'ok': False, 'error': 'url required'})

        # Normalize: 'owner/repo' → 'https://github.com/owner/repo'
        if not raw_url.startswith(('http://', 'https://', 'git@')):
            if '/' in raw_url and not raw_url.startswith('/'):
                raw_url = 'https://github.com/' + raw_url
            else:
                return self._send_json(400, {'ok': False,
                    'error': 'unrecognized url; use https://github.com/owner/repo or owner/repo'})

        # Strip trailing .git or trailing slash for name derivation
        repo_name = (req.get('name') or '').strip()
        if not repo_name:
            tail = raw_url.rstrip('/').split('/')[-1]
            if tail.endswith('.git'): tail = tail[:-4]
            repo_name = tail
        # Sanitize name: only allow safe chars
        safe = ''.join(c for c in repo_name if c.isalnum() or c in '-_.')
        if not safe:
            return self._send_json(400, {'ok': False, 'error': 'invalid repo name'})
        repo_name = safe

        target = pathlib.Path(_cafresohq_allowed_dirs[0]) / repo_name
        try:
            target = target.resolve()
            # Must remain inside the allowed dir
            target.relative_to(pathlib.Path(_cafresohq_allowed_dirs[0]).resolve())
        except (ValueError, OSError) as e:
            return self._send_json(400, {'ok': False, 'error': f'invalid target: {e}'})

        if target.exists():
            return self._send_json(409, {'ok': False,
                'error': f'destination already exists: {target}',
                'path': str(target), 'name': repo_name})

        git_bin = shutil.which('git') or shutil.which('git.exe')
        if not git_bin:
            return self._send_json(503, {'ok': False, 'error': 'git not found in PATH'})

        depth = req.get('depth', 1)
        cmd = [git_bin, 'clone']
        if depth and int(depth) > 0:
            cmd += ['--depth', str(int(depth))]
        cmd += [raw_url, str(target)]

        sys.stderr.write(f'[projects] clone {raw_url} → {target}\n')
        try:
            proc = subprocess.run(cmd, capture_output=True, text=True, timeout=300)
        except subprocess.TimeoutExpired:
            return self._send_json(504, {'ok': False, 'error': 'clone timed out (5min)'})
        except Exception as e:
            return self._send_json(500, {'ok': False, 'error': f'spawn git: {e}'})

        if proc.returncode != 0:
            err = (proc.stderr or proc.stdout or '')[:1000]
            return self._send_json(502, {'ok': False,
                'error': f'git clone failed (exit {proc.returncode})',
                'stderr': err})

        return self._send_json(200, {
            'ok': True,
            'path': str(target),
            'name': repo_name,
            'url': raw_url,
        })

    # ---- Night Shift endpoints (Sprint 4 MVP-1) -----------------------------
    def _missions_scheduled_get(self):
        return self._send_json(200, {
            'schedules': _night_load('scheduled-missions.json', []),
            'running': list(_night_running.keys()),
            'browserActive': _night_browser_active(),
        })

    def _missions_runs(self):
        return self._send_json(200, {'runs': _night_load('mission-runs.json', [])})

    def _missions_schedule(self):
        """POST /missions/schedule — create/update one night-shift schedule.
        Bounds: duration ≤ 4h, interval ≥ 60s, ≤ 20 schedules. The runner's
        tool subset (no wallet/publish/shell) is enforced in night_runner.py."""
        length = int(self.headers.get('content-length', 0) or 0)
        try:
            req = json.loads(self.rfile.read(length).decode('utf-8')) if length else {}
        except Exception:
            return self._send_json(400, {'error': 'bad json'})
        topic = str(req.get('topic', '')).strip()
        agent_id = str(req.get('agentId', '')).strip()
        if not topic or not agent_id:
            return self._send_json(400, {'error': 'topic and agentId are required'})
        try:
            start_at = int(req.get('startAt') or 0)
            duration = int(req.get('durationMs') or 3_600_000)
            interval = int(req.get('intervalMs') or 300_000)
        except Exception:
            return self._send_json(400, {'error': 'bad startAt/durationMs/intervalMs'})
        now_ms = int(time.time() * 1000)
        duration = max(60_000, min(duration, 4 * 3_600_000))
        interval = max(60_000, min(interval, duration))
        sched = {
            'id': str(req.get('id') or '').strip() or ('nsh_%d' % now_ms),
            'type': 'project-study' if req.get('type') == 'project-study' else 'research',
            'topic': topic[:200],
            'agentId': agent_id[:80],
            'agentName': str(req.get('agentName') or '').strip()[:60],
            'vaultFolder': (str(req.get('vaultFolder') or '').strip().strip('/')
                            or 'Research/night')[:120],
            'projectName': str(req.get('projectName') or '')[:80],
            'projectPath': str(req.get('projectPath') or '')[:300],
            'startAt': start_at,
            'recurrence': 'daily' if req.get('recurrence') == 'daily' else 'once',
            'durationMs': duration, 'intervalMs': interval,
            'enabled': req.get('enabled') is not False,
            'lastRunAt': 0,
            'nextRunAt': start_at if start_at > now_ms else now_ms,
            'createdAt': now_ms,
        }
        with _night_lock:
            scheds = [s for s in _night_load('scheduled-missions.json', [])
                      if s.get('id') != sched['id']]
            if len(scheds) >= MAX_NIGHT_SCHEDULES:
                return self._send_json(400, {'error': 'too many schedules (max %d)' % MAX_NIGHT_SCHEDULES})
            scheds.append(sched)
            _night_save('scheduled-missions.json', scheds)
        return self._send_json(200, {'ok': True, 'schedule': sched})

    def _missions_delete(self):
        sid = self.path.rstrip('/').rsplit('/', 1)[-1]
        with _night_lock:
            scheds = _night_load('scheduled-missions.json', [])
            kept = [s for s in scheds if s.get('id') != sid]
            _night_save('scheduled-missions.json', kept)
        return self._send_json(200, {'ok': True, 'removed': sid,
                                     'existed': len(kept) != len(scheds)})

    def _tool_exec(self):
        """Execute a bracket-format tool call dispatched by the frontend.
        Body: {tool, arg, body?}
        Requires CAFRESOHQ_ALLOWED_DIRS to be configured (same gate as /cafresohq/stream).
        BASH additionally requires 'Bash' in CAFRESOHQ_ALLOWED_TOOLS.
        """
        if not _cafresohq_allowed_dirs:
            return self._send_json(503, {'ok': False,
                'error': 'CAFRESOHQ_ALLOWED_DIRS not set — tool execution disabled'})
        length = int(self.headers.get('content-length', 0) or 0)
        try:
            req = json.loads(self.rfile.read(length).decode('utf-8') or '{}')
        except json.JSONDecodeError:
            return self._send_json(400, {'ok': False, 'error': 'bad json'})

        tool    = (req.get('tool') or '').upper()
        arg     = (req.get('arg') or '').strip()
        content = req.get('body') or ''

        # Optional project cwd — if supplied and valid (within allowed dirs),
        # use it as the base for relative paths and BASH cwd.
        req_cwd = (req.get('cwd') or '').strip()
        tool_cwd = _cafresohq_allowed_dirs[0] if _cafresohq_allowed_dirs else os.getcwd()
        if req_cwd:
            try:
                cwd_p = pathlib.Path(_client_path(req_cwd)).resolve()
                if cwd_p.is_dir():
                    if not _ALLOWED_DIRS_EXPLICIT and _RUNTIME_ENV == 'local':
                        tool_cwd = str(cwd_p)  # local default: trust any existing dir
                    else:
                        for d in _cafresohq_allowed_dirs:
                            try:
                                cwd_p.relative_to(pathlib.Path(d).resolve())
                                tool_cwd = str(cwd_p)
                                break
                            except ValueError:
                                continue
            except OSError:
                pass

        def _resolve_arg(raw):
            """Resolve a path arg — if relative, anchor to tool_cwd."""
            p = pathlib.Path(raw)
            if not p.is_absolute():
                p = pathlib.Path(tool_cwd) / p
            return self._validate_path(str(p))

        try:
            if tool == 'FILE_READ':
                p = _resolve_arg(arg)
                if not p.exists():
                    result = f'File not found: {arg}'
                else:
                    text = p.read_text(encoding='utf-8', errors='replace')
                    if len(text) > 8000:
                        text = text[:8000] + '\n…(truncated to 8000 chars)'
                    result = text

            elif tool == 'DIR_LIST':
                p = _resolve_arg(arg or tool_cwd)
                if not p.is_dir():
                    result = f'Not a directory: {arg}'
                else:
                    entries = sorted(p.iterdir(), key=lambda x: (x.is_file(), x.name.lower()))
                    lines = []
                    for e in entries[:300]:
                        lines.append(e.name + '/' if e.is_dir() else f'{e.name}  ({e.stat().st_size} B)')
                    result = '\n'.join(lines) or '(empty directory)'
                    if len(lines) == 300:
                        result += '\n…(truncated at 300 entries)'

            elif tool == 'FILE_WRITE':
                p = _resolve_arg(arg)
                p.parent.mkdir(parents=True, exist_ok=True)
                p.write_text(content, encoding='utf-8')
                result = f'Wrote {len(content)} chars → {p}'

            elif tool == 'BASH':
                if 'Bash' not in _cafresohq_allowed_tools:
                    return self._send_json(403, {'ok': False,
                        'error': 'Bash not in CAFRESOHQ_ALLOWED_TOOLS — add it to enable shell access'})
                proc = subprocess.run(
                    arg, shell=True, capture_output=True, text=True, timeout=30,
                    cwd=tool_cwd,
                )
                out = proc.stdout
                if proc.stderr:
                    out += '\nSTDERR:\n' + proc.stderr
                if len(out) > 4000:
                    out = out[:4000] + '\n…(truncated)'
                result = (out or '(no output)')
                if proc.returncode != 0:
                    result += f'\n(exit {proc.returncode})'

            else:
                return self._send_json(400, {'ok': False, 'error': f'Unknown tool: {tool}'})

            return self._send_json(200, {'ok': True, 'result': result})

        except PermissionError as e:
            return self._send_json(403, {'ok': False, 'error': str(e)})
        except subprocess.TimeoutExpired:
            return self._send_json(200, {'ok': False, 'error': 'command timed out (30s)'})
        except Exception as e:
            return self._send_json(500, {'ok': False, 'error': str(e)})

    # ---- CafresoHQ (elevated agent: Claude Code with constrained tools) ---
    def _cafresohq_status(self):
        """Report whether the elevated endpoint is wired up. The client uses
        this to disable the 🛡 toggle when no allowlist is configured."""
        bin_ = self._claudecode_resolve()
        # Validate that every configured dir exists — a typo'd dir is a
        # silent landmine (claude would just refuse the read).
        dirs_ok = []
        dirs_bad = []
        for d in _cafresohq_allowed_dirs:
            (dirs_ok if pathlib.Path(d).is_dir() else dirs_bad).append(d)
        return self._send_json(200, {
            'configured': bool(bin_) and bool(dirs_ok) and bool(_cafresohq_allowed_tools),
            'binary': bin_ or '',
            'allowedDirs': dirs_ok,
            'badDirs': dirs_bad,
            'allowedTools': list(_cafresohq_allowed_tools),
        })

    def _codex_status(self):
        """Report whether the Codex elevated endpoint is wired up."""
        bin_ = self._codex_resolve()
        dirs_ok, dirs_bad = [], []
        for d in _cafresohq_allowed_dirs:
            (dirs_ok if pathlib.Path(d).is_dir() else dirs_bad).append(d)
        return self._send_json(200, {
            'configured': bool(bin_) and bool(dirs_ok),
            'binary': bin_ or '',
            'override': _codex_bin or '',
            'allowedDirs': dirs_ok,
            'badDirs': dirs_bad,
            'allowedTools': list(_cafresohq_allowed_tools),
        })

    # ──────────────────────────────────────────────────────────────────────
    # Agents — Hermes ships built-in (default runtime); Claude Code + Codex are
    # optional and installed on-demand so the image stays lean and provisioning
    # is fast. Users add the agents they want from HQ Settings.
    # ──────────────────────────────────────────────────────────────────────
    def _agent_version(self, bin_, *args):
        if not bin_:
            return ''
        try:
            r = subprocess.run([bin_, *args], capture_output=True, text=True, timeout=6)
            out = (r.stdout or r.stderr or '').strip()
            return out.splitlines()[0][:80] if out else ''
        except Exception:
            return ''

    def _agent_auth_detect(self, aid):
        """Best-effort login/credential detection for an agent CLI on THIS host.

        Returns (authenticated: bool, auth: str) where auth is the mechanism
        found: 'oauth' (CLI login), 'api-key' (env var), 'config' (hermes
        config.yaml present), or '' when nothing was found. File checks only —
        no subprocesses, no key material is read or returned. A False here is
        a hint, not a verdict (e.g. macOS keychain-stored claude creds have no
        file to see), so the UI words it as "needs login", never "broken".
        """
        home = pathlib.Path.home()
        try:
            if aid == 'claude-code':
                if (home / '.claude' / '.credentials.json').is_file():
                    return True, 'oauth'
                cfg = home / '.claude.json'
                if cfg.is_file():
                    try:
                        if '"oauthAccount"' in cfg.read_text(encoding='utf-8', errors='ignore'):
                            return True, 'oauth'
                    except OSError:
                        pass
                if os.environ.get('ANTHROPIC_API_KEY', '').strip():
                    return True, 'api-key'
            elif aid == 'codex':
                if (home / '.codex' / 'auth.json').is_file():
                    return True, 'oauth'
                if os.environ.get('OPENAI_API_KEY', '').strip():
                    return True, 'api-key'
            elif aid == 'gemini':
                gdir = home / '.gemini'
                if ((gdir / 'oauth_creds.json').is_file()
                        or (gdir / 'google_accounts.json').is_file()):
                    return True, 'oauth'
                if (os.environ.get('GEMINI_API_KEY', '').strip()
                        or os.environ.get('GOOGLE_API_KEY', '').strip()):
                    return True, 'api-key'
            elif aid == 'hermes':
                hh = pathlib.Path(os.environ.get('HERMES_HOME', '').strip()
                                  or (home / '.hermes'))
                if (hh / 'config.yaml').is_file():
                    return True, 'config'
        except OSError:
            pass
        return False, ''

    def _hermes_gateway_running(self):
        """True when the Hermes gateway is accepting connections on loopback."""
        try:
            with socket.create_connection(('127.0.0.1', HERMES_PORT), timeout=0.8):
                return True
        except OSError:
            return False

    def _agents_status(self):
        """GET /agents — which agents are available on this serve.py host.
        Version probes (one subprocess per installed CLI) run concurrently so the
        endpoint stays snappy even with several agents installed. Each agent also
        carries best-effort login detection (authenticated/auth) so the UI can
        sync the user's existing local CLIs — and hermes reports whether its
        gateway is actually up (running)."""
        specs = [
            ('hermes',      'Hermes Agent', True,  False, self._hermes_resolve,
             'Nous Research agent — your container’s default runtime.'),
            ('claude-code', 'Claude Code',  False, True,  self._claudecode_resolve,
             'Anthropic’s coding agent CLI (BYO key).'),
            ('codex',       'Codex',        False, True,  self._codex_resolve,
             'OpenAI’s coding agent CLI (BYO key).'),
            ('gemini',      'Gemini',       False, True,  self._gemini_resolve,
             'Google’s Gemini agent CLI (BYO key).'),
        ]
        bins = {aid: resolve() for (aid, _l, _d, _r, resolve, _ds) in specs}
        # Probe versions concurrently (bounded; each call self-limits to ~6s).
        from concurrent.futures import ThreadPoolExecutor
        with ThreadPoolExecutor(max_workers=len(specs)) as ex:
            futs = {aid: ex.submit(self._agent_version, b, '--version')
                    for aid, b in bins.items() if b}
            versions = {aid: f.result() for aid, f in futs.items()}
        auth = {aid: self._agent_auth_detect(aid) for aid in bins}
        hermes_up = bool(bins.get('hermes')) and self._hermes_gateway_running()
        agents = [
            {'id': aid, 'label': label, 'installed': bool(bins[aid]),
             'default': dflt, 'removable': rem,
             'version': versions.get(aid, ''), 'desc': desc,
             'authenticated': auth[aid][0], 'auth': auth[aid][1],
             **({'running': hermes_up} if aid == 'hermes' else {})}
            for (aid, label, dflt, rem, _resolve, desc) in specs
        ]
        return self._send_json(200, {'agents': agents})

    def _agents_install(self):
        """POST /agents/install { agent: 'claude-code'|'codex'|'gemini'|'hermes' }
        — install an agent CLI onto the serve.py host. Node CLIs go through npm;
        Hermes (a unix-only Python package) goes through pip. Fixed allowlist (no
        arbitrary packages). Synchronous; can take ~30-90s. Returns { ok, agent,
        installed, version }."""
        AGENT_NPM = {
            'claude-code': '@anthropic-ai/claude-code@latest',
            'codex':       '@openai/codex@latest',
            'gemini':      '@google/gemini-cli@latest',
        }
        # Hermes is unix-only (gateway/pty_bridge import termios/pty/fcntl) — it
        # only installs where serve.py runs on a unix host (container or WSL).
        AGENT_PIP = {
            'hermes':      'hermes-agent',
        }
        length = int(self.headers.get('content-length', 0) or 0)
        try:
            body = json.loads(self.rfile.read(length) or b'{}')
        except Exception:
            return self._send_json(400, {'error': 'bad json'})
        agent = str(body.get('agent', '')).strip().lower()

        if agent in AGENT_NPM:
            npm = shutil.which('npm')
            if not npm:
                return self._send_json(503, {'error': 'npm not available on this host'})
            cmd = [npm, 'install', '-g', AGENT_NPM[agent]]
        elif agent in AGENT_PIP:
            if sys.platform == 'win32':
                return self._send_json(400, {'error':
                    'hermes is unix-only and cannot install on native Windows — '
                    'run CafresoHQ in WSL (use Start-CafresoHQ) and install there'})
            # Prefer `pip` if present, else `python3 -m pip` (pip often isn't on
            # the WSL PATH). --user keeps it in the invoking user's site dir.
            pip = shutil.which('pip3') or shutil.which('pip')
            base = [pip] if pip else [sys.executable or 'python3', '-m', 'pip']
            cmd = base + ['install', '--user', '--upgrade', AGENT_PIP[agent]]
        else:
            return self._send_json(400, {
                'error': 'agent must be claude-code, codex, gemini, or hermes'})

        # Inherit the environment so npm/pip use the prefix+cache the image set
        # (Dockerfile + entrypoint export npm_config_prefix/cache + HOME); locally
        # this is the user's own npm/pip config. Guard HOME in case it's unset.
        install_env = dict(os.environ)
        install_env.setdefault('HOME', '/root' if sys.platform != 'win32'
                               else install_env.get('USERPROFILE', ''))

        # Run the install in a BACKGROUND thread and 202 immediately — the old
        # synchronous path parked a request thread for up to 600s, starving the
        # whole server while npm ran. The client polls /agents/install/status.
        with _INSTALL_JOBS_LOCK:
            cur = _INSTALL_JOBS.get(agent)
            if cur and cur.get('status') == 'running':
                return self._send_json(202, {'ok': True, 'agent': agent,
                                             'status': 'running',
                                             'note': 'install already in progress'})
            _INSTALL_JOBS[agent] = {'agent': agent, 'status': 'running',
                                    'started': time.time()}

        resolvers = {
            'claude-code': self._claudecode_resolve,
            'codex':       self._codex_resolve,
            'gemini':      self._gemini_resolve,
            'hermes':      self._hermes_resolve,
        }
        resolve = resolvers[agent]
        version_of = self._agent_version

        def _worker():
            result = {'status': 'error', 'error': 'unknown'}
            try:
                proc = subprocess.run(cmd, capture_output=True, text=True,
                                      timeout=600, env=install_env)
                if proc.returncode != 0:
                    full = ((proc.stderr or '') + '\n' + (proc.stdout or '')).strip()
                    # Persist the FULL log — truncating to 600 chars hid the real cause.
                    try:
                        _sd = os.environ.get('CAFRESOHQ_HQ_STATE_DIR', '/data/hq-state')
                        os.makedirs(_sd, exist_ok=True)
                        with open(os.path.join(_sd, 'agent-install-errors.log'), 'a') as _f:
                            _f.write('=== %s :: %s ===\n%s\n\n' % (agent, ' '.join(cmd), full))
                    except Exception:
                        pass
                    low = full.lower()
                    if 'eacces' in low or 'permission denied' in low:
                        hint = ('permission error writing the install prefix/cache — the host '
                                'must pre-create + chmod the npm/pip target dir')
                    elif any(s in low for s in ('enotfound', 'eai_again', 'etimedout',
                             'getaddrinfo', 'could not resolve', 'network', 'timed out')):
                        hint = ('network/DNS — this host cannot reach the package registry; '
                                'open egress to registry.npmjs.org:443 / pypi.org')
                    elif 'enospc' in low or 'no space' in low:
                        hint = 'no disk space left to complete the install'
                    else:
                        hint = full[-600:] or 'install failed'
                    result = {'status': 'error', 'error': hint}
                else:
                    bin_ = resolve()
                    result = {'status': 'done', 'ok': True,
                              'installed': bool(bin_),
                              'version': version_of(bin_, '--version')}
            except subprocess.TimeoutExpired:
                result = {'status': 'error', 'error':
                          'install timed out (>600s) — usually blocked network egress to '
                          'the package registry (open outbound to registry.npmjs.org:443 / pypi.org)'}
            except Exception as e:
                result = {'status': 'error', 'error': str(e)}
            result['agent'] = agent
            result['finished'] = time.time()
            with _INSTALL_JOBS_LOCK:
                _INSTALL_JOBS[agent] = result

        threading.Thread(target=_worker, daemon=True, name=f'install-{agent}').start()
        return self._send_json(202, {'ok': True, 'agent': agent, 'status': 'started',
                                     'note': 'installing in background — poll '
                                             '/agents/install/status?agent=' + agent})

    def _agents_install_status(self):
        """GET /agents/install/status?agent=<id> → the background job's state:
        {status: running|done|error, …}. 404 when no install was started."""
        qs = urllib.parse.parse_qs(urllib.parse.urlparse(self.path).query)
        agent = (qs.get('agent', [''])[0] or '').strip().lower()
        with _INSTALL_JOBS_LOCK:
            job = dict(_INSTALL_JOBS.get(agent) or {})
        if not job:
            return self._send_json(404, {'agent': agent, 'status': 'none'})
        return self._send_json(200, job)

    # ──────────────────────────────────────────────────────────────────────
    # Browser shim — gives agents a way to fetch URLs and grab screenshots.
    #
    # /browser/fetch?url=...        Pure-stdlib HTTP GET + readable text
    #                               extraction. Always available; works for
    #                               any agent regardless of whether a
    #                               browser is running. Suitable for "read
    #                               this article" or "scrape this docs page"
    #                               jobs that don't need JS.
    #
    # /browser/screenshot?url=...   Talks to a Chromium-based browser via
    #                               the Chrome DevTools Protocol (CDP). The
    #                               user must have launched Brave/Chrome
    #                               with --remote-debugging-port=9222 OR
    #                               CAFRESOHQ_BROWSER_CDP_URL must be set.
    #                               Returns base64 PNG inside JSON. If no
    #                               CDP target is reachable we return 503
    #                               with a clear "how to enable" hint
    #                               instead of crashing — the rest of the
    #                               app keeps working.
    #
    # /browser/status               One-shot probe — tells the UI whether
    #                               screenshots are available so we can
    #                               surface the right hint.
    # ──────────────────────────────────────────────────────────────────────
    def _browser_status(self):
        cdp_ok, cdp_info = self._browser_cdp_probe()
        return self._send_json(200, {
            'fetchAvailable': True,            # always works (urllib)
            'screenshotAvailable': cdp_ok,
            'cdp': cdp_info,
            'cdpHint': (
                'Launch Brave/Chrome with --remote-debugging-port=9222, or '
                'set CAFRESOHQ_BROWSER_CDP_URL=http://host:port to enable '
                'screenshots and JS-rendered fetch.'
            ),
        })

    def _browser_cdp_probe(self):
        """Find a CDP endpoint. Tries CAFRESOHQ_BROWSER_CDP_URL first,
        then localhost:9222 (Brave/Chrome default). Returns (ok, info)."""
        candidates = []
        env_url = os.environ.get('CAFRESOHQ_BROWSER_CDP_URL', '').strip()
        if env_url:
            candidates.append(env_url.rstrip('/'))
        candidates.append('http://127.0.0.1:9222')
        for base in candidates:
            try:
                u = urllib.parse.urlparse(base)
                conn = http.client.HTTPConnection(u.hostname, u.port or 80, timeout=2)
                conn.request('GET', '/json/version')
                resp = conn.getresponse()
                if resp.status == 200:
                    info = json.loads(resp.read().decode('utf-8', 'replace'))
                    conn.close()
                    return True, {'base': base, 'browser': info.get('Browser', '?')}
                conn.close()
            except Exception:
                continue
        return False, {'base': None, 'browser': None}

    def _browser_fetch(self):
        """GET /browser/fetch?url=...&max_chars=N
        Returns { url, status, title, text, length } where text is the
        page's readable content with HTML stripped + whitespace normalised.
        No JS execution — for that, use /browser/screenshot which goes
        through CDP. Hard cap on size (default 50KB) so a giant page
        doesn't blow up the agent's context."""
        qs = urllib.parse.urlparse(self.path).query
        params = urllib.parse.parse_qs(qs)
        url = (params.get('url') or [''])[0].strip()
        try:
            max_chars = int((params.get('max_chars') or ['50000'])[0])
        except ValueError:
            max_chars = 50000
        max_chars = max(500, min(max_chars, 200000))
        if not url:
            return self._send_json(400, {'error': 'url query parameter required'})
        if not (url.startswith('http://') or url.startswith('https://')):
            return self._send_json(400, {'error': 'url must start with http:// or https://'})
        try:
            req = urllib.request.Request(url, headers={
                'User-Agent': 'CafresoHQ/1.0 (+browser-shim)',
                'Accept': 'text/html,application/xhtml+xml,*/*;q=0.8',
                'Accept-Language': 'en-US,en;q=0.9',
            })
            with urllib.request.urlopen(req, timeout=20) as r:
                status = r.status
                ctype = r.headers.get('Content-Type', '')
                # Read up to 2MB raw before processing
                raw = r.read(2 * 1024 * 1024)
        except urllib.error.HTTPError as e:
            return self._send_json(200, {
                'url': url, 'status': e.code,
                'error': f'HTTP {e.code}: {e.reason}',
                'text': '', 'title': '', 'length': 0,
            })
        except Exception as e:
            return self._send_json(502, {'error': f'fetch failed: {type(e).__name__}: {e}'})
        # Decode (best-effort)
        encoding = 'utf-8'
        if 'charset=' in ctype.lower():
            try: encoding = ctype.split('charset=', 1)[1].split(';')[0].strip() or 'utf-8'
            except Exception: pass
        try: html = raw.decode(encoding, errors='replace')
        except Exception: html = raw.decode('utf-8', errors='replace')
        # Title
        m_title = re.search(r'<title[^>]*>([^<]*)</title>', html, re.IGNORECASE)
        title = (m_title.group(1).strip() if m_title else '')[:200]
        # Strip scripts, styles, noscript, then tags. Cheap but effective.
        cleaned = re.sub(r'<script[\s\S]*?</script>', ' ', html, flags=re.IGNORECASE)
        cleaned = re.sub(r'<style[\s\S]*?</style>', ' ', cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(r'<noscript[\s\S]*?</noscript>', ' ', cleaned, flags=re.IGNORECASE)
        # Preserve paragraph breaks before stripping tags
        cleaned = re.sub(r'</(p|div|h[1-6]|li|br|tr)>', '\n', cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(r'<br\s*/?>', '\n', cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(r'<[^>]+>', ' ', cleaned)
        # Collapse whitespace; keep blank lines
        cleaned = re.sub(r'[ \t]+', ' ', cleaned)
        cleaned = re.sub(r'\n[ \t]*\n+', '\n\n', cleaned)
        # HTML entity decode
        try:
            import html as _htmlmod
            cleaned = _htmlmod.unescape(cleaned)
        except Exception:
            pass
        cleaned = cleaned.strip()
        truncated = len(cleaned) > max_chars
        if truncated:
            cleaned = cleaned[:max_chars] + f'\n\n[…truncated, {len(cleaned) - max_chars} more chars]'
        return self._send_json(200, {
            'url': url, 'status': status, 'title': title,
            'text': cleaned, 'length': len(cleaned), 'truncated': truncated,
        })

    def _browser_screenshot(self):
        """GET /browser/screenshot?url=...&width=1280&height=800&fullPage=0
        Drives a Chromium-based browser via CDP to navigate + screenshot.
        Returns { url, width, height, png } where png is a data: URL.
        Requires a running browser with --remote-debugging-port=9222
        (or CAFRESOHQ_BROWSER_CDP_URL pointing at one).

        We use the lightweight HTTP+WS subset of CDP via stdlib only:
          1. GET /json/version  → discover webSocketDebuggerUrl
          2. PUT /json/new       → create a fresh tab (so we don't
                                   stomp on the user's actual browsing)
          3. WS dance: Page.navigate → Page.captureScreenshot
          4. DELETE /json/close/<id> → tidy up
        """
        cdp_ok, cdp_info = self._browser_cdp_probe()
        if not cdp_ok:
            return self._send_json(503, {
                'error': 'No CDP-enabled browser detected.',
                'hint': (
                    'Launch Brave or Chrome with --remote-debugging-port=9222 '
                    '(or set CAFRESOHQ_BROWSER_CDP_URL to a different host:port).'
                ),
            })
        qs = urllib.parse.urlparse(self.path).query
        params = urllib.parse.parse_qs(qs)
        url = (params.get('url') or [''])[0].strip()
        width = int((params.get('width') or ['1280'])[0])
        height = int((params.get('height') or ['800'])[0])
        full_page = (params.get('fullPage') or ['0'])[0] in ('1', 'true', 'yes')
        if not url:
            return self._send_json(400, {'error': 'url query parameter required'})
        try:
            png_b64 = self._cdp_screenshot(cdp_info['base'], url, width, height, full_page)
        except Exception as e:
            return self._send_json(502, {'error': f'screenshot failed: {type(e).__name__}: {e}'})
        return self._send_json(200, {
            'url': url, 'width': width, 'height': height,
            'png': 'data:image/png;base64,' + png_b64,
        })

    def _cdp_screenshot(self, base, url, width, height, full_page):
        """Open a fresh tab via CDP, navigate, screenshot, close tab.
        Pure stdlib WebSocket — small enough to do by hand for one
        request/response per CDP method. Cap timeout per step."""
        # 1. Open a new tab.
        u = urllib.parse.urlparse(base)
        conn = http.client.HTTPConnection(u.hostname, u.port or 80, timeout=10)
        conn.request('PUT', '/json/new?about:blank')
        resp = conn.getresponse()
        if resp.status not in (200, 201):
            raise RuntimeError(f'/json/new returned {resp.status}')
        tab = json.loads(resp.read().decode('utf-8', 'replace'))
        conn.close()
        ws_url = tab['webSocketDebuggerUrl']
        tab_id = tab['id']
        try:
            # 2. WS dance.
            ws = _CDPWebSocket(ws_url, timeout=20)
            ws.connect()
            # Set viewport
            ws.call('Emulation.setDeviceMetricsOverride',
                    {'width': width, 'height': height, 'deviceScaleFactor': 1, 'mobile': False})
            ws.call('Page.enable', {})
            ws.call('Page.navigate', {'url': url})
            # Wait for load (up to 12s) — poll lifecycle events.
            ws.wait_for_event('Page.loadEventFired', timeout=12)
            # Tiny grace period for late-rendering JS.
            time.sleep(0.6)
            cap = ws.call('Page.captureScreenshot',
                          {'format': 'png', 'captureBeyondViewport': bool(full_page)})
            ws.close()
            return cap['result']['data']
        finally:
            # 3. Close tab.
            try:
                conn = http.client.HTTPConnection(u.hostname, u.port or 80, timeout=5)
                conn.request('GET', f'/json/close/{tab_id}')
                conn.getresponse().read()
                conn.close()
            except Exception:
                pass

    def _cafresohq_stream(self):
        """Streaming chat for ELEVATED agents. Same SSE-shaped output as
        /claudecode/stream so the client uses one parser. Differences:
          - tools enabled (--allowed-tools <server-side allowlist>)
          - working set bound to --add-dir <each allowed dir>
          - refuses if either allowlist is empty (safe default)
        """
        bin_ = self._claudecode_resolve()
        if not bin_:
            return self._send_json(503, {
                'error': 'claude CLI not found — install Claude Code or set CAFRESOHQ_CLAUDE_BIN'
            })
        if not _cafresohq_allowed_dirs:
            return self._send_json(503, {
                'error': 'CAFRESOHQ_ALLOWED_DIRS not set — elevated endpoint disabled'
            })
        if not _cafresohq_allowed_tools:
            return self._send_json(503, {
                'error': 'CAFRESOHQ_ALLOWED_TOOLS empty — elevated endpoint disabled'
            })
        # Refuse if any configured dir is missing — easier to fix typos than
        # to debug a silent permission error.
        for d in _cafresohq_allowed_dirs:
            if not pathlib.Path(d).is_dir():
                return self._send_json(503, {
                    'error': f'CAFRESOHQ_ALLOWED_DIRS includes missing dir: {d}'
                })

        length = int(self.headers.get('content-length', 0) or 0)
        try:
            body = json.loads(self.rfile.read(length).decode('utf-8') or '{}')
        except json.JSONDecodeError:
            return self._send_json(400, {'error': 'bad json'})

        messages = body.get('messages') or []
        system   = (body.get('system') or '').strip()
        model    = (body.get('model') or '').strip()
        agent    = (body.get('agentName') or 'elevated-agent').strip()[:60]

        history_lines = []
        for m in messages:
            role = (m.get('role') or 'user').upper()
            content = (m.get('content') or '').strip()
            if not content: continue
            history_lines.append(f'{role}: {content}')
        prompt = '\n\n'.join(history_lines)
        if not prompt:
            return self._send_json(400, {'error': 'no prompt'})

        # Compose the Claude Code command. Critical bits:
        #   --allowed-tools <comma-list>  bounds what Claude can DO
        #   --add-dir <abs-path>          bounds where it can READ/WRITE
        # Both come from server-side env, NOT from the request body.
        cmd = [bin_, '--print',
               '--output-format', 'stream-json',
               '--verbose',
               '--allowed-tools', ','.join(_cafresohq_allowed_tools),
               '--input-format', 'text']
        for d in _cafresohq_allowed_dirs:
            cmd += ['--add-dir', d]
        if model:
            cmd += ['--model', model]
        # Always prepend a guard-rail system note so the elevated agent knows
        # it's running under HQ's authority and what its boundaries are.
        guard = (
            f'You are {agent}, an elevated HQ agent with computer access. '
            f'You are restricted to these directories: {", ".join(_cafresohq_allowed_dirs)}. '
            'When you intend to perform an action with side effects, first emit '
            '[NEEDS_APPROVAL: <one-line description>] and stop until the boss replies.'
        )
        full_system = guard + ('\n\n' + system if system else '')
        cmd += ['--append-system-prompt', full_system]

        # Audit: log every elevated invocation server-side so there's a
        # tamper-resistant record outside the browser.
        sys.stderr.write(f'[cafresohq] elevated stream: agent={agent} model={model or "(default)"} '
                         f'dirs={_cafresohq_allowed_dirs} tools={_cafresohq_allowed_tools}\n')

        # Use project cwd if provided and valid, otherwise first allowed dir.
        agent_cwd = _cafresohq_allowed_dirs[0] if _cafresohq_allowed_dirs else None
        req_cwd = (body.get('cwd') or '').strip()
        if req_cwd:
            try:
                cwd_p = pathlib.Path(_client_path(req_cwd)).resolve()
                for d in _cafresohq_allowed_dirs:
                    try:
                        cwd_p.relative_to(pathlib.Path(d).resolve())
                        if cwd_p.is_dir():
                            agent_cwd = str(cwd_p)
                        break
                    except ValueError:
                        continue
            except OSError:
                pass

        # Ensure Git Bash is in PATH so Claude Code's Bash tool can find
        # bash.exe on Windows. Without this, Bash exits 255 with
        # "The system cannot find the path specified."
        import copy as _copy
        agent_env = _copy.deepcopy(os.environ)
        _git_bash_dirs = [
            r'C:\Program Files\Git\usr\bin',
            r'C:\Program Files\Git\bin',
            r'C:\Program Files\Git\mingw64\bin',
        ]
        _path_parts = agent_env.get('PATH', '').split(os.pathsep)
        for _d in reversed(_git_bash_dirs):
            if os.path.isdir(_d) and _d not in _path_parts:
                agent_env['PATH'] = _d + os.pathsep + agent_env.get('PATH', '')

        try:
            proc = subprocess.Popen(
                cmd,
                stdin=subprocess.PIPE,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True, bufsize=1, encoding='utf-8',
                cwd=agent_cwd,
                env=agent_env,
            )
        except FileNotFoundError as e:
            return self._send_json(500, {'error': f'spawn: {e}'})

        try:
            proc.stdin.write(prompt)
            proc.stdin.close()
        except Exception as e:
            try: proc.kill()
            except Exception: pass
            return self._send_json(500, {'error': f'stdin: {e}'})

        # SSE response — same shape /claudecode/stream emits.
        self.send_response(200)
        self.send_header('content-type', 'text/event-stream')
        self.send_header('cache-control', 'no-store')
        self.end_headers()

        def write_sse(obj):
            try:
                self.wfile.write(b'data: ' + json.dumps(obj).encode('utf-8') + b'\n\n')
                self.wfile.flush()
            except (BrokenPipeError, ConnectionResetError):
                return False
            return True

        stderr_buf = []
        def _drain_err():
            try:
                for line in proc.stderr:
                    stderr_buf.append(line)
            except Exception:
                pass
        threading.Thread(target=_drain_err, daemon=True).start()

        in_tokens = 0
        out_tokens = 0
        try:
            for line in proc.stdout:
                line = line.strip()
                if not line: continue
                try:
                    ev = json.loads(line)
                except json.JSONDecodeError:
                    continue
                t = ev.get('type')
                if t == 'assistant':
                    msg = ev.get('message') or {}
                    for block in (msg.get('content') or []):
                        if block.get('type') == 'text':
                            text = block.get('text') or ''
                            if text:
                                ok = write_sse({
                                    'choices': [{
                                        'index': 0,
                                        'delta': {'content': text},
                                    }],
                                })
                                if not ok: break
                    u = msg.get('usage')
                    if u:
                        in_tokens  = u.get('input_tokens',  in_tokens)
                        out_tokens = u.get('output_tokens', out_tokens)
                elif t == 'result':
                    u = ev.get('usage') or {}
                    in_tokens  = u.get('input_tokens',  in_tokens)
                    out_tokens = u.get('output_tokens', out_tokens)
                elif t == 'error':
                    write_sse({'choices': [{'index': 0, 'delta': {
                        'content': f"\n\n⚠ CafresoHQ error: {ev.get('message') or ev}"}}]})
            if in_tokens or out_tokens:
                write_sse({
                    'choices': [],
                    'usage': {
                        'prompt_tokens': in_tokens,
                        'completion_tokens': out_tokens,
                        'total_tokens': in_tokens + out_tokens,
                    },
                })
        finally:
            try: proc.wait(timeout=2)
            except Exception:
                try: proc.kill()
                except Exception: pass
            if proc.returncode and proc.returncode != 0 and not (in_tokens or out_tokens):
                err_text = ('\n'.join(stderr_buf))[:600] or f'exit {proc.returncode}'
                write_sse({'choices': [{'index': 0, 'delta': {
                    'content': f"\n\n⚠ CafresoHQ exited {proc.returncode}: {err_text}"}}]})

    # ---- Codex CLI elevated agent ----------------------------------------
    def _codex_resolve(self):
        if _codex_bin and pathlib.Path(_codex_bin).is_file():
            return _codex_bin
        # On Windows, prefer codex.cmd over the extensionless npm wrapper.
        # The extensionless file is a shell script that Windows can't exec
        # directly — subprocess.Popen on it raises OSError 193 ("not a valid
        # Win32 application"), which used to escape the except FileNotFoundError
        # guard and kill the HTTP connection (→ browser "Failed to fetch").
        if sys.platform == 'win32':
            return (shutil.which('codex.cmd')
                    or shutil.which(_codex_bin or 'codex')
                    or shutil.which('codex')
                    or '')
        return shutil.which(_codex_bin or 'codex') or shutil.which('codex.cmd') or ''

    def _hermes_resolve(self):
        """Find the hermes binary. Returns absolute path or ''.
        CAFRESOHQ_HERMES_BIN overrides; otherwise PATH lookup. In the container
        and in WSL/unix this resolves natively; on native Windows it returns ''
        (hermes is unix-only — run the stack in WSL instead)."""
        if _hermes_bin and pathlib.Path(_hermes_bin).is_file():
            return _hermes_bin
        return shutil.which(_hermes_bin or 'hermes') or ''

    def _gemini_resolve(self):
        """Find the gemini binary (npm @google/gemini-cli). Returns path or ''.
        CAFRESOHQ_GEMINI_BIN overrides. On Windows prefer the .cmd npm wrapper
        (same OSError-193 reasoning as _codex_resolve)."""
        if _gemini_bin and pathlib.Path(_gemini_bin).is_file():
            return _gemini_bin
        if sys.platform == 'win32':
            return (shutil.which('gemini.cmd')
                    or shutil.which(_gemini_bin or 'gemini')
                    or shutil.which('gemini')
                    or '')
        return shutil.which(_gemini_bin or 'gemini') or shutil.which('gemini.cmd') or ''

    def _codex_configure(self):
        global _codex_bin
        length = int(self.headers.get('content-length', 0) or 0)
        try:
            body = json.loads(self.rfile.read(length).decode('utf-8') or '{}')
        except json.JSONDecodeError:
            return self._send_json(400, {'error': 'bad json'})
        path = (body.get('binary') or '').strip()
        if path and not pathlib.Path(path).is_file():
            return self._send_json(400, {'error': f'not a file: {path}'})
        _codex_bin = path
        return self._codex_status()

    def _codex_stream(self):
        """Streaming elevated agent backed by Codex CLI (codex exec --json).
        Auth comes from ~/.codex/config.toml — serve.py never touches API keys.
        Same SSE output shape as /cafresohq/stream so the UI uses one parser.
        """
        bin_ = self._codex_resolve()
        if not bin_:
            return self._send_json(503, {
                'error': 'codex CLI not found — install via npm i -g @openai/codex or set CAFRESOHQ_CODEX_BIN'
            })
        if not _cafresohq_allowed_dirs:
            return self._send_json(503, {
                'error': 'CAFRESOHQ_ALLOWED_DIRS not set — codex endpoint disabled'
            })

        length = int(self.headers.get('content-length', 0) or 0)
        try:
            body = json.loads(self.rfile.read(length).decode('utf-8') or '{}')
        except json.JSONDecodeError:
            return self._send_json(400, {'error': 'bad json'})

        for d in _cafresohq_allowed_dirs:
            if not pathlib.Path(d).is_dir():
                return self._send_json(503, {
                    'error': f'CAFRESOHQ_ALLOWED_DIRS includes missing dir: {d}'
                })

        messages = body.get('messages') or []
        model    = (body.get('model') or '').strip()
        system   = (body.get('system') or '').strip()
        agent    = (body.get('agentName') or body.get('agent') or 'elevated-agent').strip()[:60]

        # Build prompt from message history (same as cafresohq handler)
        history_lines = []
        for m in messages:
            role = (m.get('role') or 'user').upper()
            content = (m.get('content') or '').strip()
            if not content: continue
            history_lines.append(f'{role}: {content}')
        prompt = '\n\n'.join(history_lines)
        if not prompt:
            return self._send_json(400, {'error': 'no prompt'})

        # Use project cwd if provided and valid, otherwise first allowed dir.
        primary_dir = _cafresohq_allowed_dirs[0]
        req_cwd = (body.get('cwd') or '').strip()
        if req_cwd:
            try:
                cwd_p = pathlib.Path(_client_path(req_cwd)).resolve()
                for d in _cafresohq_allowed_dirs:
                    try:
                        cwd_p.relative_to(pathlib.Path(d).resolve())
                        if cwd_p.is_dir():
                            primary_dir = str(cwd_p)
                        break
                    except ValueError:
                        continue
            except OSError:
                pass
        cmd = [bin_, 'exec',
               '--json',
               '--skip-git-repo-check',
               '--sandbox', 'workspace-write',
               '-C', primary_dir]
        for d in _cafresohq_allowed_dirs[1:]:
            cmd += ['--add-dir', d]
        # Codex talks directly to OpenAI (default provider). Auth via
        # OPENAI_API_KEY (user-supplied through HQ settings / operator env).
        wire_model = model or 'gpt-4.1'
        cmd += ['-c', 'model_provider="openai"']
        cmd += ['--model', wire_model]
        # OCI cloud requirements (synced into ~/.codex/cloud-requirements-cache.json)
        # restrict approval_policy to "untrusted" — anything else (including
        # "never") is rejected as `Configured value … is disallowed by
        # requirements; falling back to required value UnlessTrusted`. That
        # rejection used to surface as an `item.completed type=error` event
        # and the user saw "Codex returned no content". Use the allowed
        # value so the request actually runs.
        cmd += ['-c', 'approval_policy="untrusted"']

        guard = (
            f'You are {agent}, an elevated HQ agent with computer access. '
            f'You are restricted to these directories: {", ".join(_cafresohq_allowed_dirs)}. '
            'When you intend to perform an action with side effects, first emit '
            '[NEEDS_APPROVAL: <one-line description>] and stop until the boss replies.'
        )
        full_prompt = (guard + '\n\n' + (system + '\n\n' if system else '')) + prompt

        sys.stderr.write(f'[codex] stream: agent={agent} model={model or "(default)"} '
                         f'wire_model={wire_model or "(default)"} '
                         f'dirs={_cafresohq_allowed_dirs}\n')

        # Inherit env + Git Bash paths so shell tools work on Windows.
        # Codex talks to OpenAI directly via OPENAI_API_KEY (BYOK / operator env).
        import copy as _copy
        agent_env = _copy.deepcopy(os.environ)
        path_key = next((k for k in agent_env.keys() if k.lower() == 'path'), 'Path')
        path_value = agent_env.get(path_key, '')
        for _d in [r'C:\Program Files\Git\usr\bin',
                   r'C:\Program Files\Git\bin',
                   r'C:\Program Files\Git\mingw64\bin']:
            if os.path.isdir(_d) and _d not in path_value:
                path_value = _d + os.pathsep + path_value
        path_value = os.pathsep.join(
            p for p in path_value.split(os.pathsep)
            if p and r'\.codex\tmp\arg0' not in p.lower()
        )
        for _k in [k for k in list(agent_env.keys()) if k.lower() == 'path']:
            agent_env.pop(_k, None)
        agent_env['Path'] = path_value
        # Codex uses OpenAI directly via OPENAI_API_KEY; no base_url override.
        agent_env.pop('OPENAI_BASE_URL', None)

        try:
            proc = subprocess.Popen(
                cmd,
                stdin=subprocess.PIPE,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True, bufsize=1, encoding='utf-8',
                cwd=primary_dir,
                env=agent_env,
            )
        except (FileNotFoundError, OSError) as e:
            return self._send_json(500, {'error': f'spawn codex: {e}'})

        try:
            proc.stdin.write(full_prompt)
            proc.stdin.close()
        except Exception as e:
            # Codex died before we could send the prompt — capture stderr so
            # the user sees the real failure (bad profile, bad config, etc.)
            stderr_text = ''
            try:
                stderr_text = proc.stderr.read()[:1500]
            except Exception:
                pass
            try: proc.kill()
            except Exception: pass
            return self._send_json(500, {
                'error': f'codex exited before prompt could be sent: {e}',
                'stderr': stderr_text,
                'cmd': ' '.join(cmd),
            })

        # SSE response
        self.send_response(200)
        self.send_header('content-type', 'text/event-stream')
        self.send_header('cache-control', 'no-store')
        self.end_headers()

        def write_sse(obj):
            try:
                self.wfile.write(b'data: ' + json.dumps(obj).encode('utf-8') + b'\n\n')
                self.wfile.flush()
            except (BrokenPipeError, ConnectionResetError):
                return False
            return True

        stderr_buf = []
        def _drain_err():
            try:
                for line in proc.stderr:
                    stderr_buf.append(line)
            except Exception:
                pass
        threading.Thread(target=_drain_err, daemon=True).start()

        item_text_seen = {}

        def _content_text(value):
            if value is None:
                return ''
            if isinstance(value, str):
                return value
            if isinstance(value, list):
                parts = []
                for block in value:
                    if isinstance(block, str):
                        parts.append(block)
                    elif isinstance(block, dict):
                        parts.append(block.get('text') or block.get('content') or block.get('output_text') or '')
                return ''.join(parts)
            if isinstance(value, dict):
                return (
                    value.get('text')
                    or value.get('content')
                    or value.get('message')
                    or _content_text(value.get('content_parts'))
                )
            return ''

        # Codex --json event types we care about:
        #   agent_message / message  → text to show the user
        #   function_call            → show tool invocation inline
        #   function_call_output     → show tool result inline
        #   error                    → surface as warning
        #   turn.failed              → surface as error
        #   thread.started / turn.started / turn.completed → silently ignored
        events_seen = 0
        text_emitted = False
        try:
            for line in proc.stdout:
                line = line.strip()
                if not line:
                    continue
                try:
                    ev = json.loads(line)
                except json.JSONDecodeError:
                    continue

                events_seen += 1
                t = ev.get('type', '')

                # Assistant text output
                if t in ('agent_message', 'message'):
                    content = ev.get('content') or ev.get('message') or ''
                    text = _content_text(content)
                    if text:
                        text_emitted = True
                        write_sse({'choices': [{'index': 0, 'delta': {'content': text}}]})

                elif t in ('agent_message_delta', 'message_delta', 'response.output_text.delta'):
                    text = ev.get('delta') or ev.get('text') or ev.get('content') or ''
                    if text:
                        text_emitted = True
                        write_sse({'choices': [{'index': 0, 'delta': {'content': str(text)}}]})

                elif t in ('item.updated', 'item.completed'):
                    item = ev.get('item') or {}
                    item_type = item.get('type', '')
                    # Codex 0.128 emits `type: "agent_message"` with a flat
                    # `text` field for assistant replies (older versions used
                    # `type: "message"` with `content` blocks). Surface errors
                    # too — those are stream-disconnect / config-rejection
                    # diagnostics the user needs to see.
                    if item_type == 'error':
                        msg = item.get('message') or item.get('text') or json.dumps(item)
                        # Suppress informational warnings that don't block
                        # the run — the cloud-requirements override message
                        # always fires for `codex exec` because the CLI
                        # defaults approval_policy to `Never` and the cloud
                        # gate downgrades it to UnlessTrusted. Codex still
                        # produces an answer in the next item.completed,
                        # so showing the warning to the user is just noise.
                        _is_warning = (
                            'disallowed by requirements' in msg
                            or 'falling back to required value' in msg
                        )
                        if not _is_warning:
                            text_emitted = True
                            write_sse({'choices': [{'index': 0, 'delta': {
                                'content': f'\n\n⚠ Codex item error: {msg}'}}]})
                    elif (item.get('role') == 'assistant'
                          or item_type in ('message', 'agent_message')):
                        text = (item.get('text')
                                or _content_text(item.get('content') or item.get('message')))
                        item_id = item.get('id') or ev.get('item_id') or 'assistant'
                        prior = item_text_seen.get(item_id, '')
                        if text and text.startswith(prior):
                            delta = text[len(prior):]
                        elif text and text != prior:
                            delta = text
                        else:
                            delta = ''
                        item_text_seen[item_id] = text or prior
                        if delta:
                            text_emitted = True
                            write_sse({'choices': [{'index': 0, 'delta': {'content': delta}}]})

                elif t in ('turn.completed', 'response.completed'):
                    text = _content_text(ev.get('last_agent_message') or ev.get('message') or '')
                    if text:
                        text_emitted = True
                        write_sse({'choices': [{'index': 0, 'delta': {'content': text}}]})

                # Tool call — show what Codex is doing
                elif t == 'function_call':
                    name = ev.get('name') or ev.get('function') or 'tool'
                    args = ev.get('arguments') or ev.get('input') or {}
                    if isinstance(args, dict):
                        cmd_str = args.get('cmd') or args.get('command') or json.dumps(args)
                    else:
                        cmd_str = str(args)
                    text_emitted = True
                    write_sse({'choices': [{'index': 0, 'delta': {
                        'content': f'\n[{name.upper()}: {cmd_str}]\n'}}]})

                # Tool result
                elif t == 'function_call_output':
                    output = ev.get('output') or ''
                    if output:
                        text_emitted = True
                        write_sse({'choices': [{'index': 0, 'delta': {
                            'content': f'\n📡 {name if (name := ev.get("name","tool")) else "result"}("{cmd_str[:60] if (cmd_str := str(ev.get("call_id",""))) else "..."}" →\n{output}\n'}}]})

                # Errors
                elif t in ('error', 'turn.failed'):
                    msg = ev.get('message') or (ev.get('error') or {}).get('message') or str(ev)
                    text_emitted = True
                    write_sse({'choices': [{'index': 0, 'delta': {
                        'content': f'\n\n⚠ Codex error: {msg}'}}]})

        finally:
            try: proc.wait(timeout=2)
            except Exception:
                try: proc.kill()
                except Exception: pass

            stderr_text = (''.join(stderr_buf)).strip()
            rc = proc.returncode

            def _summarize_stderr(s):
                """Codex sometimes dumps the entire upstream HTTP body into
                stderr. Trim that to the actual error summary so the chat
                surface doesn't get blasted with HTML."""
                if not s: return ''
                # Known noise: a model-list refresh failure can include a large
                # HTML payload. Show only the reason line.
                m = re.search(r'failed to refresh available models:[^\n]+', s)
                refresh_hint = ''
                if m:
                    refresh_hint = (
                        '\n\nNote: the model-list refresh failed. Check that '
                        'OPENAI_API_KEY is set and the selected model is valid.'
                    )
                # Drop any line that looks like raw HTML / JSON body, or
                # known-benign codex startup chatter ("Reading prompt from
                # stdin...", model-refresh notices that we already handle
                # via the proxy). These are not user-facing errors.
                _BENIGN_PATTERNS = (
                    'Reading prompt from stdin',
                    'codex_models_manager',
                    'failed to refresh available models',
                )
                cleaned = []
                for line in s.splitlines():
                    if line.startswith(' ') or line.startswith('"') or line.startswith('}') or line.startswith('{'):
                        continue
                    if '<!DOCTYPE' in line or '<html' in line:
                        continue
                    if any(p in line for p in _BENIGN_PATTERNS):
                        continue
                    cleaned.append(line)
                summary = '\n'.join(cleaned)[:600].rstrip()
                return summary + refresh_hint

            # Surface failures explicitly. Three cases that all used to look
            # like a silent successful run with only [DONE]:
            #   1. non-zero exit              → exited N: <stderr>
            #   2. zero exit + no events      → no output (config/sandbox issue)
            #   3. zero exit + stderr noise   → completed with warnings
            if rc and rc not in (0, None):
                err_text = _summarize_stderr(stderr_text) or f'exit {rc}'
                write_sse({'choices': [{'index': 0, 'delta': {
                    'content': f'\n\n⚠ Codex exited {rc}: {err_text}'}}]})
            elif not text_emitted:
                hint = _summarize_stderr(stderr_text) or '(no stdout, no stderr — likely the OS blocked the nested codex spawn)'
                write_sse({'choices': [{'index': 0, 'delta': {
                    'content': f'\n\n⚠ Codex returned no content (events seen: {events_seen}, exit: {rc}). {hint}'}}]})
            elif stderr_text:
                summary = _summarize_stderr(stderr_text)
                if summary:
                    write_sse({'choices': [{'index': 0, 'delta': {
                        'content': f'\n\n_(codex stderr: {summary[:300]})_'}}]})

            try: self.wfile.write(b'data: [DONE]\n\n'); self.wfile.flush()
            except (BrokenPipeError, ConnectionResetError): pass

    # ---- Filesystem browser (used by Projects path picker) ---------------
    def _fs_browse(self):
        """GET /fs/browse?path=<dir>
        Returns a directory listing for the path picker popup.
        In container mode: restricted to CAFRESOHQ_ALLOWED_DIRS subtrees.
        In local mode: any readable path is allowed.
        Response: {path, parent, entries:[{name, type:'dir'|'file', path}]}
        """
        qs = urllib.parse.urlparse(self.path).query
        params = urllib.parse.parse_qs(qs)
        req_path = (params.get('path') or [''])[0].strip()

        # Default starting location: first allowed dir → home → cwd
        if not req_path:
            if _cafresohq_allowed_dirs:
                req_path = _cafresohq_allowed_dirs[0]
            else:
                try:
                    req_path = str(pathlib.Path.home())
                except Exception:
                    req_path = os.getcwd()

        try:
            p = pathlib.Path(_client_path(req_path)).resolve()
        except Exception as e:
            return self._send_json(400, {'error': f'invalid path: {e}'})

        if not p.is_dir():
            return self._send_json(400, {'error': 'not a directory'})

        # Only browse within allowed dirs — enforced in EVERY mode (was
        # container-only, which left local/BYO reads unbounded).
        if not _within_allowed_dirs(p):
            return self._send_json(403, {
                'error': 'path is outside CAFRESOHQ_ALLOWED_DIRS',
                'allowed': _cafresohq_allowed_dirs,
            })

        try:
            raw = list(p.iterdir())
        except PermissionError:
            return self._send_json(403, {'error': 'permission denied'})
        except Exception as e:
            return self._send_json(500, {'error': f'listing failed: {e}'})

        # Build entries — guard each is_dir() individually because Windows
        # junction points / reparse points can raise ValueError for some paths.
        entries = []
        for item in raw:
            if item.name.startswith('.'):
                continue
            try:
                is_dir = item.is_dir()
            except Exception:
                continue  # skip inaccessible / broken entries silently
            entries.append({
                'name': item.name,
                'type': 'dir' if is_dir else 'file',
                'path': str(item),
            })
        # Dirs first, then files, both alphabetical
        entries.sort(key=lambda e: (e['type'] != 'dir', e['name'].lower()))

        parent = str(p.parent) if str(p.parent) != str(p) else None
        return self._send_json(200, {
            'path':    str(p),
            'parent':  parent,
            'entries': entries,
        })

    def _fs_collect(self):
        """GET /fs/collect?path=<dir>
        Walk a built-site directory and return every file as base64 + a guessed
        content-type, so the authenticated shell can upload the site to the
        cafresohq_state canister's site host (Publish-to-Canister). Same
        allowed-dirs guard as the other /fs endpoints (via _validate_path), so
        it can't read outside the sandbox. Skips hidden/.url files and files
        larger than the canister's ~2 MiB per-file cap (reported in `skipped`).
        Response: {root, files:[{path, contentType, size, b64}], skipped:[...]}
        """
        import mimetypes, base64
        qs = urllib.parse.urlparse(self.path).query
        req_path = (urllib.parse.parse_qs(qs).get('path') or [''])[0].strip()
        if not req_path:
            return self._send_json(400, {'error': 'path required'})
        try:
            root = self._validate_path(req_path)
        except PermissionError as e:
            return self._send_json(403, {'error': str(e)})
        except Exception as e:
            return self._send_json(400, {'error': f'invalid path: {e}'})
        if not root.is_dir():
            return self._send_json(400, {'error': 'not a directory'})

        MAX_FILE = 2_000_000       # matches cafresohq_state MAX_SITE_FILE_BYTES
        MAX_FILES = 300            # sanity cap on a single publish
        files, skipped = [], []
        for dirpath, dirnames, filenames in os.walk(root):
            # prune hidden dirs (node_modules-style publishing is out of scope for MVP)
            dirnames[:] = [d for d in dirnames if not d.startswith('.')]
            for name in filenames:
                if name.startswith('.') or name.lower().endswith('.url'):
                    continue
                full = pathlib.Path(dirpath) / name
                try:
                    rel = full.relative_to(root).as_posix()
                    size = full.stat().st_size
                except Exception:
                    continue
                if size > MAX_FILE:
                    skipped.append({'path': rel, 'reason': 'over 2 MiB'})
                    continue
                if len(files) >= MAX_FILES:
                    skipped.append({'path': rel, 'reason': 'file limit reached'})
                    continue
                ctype = mimetypes.guess_type(name)[0] or 'application/octet-stream'
                if ctype.startswith('text/') or ctype in ('application/javascript', 'application/json'):
                    ctype = ctype + '; charset=utf-8'
                try:
                    data = full.read_bytes()
                except Exception as e:
                    skipped.append({'path': rel, 'reason': str(e)})
                    continue
                files.append({
                    'path': rel,
                    'contentType': ctype,
                    'size': size,
                    'b64': base64.b64encode(data).decode('ascii'),
                })
        return self._send_json(200, {'root': str(root), 'files': files, 'skipped': skipped})

    def _fs_file(self):
        """GET /fs/file?path=<file>
        Serve a single file's raw bytes with a guessed content-type, INLINE so the
        browser renders it directly — the enabler for the artifact preview pane
        (HTML / PDF / image / SVG built by an agent). Read-only. Same allowed-dirs
        gate as /fs/browse (container mode), so it can't traverse outside the
        sandbox. Capped at 50 MiB to keep a runaway file from hanging the preview.
        """
        import mimetypes
        qs = urllib.parse.urlparse(self.path).query
        req_path = (urllib.parse.parse_qs(qs).get('path') or [''])[0].strip()
        if not req_path:
            return self._send_json(400, {'error': 'path required'})
        try:
            p = pathlib.Path(_client_path(req_path)).resolve()
        except Exception as e:
            return self._send_json(400, {'error': f'invalid path: {e}'})
        if not p.is_file():
            return self._send_json(404, {'error': 'not a file'})
        # Serve only within CAFRESOHQ_ALLOWED_DIRS — every mode (was container-
        # only, which exposed unauthenticated arbitrary file read in local/BYO).
        if not _within_allowed_dirs(p):
            return self._send_json(403, {'error': 'path is outside CAFRESOHQ_ALLOWED_DIRS',
                                         'allowed': _cafresohq_allowed_dirs})
        try:
            if p.stat().st_size > 50 * 1024 * 1024:
                return self._send_json(413, {'error': 'file too large to preview (>50 MiB)'})
            data = p.read_bytes()
        except PermissionError:
            return self._send_json(403, {'error': 'permission denied'})
        except Exception as e:
            return self._send_json(500, {'error': f'read failed: {e}'})
        ctype = mimetypes.guess_type(str(p))[0] or 'application/octet-stream'
        import hashlib as _hl
        fhash = _hl.sha1(data).hexdigest()[:16]
        try:
            fmtime = int(p.stat().st_mtime)
        except Exception:
            fmtime = 0
        self.send_response(200)
        self.send_header('Content-Type', ctype)
        self.send_header('Content-Length', str(len(data)))
        self.send_header('Content-Disposition', 'inline; filename="%s"' % p.name.replace('"', ''))
        # Conflict-safety: the editor captures these on read and re-checks before
        # saving / on an agent write, so co-editing never silently clobbers.
        self.send_header('X-File-Mtime', str(fmtime))
        self.send_header('X-File-Hash', fhash)
        self.send_header('Access-Control-Expose-Headers', 'X-File-Mtime, X-File-Hash')
        self.send_header('Cache-Control', 'no-store')
        self.send_header('Access-Control-Allow-Origin', '*')
        self.end_headers()
        try:
            self.wfile.write(data)
        except Exception:
            pass

    def _fs_stat(self):
        """GET /fs/stat?path=<file>  ->  {ok, mtime, hash, size}
        Cheap conflict-check for the editor (no content body): compare against the
        X-File-Mtime / X-File-Hash captured at read time before saving or on an
        agent write, so co-editing never silently clobbers. Same allowed-dirs
        guard as /fs/file; hash matches /fs/file's (sha1, first 16 hex)."""
        import hashlib as _hl
        qs = urllib.parse.urlparse(self.path).query
        req_path = (urllib.parse.parse_qs(qs).get('path') or [''])[0].strip()
        if not req_path:
            return self._send_json(400, {'error': 'path required'})
        try:
            p = pathlib.Path(_client_path(req_path)).resolve()
        except Exception as e:
            return self._send_json(400, {'error': f'invalid path: {e}'})
        if not _within_allowed_dirs(p):
            return self._send_json(403, {'error': 'path is outside CAFRESOHQ_ALLOWED_DIRS'})
        if not p.is_file():
            return self._send_json(404, {'ok': False, 'error': 'not a file'})
        try:
            st = p.stat()
            h = _hl.sha1(p.read_bytes()).hexdigest()[:16]
        except Exception as e:
            return self._send_json(500, {'error': str(e)})
        return self._send_json(200, {'ok': True, 'mtime': int(st.st_mtime), 'hash': h, 'size': st.st_size})

    def _fs_site(self):
        """GET /fs/site/<b64root>/<relpath>
        Serve a multi-file site from a directory so an agent-built page's
        RELATIVE refs (<link href=styles.css>, <script src=app.js>, <img
        src=assets/x.png>) resolve — the preview pane points an iframe at
        index.html here and the browser fetches every sibling/sub-asset back
        through the same /fs/site/<b64root>/ prefix. <b64root> is urlsafe
        base64 of the absolute site root; <relpath> is resolved under it. Both
        the root and the final file are re-checked against CAFRESOHQ_ALLOWED_DIRS
        via _validate_path, so it can't escape the sandbox. Read-only, embeddable
        (no X-Frame-Options), 50 MiB/file cap. An empty relpath or a directory
        falls back to index.html.
        """
        import base64, mimetypes
        rest = self.path[len('/fs/site/'):]
        rest = rest.split('?', 1)[0].split('#', 1)[0]
        if '/' in rest:
            b64root, relpath = rest.split('/', 1)
        else:
            b64root, relpath = rest, ''
        if not b64root:
            return self._send_json(400, {'error': 'site root required'})
        try:
            pad = '=' * (-len(b64root) % 4)
            root = base64.urlsafe_b64decode(b64root + pad).decode('utf-8')
        except Exception as e:
            return self._send_json(400, {'error': f'bad site root: {e}'})
        rel = urllib.parse.unquote(relpath).strip('/')
        if not rel:
            rel = 'index.html'
        try:
            root_p = self._validate_path(root)
        except PermissionError as e:
            return self._send_json(403, {'error': str(e)})
        except Exception as e:
            return self._send_json(400, {'error': f'invalid root: {e}'})
        try:
            target = self._validate_path(str(root_p / rel))
        except PermissionError as e:
            return self._send_json(403, {'error': str(e)})
        except Exception as e:
            return self._send_json(400, {'error': f'invalid path: {e}'})
        if target.is_dir():
            try:
                target = self._validate_path(str(target / 'index.html'))
            except Exception:
                return self._send_json(404, {'error': 'no index.html in directory'})
        if not target.is_file():
            return self._send_json(404, {'error': 'not found'})
        try:
            if target.stat().st_size > 50 * 1024 * 1024:
                return self._send_json(413, {'error': 'file too large to preview (>50 MiB)'})
            data = target.read_bytes()
        except PermissionError:
            return self._send_json(403, {'error': 'permission denied'})
        except Exception as e:
            return self._send_json(500, {'error': f'read failed: {e}'})
        ctype = mimetypes.guess_type(str(target))[0] or 'application/octet-stream'
        self.send_response(200)
        self.send_header('Content-Type', ctype)
        self.send_header('Content-Length', str(len(data)))
        self.send_header('Cache-Control', 'no-store')
        self.send_header('Access-Control-Allow-Origin', '*')
        self.end_headers()
        try:
            self.wfile.write(data)
        except Exception:
            pass

    def _fs_upload(self):
        """POST /fs/upload?path=<dir>   (multipart/form-data)
        Drop files into a project's working tree so the agents (FILE_READ /
        DIR_LIST) and the preview pane can immediately use them — the
        user-facing half of the "share files with agents" loop. Writes into
        the REAL filesystem, guarded by the same CAFRESOHQ_ALLOWED_DIRS
        whitelist as FILE_WRITE (via _validate_path), so it can't escape the
        sandbox. Filenames are flattened + sanitized; the target dir is created
        if missing; 50 MiB cap per request. ?path= defaults to the first
        allowed dir.
        """
        # Same enable-gate as /tools/exec FILE_WRITE: allow in unrestricted
        # local dev, otherwise require an explicit allow-list.
        _unrestricted_local = (_RUNTIME_ENV == 'local' and not _ALLOWED_DIRS_EXPLICIT)
        if not _cafresohq_allowed_dirs and not _unrestricted_local:
            return self._send_json(503, {'error': 'CAFRESOHQ_ALLOWED_DIRS not set — upload disabled'})

        ctype = self.headers.get('content-type', '')
        if 'multipart/form-data' not in ctype:
            return self._send_json(400, {'error': 'expected multipart/form-data'})
        length = int(self.headers.get('content-length', 0) or 0)
        if length <= 0:
            return self._send_json(400, {'error': 'empty upload'})
        if length > 50 * 1024 * 1024:
            return self._send_json(413, {'error': 'upload too large (50 MiB max)'})

        qs = urllib.parse.parse_qs(urllib.parse.urlparse(self.path).query)
        raw_dir = (qs.get('path') or [''])[0].strip()
        if not raw_dir:
            raw_dir = _cafresohq_allowed_dirs[0] if _cafresohq_allowed_dirs else os.getcwd()
        try:
            target_dir = self._validate_path(raw_dir)
        except PermissionError as e:
            return self._send_json(403, {'error': str(e)})
        except Exception as e:
            # Malformed path (e.g. embedded null byte) — resolve() raises
            # ValueError. Mirror _fs_browse/_fs_file: clean 400, not a dropped
            # connection.
            return self._send_json(400, {'error': f'invalid path: {e}'})
        if target_dir.exists() and not target_dir.is_dir():
            return self._send_json(400, {'error': 'target path is not a directory'})

        raw = self.rfile.read(length)
        import email.parser as _ep
        import re as _re
        msg = _ep.BytesParser().parsebytes(
            b'Content-Type: ' + ctype.encode('latin-1', 'replace') + b'\r\n\r\n' + raw)
        if not msg.is_multipart():
            return self._send_json(400, {'error': 'bad multipart body'})

        try:
            target_dir.mkdir(parents=True, exist_ok=True)
        except Exception as e:
            return self._send_json(500, {'error': f'mkdir failed: {e}'})

        saved, errors = [], []
        for part in msg.get_payload():
            fname = part.get_filename()
            if not fname:
                continue
            # Flatten any path components, then strip anything but a safe set.
            fname = fname.replace('\\', '/').split('/')[-1]
            fname = _re.sub(r'[^\w .()\[\]\-]+', '_', fname).strip()
            if not fname or fname.startswith('.'):
                continue
            dest = (target_dir / fname).resolve()
            # Defense-in-depth: re-assert the whitelist on the final path even
            # though a sanitized basename can't traverse.
            try:
                self._validate_path(str(dest))
            except PermissionError:
                errors.append({'path': fname, 'error': 'outside allowed dirs'})
                continue
            data = part.get_payload(decode=True) or b''
            try:
                dest.write_bytes(data)
                saved.append({'path': str(dest), 'name': fname, 'size': len(data)})
            except Exception as e:
                errors.append({'path': fname, 'error': str(e)})
        if not saved and errors:
            return self._send_json(500, {'error': errors[0]['error'], 'failed': errors})
        return self._send_json(200, {'uploaded': saved, 'count': len(saved),
                                     'dir': str(target_dir),
                                     **({'failed': errors} if errors else {})})

    # ---- Working-tree file manager (mkdir / rename / delete) -------------
    # All three share FILE_WRITE's posture: open in unrestricted local dev,
    # else confined to CAFRESOHQ_ALLOWED_DIRS via _validate_path (resolve +
    # .relative_to). They make the container feel like an OS the user owns.
    def _fs_mutate_ok(self):
        _unrestricted_local = (_RUNTIME_ENV == 'local' and not _ALLOWED_DIRS_EXPLICIT)
        if not _cafresohq_allowed_dirs and not _unrestricted_local:
            self._send_json(503, {'error': 'CAFRESOHQ_ALLOWED_DIRS not set — file ops disabled'})
            return False
        return True

    def _fs_json_body(self):
        length = int(self.headers.get('content-length', 0) or 0)
        try:
            return json.loads(self.rfile.read(length).decode('utf-8') or '{}')
        except Exception:
            return None

    def _fs_mkdir(self):
        """POST /fs/mkdir  {path}  — create a folder in the working tree."""
        if not self._fs_mutate_ok():
            return
        req = self._fs_json_body()
        if req is None:
            return self._send_json(400, {'error': 'bad json'})
        raw = (req.get('path') or '').strip()
        if not raw:
            return self._send_json(400, {'error': 'path required'})
        try:
            target = self._validate_path(raw)
        except PermissionError as e:
            return self._send_json(403, {'error': str(e)})
        except Exception as e:
            return self._send_json(400, {'error': f'invalid path: {e}'})
        if target.exists():
            if target.is_dir():
                return self._send_json(200, {'ok': True, 'path': str(target), 'existed': True})
            return self._send_json(409, {'error': 'a file with that name exists'})
        try:
            target.mkdir(parents=True, exist_ok=True)
        except (NotADirectoryError, FileExistsError) as e:
            return self._send_json(409, {'error': f'cannot create folder here: {e}'})
        except Exception as e:
            return self._send_json(500, {'error': str(e)})
        return self._send_json(200, {'ok': True, 'path': str(target)})

    def _fs_rename(self):
        """POST /fs/rename  {from, to}  — rename/move within the working tree."""
        if not self._fs_mutate_ok():
            return
        req = self._fs_json_body()
        if req is None:
            return self._send_json(400, {'error': 'bad json'})
        src = (req.get('from') or '').strip()
        dst = (req.get('to') or '').strip()
        if not src or not dst:
            return self._send_json(400, {'error': 'from and to required'})
        try:
            sp = self._validate_path(src)
            dp = self._validate_path(dst)
        except PermissionError as e:
            return self._send_json(403, {'error': str(e)})
        except Exception as e:
            return self._send_json(400, {'error': f'invalid path: {e}'})
        if not sp.exists():
            return self._send_json(404, {'error': 'source not found'})
        # Never move a workspace root itself (parallels _fs_delete; matters when
        # more than one allowed dir is configured).
        roots = {str(pathlib.Path(d).resolve()) for d in _cafresohq_allowed_dirs}
        if str(sp) in roots:
            return self._send_json(403, {'error': 'refusing to move a workspace root'})
        if dp.exists():
            return self._send_json(409, {'error': 'target already exists'})
        try:
            dp.parent.mkdir(parents=True, exist_ok=True)
            os.replace(str(sp), str(dp))
        except (NotADirectoryError, FileExistsError) as e:
            return self._send_json(409, {'error': f'cannot move there: {e}'})
        except Exception as e:
            return self._send_json(500, {'error': str(e)})
        return self._send_json(200, {'ok': True, 'from': str(sp), 'to': str(dp)})

    def _fs_delete(self):
        """POST /fs/delete  {path}  — delete a file or directory (recursive)."""
        if not self._fs_mutate_ok():
            return
        req = self._fs_json_body()
        if req is None:
            return self._send_json(400, {'error': 'bad json'})
        raw = (req.get('path') or '').strip()
        if not raw:
            return self._send_json(400, {'error': 'path required'})
        try:
            target = self._validate_path(raw)
        except PermissionError as e:
            return self._send_json(403, {'error': str(e)})
        except Exception as e:
            return self._send_json(400, {'error': f'invalid path: {e}'})
        if not target.exists() and not target.is_symlink():
            return self._send_json(404, {'error': 'not found'})
        # Never delete an allowed-dir root itself.
        roots = {str(pathlib.Path(d).resolve()) for d in _cafresohq_allowed_dirs}
        if str(target) in roots:
            return self._send_json(403, {'error': 'refusing to delete a workspace root'})
        try:
            # Follow-the-link guard: a symlinked dir is unlinked (remove the
            # link), never rmtree'd (which would wipe the link's target).
            if target.is_dir() and not target.is_symlink():
                shutil.rmtree(str(target))
            else:
                target.unlink()
        except Exception as e:
            return self._send_json(500, {'error': str(e)})
        return self._send_json(200, {'ok': True, 'deleted': str(target)})

    # ---- Project Terminal (Claude Code / Codex CLI runner) ---------------
    def _terminal_spawn(self):
        """GET /terminal/spawn?cli=claude|codex&cwd=<path>
        Opens the CLI in a new OS terminal window so the user sees the
        full interactive TUI (welcome screen, coloured prompts, etc.).
        Windows: tries Windows Terminal → pwsh → PowerShell → cmd.
        Linux/macOS: tries x-terminal-emulator → gnome-terminal → xterm.
        """
        qs = urllib.parse.urlparse(self.path).query
        params = urllib.parse.parse_qs(qs)
        cli  = (params.get('cli')  or ['claude'])[0].strip().lower()
        cwd  = (params.get('cwd')  or ['']      )[0].strip()
        if cli not in ('claude', 'codex'):
            return self._send_json(400, {'error': 'cli must be claude or codex'})
        if not cwd:
            return self._send_json(400, {'error': 'cwd required'})
        cwd_path = pathlib.Path(_client_path(cwd)).resolve()
        if not cwd_path.is_dir():
            return self._send_json(400, {'error': f'directory not found: {cwd}'})
        bin_ = (self._claudecode_resolve() if cli == 'claude' else self._codex_resolve())
        if not bin_:
            return self._send_json(503, {'error': f'{cli} CLI not found'})
        try:
            if sys.platform == 'win32':
                wt  = shutil.which('wt')
                ps  = shutil.which('pwsh') or shutil.which('powershell')
                if wt and ps:
                    subprocess.Popen(
                        ['wt', '-d', str(cwd_path), ps, '-NoExit', '-Command', cli],
                        creationflags=subprocess.DETACHED_PROCESS | subprocess.CREATE_NEW_PROCESS_GROUP,
                    )
                elif ps:
                    subprocess.Popen(
                        [ps, '-NoExit', '-Command',
                         f'Set-Location "{cwd_path}"; {cli}'],
                        creationflags=subprocess.CREATE_NEW_CONSOLE | subprocess.DETACHED_PROCESS,
                    )
                else:
                    subprocess.Popen(
                        ['cmd', '/k', f'cd /d "{cwd_path}" && {cli}'],
                        creationflags=subprocess.CREATE_NEW_CONSOLE | subprocess.DETACHED_PROCESS,
                    )
            else:
                launched = False
                for term in ['x-terminal-emulator', 'gnome-terminal', 'kitty', 'xterm']:
                    if shutil.which(term):
                        subprocess.Popen([term, '-e', cli], cwd=str(cwd_path))
                        launched = True
                        break
                if not launched:
                    return self._send_json(503, {'error': 'no terminal emulator found'})
        except Exception as e:
            return self._send_json(500, {'error': f'spawn failed: {e}'})
        return self._send_json(200, {'ok': True, 'cli': cli, 'cwd': str(cwd_path)})

    def _app_origins(self):
        """Browser origins permitted to open the terminal WebSocket and fetch the
        PTY nonce: same-origin, the production gateway, localhost dev, plus any
        canister origins configured via CAFRESOHQ_ALLOWED_WS_ORIGINS (cross-origin
        frontend/backend split)."""
        host   = self.headers.get('Host', '').strip()
        scheme = 'https' if isinstance(self.connection, ssl.SSLSocket) else 'http'
        origins = {
            'https://hq.cafreso.com',        # production Caddy gateway
            'https://hq-ui.cafreso.com',     # canister UI shell (cross-origin split)
            'https://ai.cafreso.com',        # SvelteKit frontend
            'http://localhost:8787',         # local dev (same-origin)
            'http://127.0.0.1:8787',         # local dev (alt)
            'http://localhost:5173',         # vite dev server
            'http://localhost:5174',         # vite dev server (alt)
            f'{scheme}://{host}',            # same-origin (any host the client used)
        }
        return origins | _extra_app_origins

    def _terminal_nonce(self):
        """GET /terminal/nonce → {"nonce": "<hex>"}
        Returns the per-process nonce that /terminal/pty requires as a query
        param. Same-origin callers (no Origin header on same-origin XHR, or an
        Origin in the allowlist) get it; a cross-origin Origin that isn't in the
        allowlist is refused, so only the HQ app — same-origin or an approved
        canister origin (CAFRESOHQ_ALLOWED_WS_ORIGINS) — can obtain it.
        """
        _origin = self.headers.get('Origin', '').strip()
        if _origin and _origin not in self._app_origins():
            return self._send_json(403, {'error': f'origin not allowed: {_origin}'})
        return self._send_json(200, {'nonce': _PTY_NONCE})

    def _terminal_pty_ws(self):
        """GET /terminal/pty?cli=claude|codex&cwd=<path>&cols=120&rows=30
        Upgrades to WebSocket, spawns a PTY running the CLI, and bridges
        stdin/stdout between the browser (xterm.js) and the process.

        Windows: uses pywinpty (existing behaviour, unchanged).
        Linux/macOS: uses Python stdlib pty + subprocess (no extra packages).

        Browser → server frames: raw keystroke bytes (text or binary frames),
        or a JSON resize message: {"type":"resize","cols":N,"rows":N}.
        Server → browser frames: binary frames containing PTY output bytes.
        """
        _is_win = sys.platform == 'win32'
        if _is_win:
            try:
                import winpty as _winpty
            except ImportError:
                return self._send_json(503, {'error': 'pywinpty not installed — run: pip install pywinpty'})

        qs = urllib.parse.urlparse(self.path).query
        params = urllib.parse.parse_qs(qs)
        cli  = (params.get('cli')  or ['claude'])[0].strip().lower()
        cwd  = (params.get('cwd')  or ['']      )[0].strip()
        cols = max(10, min(500, int((params.get('cols') or ['120'])[0])))
        rows = max(5,  min(200, int((params.get('rows') or ['30'] )[0])))

        if cli not in ('claude', 'codex', 'hermes', 'gemini'):
            return self._send_json(400, {'error': 'cli must be claude, codex, hermes, or gemini'})
        if not cwd:
            return self._send_json(400, {'error': 'cwd required'})
        cwd_path = pathlib.Path(_client_path(cwd)).resolve()
        if not cwd_path.is_dir():
            return self._send_json(400, {'error': f'directory not found: {cwd}'})
        if cli == 'claude':
            bin_ = self._claudecode_resolve()
        elif cli == 'codex':
            bin_ = self._codex_resolve()
        elif cli == 'gemini':
            bin_ = self._gemini_resolve()
        else:  # hermes — the default agent (unix-only; native here or via WSL)
            bin_ = self._hermes_resolve()
        if not bin_:
            return self._send_json(503, {'error': f'{cli} CLI not found'})
        # Hermes opens its interactive agent via the `chat` subcommand; the
        # workspace cwd scopes the agent to that project. Others just exec.
        cli_extra_args = ['chat'] if cli == 'hermes' else []

        # ── Security: Origin + nonce checks ────────────────────────────────
        # Origin validation — reject cross-origin WS initiations.
        # Browsers always send Origin on WebSocket upgrade; non-browser clients
        # may omit it (curl, CLI tools, unit tests) — those are allowed through
        # so local dev tooling isn't broken.
        _origin = self.headers.get('Origin', '').strip()
        if _origin and _origin not in self._app_origins():
            return self.send_error(403, f'WebSocket origin not allowed: {_origin}')

        # Nonce validation — the frontend fetches /terminal/nonce first and
        # appends ?nonce=<value> to the WS URL.  Any connection without the
        # correct nonce is rejected (attacker can't fetch the nonce cross-origin
        # because /terminal/nonce is same-origin-only for browser XHR).
        _provided_nonce = (params.get('nonce') or [''])[0]
        if not _provided_nonce or not secrets.compare_digest(_provided_nonce, _PTY_NONCE):
            return self.send_error(403, 'Missing or invalid nonce — fetch /terminal/nonce first')

        # ── WebSocket handshake ─────────────────────────────────────────────
        ws_key = self.headers.get('Sec-WebSocket-Key', '')
        if not ws_key:
            return self.send_error(400, 'missing Sec-WebSocket-Key')
        accept = base64.b64encode(
            hashlib.sha1((ws_key + self._WS_GUID).encode()).digest()
        ).decode()
        self.wfile.write((
            'HTTP/1.1 101 Switching Protocols\r\n'
            'Upgrade: websocket\r\n'
            'Connection: Upgrade\r\n'
            f'Sec-WebSocket-Accept: {accept}\r\n\r\n'
        ).encode('latin-1'))
        self.wfile.flush()
        self.close_connection = True
        client_sock = self.connection
        client_sock.settimeout(None)

        # ── Build PTY environment ───────────────────────────────────────────
        import copy as _copy
        agent_env = _copy.deepcopy(os.environ)
        for _d in [r'C:\Program Files\Git\usr\bin',
                   r'C:\Program Files\Git\bin',
                   r'C:\Program Files\Git\mingw64\bin']:
            if os.path.isdir(_d) and _d not in agent_env.get('PATH', ''):
                agent_env['PATH'] = _d + os.pathsep + agent_env.get('PATH', '')

        def _ws_recv_frame(sock):
            """Read one WebSocket frame; return (opcode, payload) or (None,None)."""
            def _read(n):
                buf = b''
                while len(buf) < n:
                    chunk = sock.recv(n - len(buf))
                    if not chunk:
                        raise ConnectionError('ws closed')
                    buf += chunk
                return buf
            try:
                b0, b1 = struct.unpack('!BB', _read(2))
                opcode  = b0 & 0x0F
                masked  = bool(b1 & 0x80)
                length  = b1 & 0x7F
                if length == 126:
                    (length,) = struct.unpack('!H', _read(2))
                elif length == 127:
                    (length,) = struct.unpack('!Q', _read(8))
                mask = _read(4) if masked else b'\x00\x00\x00\x00'
                payload = bytearray(_read(length))
                if masked:
                    for i in range(len(payload)):
                        payload[i] ^= mask[i % 4]
                return opcode, bytes(payload)
            except Exception:
                return None, None

        # ── Optional init frame — inject BYOK keys before spawning ──────────
        # The browser sends {"type":"init","anthropic_key":"...","openai_key":"..."}
        # as the very first WebSocket frame (in ws.onopen) so the shell inherits
        # the user's stored API keys automatically.  If the frame doesn't arrive
        # within 2 s (or the user has no stored keys), we just proceed without it.
        # Server-side env vars always win: BYOK is only injected when the
        # container env is blank (same rule as /terminal/stream).
        pending_frame = None   # replayed into ws_to_pty if it isn't an init frame
        client_sock.settimeout(2.0)
        _init_opcode, _init_payload = _ws_recv_frame(client_sock)
        client_sock.settimeout(None)
        if _init_opcode in (0x1, 0x2) and _init_payload:
            try:
                _init_msg = json.loads(_init_payload)
                if isinstance(_init_msg, dict) and _init_msg.get('type') == 'init':
                    _pty_auth = (_init_msg.get('auth_method') or '').strip().lower()
                    _ak = (_init_msg.get('anthropic_key') or '').strip()
                    _ok = (_init_msg.get('openai_key')    or '').strip()
                    _gk = (_init_msg.get('gemini_key')    or '').strip()
                    if _pty_auth == 'subscription':
                        # Strip API key so CLI uses its own OAuth login
                        agent_env.pop('ANTHROPIC_API_KEY', None)
                    elif _ak and not agent_env.get('ANTHROPIC_API_KEY', '').strip():
                        agent_env['ANTHROPIC_API_KEY'] = _ak
                    if _ok and not agent_env.get('OPENAI_API_KEY', '').strip():
                        agent_env['OPENAI_API_KEY'] = _ok
                    # Gemini CLI reads GEMINI_API_KEY (and GOOGLE_API_KEY as a
                    # fallback) — set both so either lookup path works.
                    if _gk and not agent_env.get('GEMINI_API_KEY', '').strip():
                        agent_env['GEMINI_API_KEY'] = _gk
                        agent_env.setdefault('GOOGLE_API_KEY', _gk)
                    del _ak, _ok, _gk, _pty_auth  # clear plaintext key strings from Python locals
                    del _init_msg, _init_payload  # clear raw frame bytes
                else:
                    pending_frame = (_init_opcode, _init_payload)  # not init — replay
                    del _init_msg
            except (ValueError, TypeError):
                pending_frame = (_init_opcode, _init_payload)      # not JSON — replay
        elif _init_opcode is not None:
            pending_frame = (_init_opcode, _init_payload)          # non-text — replay

        # ── WebSocket frame helper (server → client, unmasked) ─────────────
        def _ws_send_raw(sock, data):
            if isinstance(data, str):
                data = data.encode('utf-8')
            n = len(data)
            if n < 126:
                hdr = struct.pack('!BB', 0x82, n)
            elif n < 65536:
                hdr = struct.pack('!BBH', 0x82, 126, n)
            else:
                hdr = struct.pack('!BBQ', 0x82, 127, n)
            sock.sendall(hdr + data)   # raises OSError on failure; callers handle it

        # ── Session resume or new spawn ─────────────────────────────────────
        session_id = (params.get('session_id') or [''])[0].strip()

        with _PTY_SESSIONS_LK:
            sess = _PTY_SESSIONS.get(session_id) if session_id else None

        if sess is not None and not sess['stop'].is_set():
            # ── Reconnect to existing PTY session ──────────────────────────
            with sess['sock_lk']:
                sess['sock']    = client_sock
                sess['expires'] = None   # cancel reaper countdown

            # Replay buffered output collected while the client was away.
            with sess['buf_lk']:
                replay, sess['buf'] = bytes(sess['buf']), bytearray()
            if replay:
                try:
                    _ws_send_raw(client_sock, replay)
                except OSError:
                    pass

            sys.stderr.write(f'[pty-ws] reconnected: {session_id[:8]}… {cli} @ {cwd_path}\n')

        else:
            # ── Spawn a fresh PTY ───────────────────────────────────────────
            stop_evt = threading.Event()
            sess = {
                'stop':    stop_evt,
                'sock':    client_sock,
                'sock_lk': threading.Lock(),
                'buf':     bytearray(),
                'buf_lk':  threading.Lock(),
                'expires': None,
            }
            if session_id:
                with _PTY_SESSIONS_LK:
                    _PTY_SESSIONS[session_id] = sess

            if _is_win:
                spawn_argv = ['cmd.exe', '/c', bin_ or cli] + cli_extra_args
                try:
                    pty_proc = _winpty.PtyProcess.spawn(
                        spawn_argv, cwd=str(cwd_path), env=agent_env,
                        dimensions=(rows, cols),
                    )
                except Exception as exc:
                    try: _ws_send_raw(client_sock, f'\r\n\x1b[31mFailed to spawn {cli}: {exc}\x1b[0m\r\n')
                    except OSError: pass
                    if session_id:
                        with _PTY_SESSIONS_LK: _PTY_SESSIONS.pop(session_id, None)
                    return
                sess['pty_proc'] = pty_proc

                def pty_to_ws():
                    while not sess['stop'].is_set():
                        try:
                            data = pty_proc.read(4096)
                        except EOFError:
                            data = None
                        except Exception:
                            data = None if not pty_proc.isalive() else b''
                        if data is None:
                            break
                        if not data:
                            continue
                        raw = data.encode('utf-8') if isinstance(data, str) else data
                        with sess['sock_lk']:
                            sock = sess['sock']
                        if sock:
                            try:
                                _ws_send_raw(sock, raw)
                                continue
                            except OSError:
                                with sess['sock_lk']:
                                    if sess['sock'] is sock:
                                        sess['sock']    = None
                                        sess['expires'] = time.time() + _PTY_SESSION_TTL
                        with sess['buf_lk']:
                            sess['buf'] += raw
                            if len(sess['buf']) > _PTY_BUF_CAP:
                                sess['buf'] = sess['buf'][-_PTY_BUF_CAP:]
                    # PTY exited — notify client and clean up.
                    sess['stop'].set()
                    if session_id:
                        with _PTY_SESSIONS_LK: _PTY_SESSIONS.pop(session_id, None)
                    with sess['sock_lk']:
                        sock = sess['sock']
                    if sock:
                        try: _ws_send_raw(sock, b'\r\n\x1b[2m[process exited]\x1b[0m\r\n')
                        except OSError: pass
                        try: sock.shutdown(socket.SHUT_RDWR)
                        except OSError: pass
                    sys.stderr.write(f'[pty-ws] PTY exited: {session_id[:8] if session_id else "?"} {cli}\n')

            else:
                import pty as _pty, fcntl, struct as _struct2, termios, select as _select
                master_fd, slave_fd = _pty.openpty()
                fcntl.ioctl(slave_fd, termios.TIOCSWINSZ,
                            _struct2.pack('HHHH', rows, cols, 0, 0))
                spawn_argv = [bin_ or cli] + cli_extra_args
                try:
                    proc = subprocess.Popen(
                        spawn_argv,
                        stdin=slave_fd, stdout=slave_fd, stderr=slave_fd,
                        close_fds=True, start_new_session=True,
                        cwd=str(cwd_path), env=agent_env,
                    )
                except Exception as exc:
                    os.close(slave_fd); os.close(master_fd)
                    try: _ws_send_raw(client_sock, f'\r\n\x1b[31mFailed to spawn {cli}: {exc}\x1b[0m\r\n')
                    except OSError: pass
                    if session_id:
                        with _PTY_SESSIONS_LK: _PTY_SESSIONS.pop(session_id, None)
                    return
                os.close(slave_fd)
                sess['proc']      = proc
                sess['master_fd'] = master_fd

                def pty_to_ws():
                    while not sess['stop'].is_set():
                        try:
                            r, _, _ = _select.select([master_fd], [], [], 0.05)
                            if r:
                                data = os.read(master_fd, 4096)
                                if not data:
                                    break
                            elif proc.poll() is not None:
                                break
                            else:
                                continue
                        except OSError:
                            break
                        with sess['sock_lk']:
                            sock = sess['sock']
                        if sock:
                            try:
                                _ws_send_raw(sock, data)
                                continue
                            except OSError:
                                with sess['sock_lk']:
                                    if sess['sock'] is sock:
                                        sess['sock']    = None
                                        sess['expires'] = time.time() + _PTY_SESSION_TTL
                        with sess['buf_lk']:
                            sess['buf'] += data
                            if len(sess['buf']) > _PTY_BUF_CAP:
                                sess['buf'] = sess['buf'][-_PTY_BUF_CAP:]
                    sess['stop'].set()
                    if session_id:
                        with _PTY_SESSIONS_LK: _PTY_SESSIONS.pop(session_id, None)
                    try: os.close(master_fd)
                    except OSError: pass
                    with sess['sock_lk']:
                        sock = sess['sock']
                    if sock:
                        try: _ws_send_raw(sock, b'\r\n\x1b[2m[process exited]\x1b[0m\r\n')
                        except OSError: pass
                        try: sock.shutdown(socket.SHUT_RDWR)
                        except OSError: pass
                    sys.stderr.write(f'[pty-ws] PTY exited: {session_id[:8] if session_id else "?"} {cli}\n')

            t_out = threading.Thread(target=pty_to_ws, daemon=True)
            t_out.start()

        # ── ws_to_pty — shared for both new and reconnected sessions ────────
        # Reads frames from this WS connection and forwards them to the PTY.
        # On disconnect, detaches the socket from the session WITHOUT killing
        # the PTY so the client can reconnect and resume.
        if _is_win:
            def ws_to_pty():
                nonlocal pending_frame
                _pf, pending_frame = pending_frame, None
                _pf_queue = [_pf] if _pf is not None else []
                while not sess['stop'].is_set():
                    if _pf_queue:
                        opcode, payload = _pf_queue.pop(0)
                    else:
                        opcode, payload = _ws_recv_frame(client_sock)
                    if opcode is None or opcode == 0x8:
                        break
                    if opcode == 0x9:   # ping → pong
                        try: client_sock.sendall(struct.pack('!BB', 0x8A, len(payload)) + payload)
                        except OSError: break
                        continue
                    if opcode in (0x1, 0x2) and payload:
                        try:
                            msg = json.loads(payload)
                            if isinstance(msg, dict) and msg.get('type') == 'resize':
                                c = max(10, min(500, int(msg.get('cols', cols))))
                                r = max(5,  min(200, int(msg.get('rows', rows))))
                                sess['pty_proc'].setwinsize(r, c)
                                continue
                        except (ValueError, TypeError):
                            pass
                        try:
                            sess['pty_proc'].write(payload.decode('utf-8', errors='replace'))
                        except Exception:
                            break
                # WS dropped — detach socket; PTY keeps running.
                with sess['sock_lk']:
                    if sess['sock'] is client_sock:
                        sess['sock']    = None
                        sess['expires'] = time.time() + _PTY_SESSION_TTL
        else:
            def ws_to_pty():
                import fcntl as _fcntl2, struct as _struct3, termios as _termios2
                nonlocal pending_frame
                _pf, pending_frame = pending_frame, None
                _pf_queue = [_pf] if _pf is not None else []
                while not sess['stop'].is_set():
                    if _pf_queue:
                        opcode, payload = _pf_queue.pop(0)
                    else:
                        opcode, payload = _ws_recv_frame(client_sock)
                    if opcode is None or opcode == 0x8:
                        break
                    if opcode == 0x9:
                        try: client_sock.sendall(_struct3.pack('!BB', 0x8A, len(payload)) + payload)
                        except OSError: break
                        continue
                    if opcode in (0x1, 0x2) and payload:
                        try:
                            msg = json.loads(payload)
                            if isinstance(msg, dict) and msg.get('type') == 'resize':
                                c = max(10, min(500, int(msg.get('cols', cols))))
                                r = max(5,  min(200, int(msg.get('rows', rows))))
                                _fcntl2.ioctl(sess['master_fd'], _termios2.TIOCSWINSZ,
                                              _struct3.pack('HHHH', r, c, 0, 0))
                                continue
                        except (ValueError, TypeError):
                            pass
                        try:
                            os.write(sess['master_fd'], payload)
                        except OSError:
                            break
                with sess['sock_lk']:
                    if sess['sock'] is client_sock:
                        sess['sock']    = None
                        sess['expires'] = time.time() + _PTY_SESSION_TTL

        t_in = threading.Thread(target=ws_to_pty, daemon=True)
        t_in.start()
        t_in.join()   # block until this WS connection closes
        sys.stderr.write(f'[pty-ws] WS detached: {session_id[:8] if session_id else "?"} {cli} @ {cwd_path}\n')

    def _terminal_status(self):
        """Report CLI availability, runtime environment, and spawn capability."""
        claude_bin = self._claudecode_resolve()
        codex_bin  = self._codex_resolve()
        hermes_bin = self._hermes_resolve()
        gemini_bin = self._gemini_resolve()
        # Spawn is only meaningful on local machines that have a display/GUI.
        spawn_ok = (_RUNTIME_ENV == 'local')
        try:
            import winpty as _wp  # noqa: F401
            pty_ok = True          # Windows + pywinpty installed
        except ImportError:
            # Linux / macOS: stdlib pty is always available
            pty_ok = (sys.platform != 'win32')
        return self._send_json(200, {
            'claude':          bool(claude_bin),
            'claudePath':      claude_bin or '',
            'codex':           bool(codex_bin),
            'codexPath':       codex_bin or '',
            'hermes':          bool(hermes_bin),
            'hermesPath':      hermes_bin or '',
            # hermesVia: how a Projects terminal would reach hermes — 'native'
            # when on PATH here; '' when unavailable (e.g. native Windows, where
            # the supported path is running the whole stack in WSL).
            'hermesVia':       ('native' if hermes_bin else ''),
            'gemini':          bool(gemini_bin),
            'geminiPath':      gemini_bin or '',
            'runtime_env':     _RUNTIME_ENV,
            'spawn_supported': spawn_ok,
            'pty_supported':   pty_ok,
        })

    def _terminal_stream(self):
        """Stream an agentic CLI session scoped to a project directory.

        Body: {messages, cli, cwd, model?, projectName?, authMethod?,
               claudeKey?, codexKey?}
          cli        — 'claude' | 'codex'
          cwd        — absolute path to the project directory
          messages   — [{role, content}] conversation history
          model      — optional model override
          authMethod — 'subscription' | 'apikey' (default varies by cli)
                       'subscription' strips any ANTHROPIC_API_KEY from the
                       subprocess env so the Claude CLI uses its own OAuth
                       credentials (Pro/Max plan).
          claudeKey  — BYOK: client's ANTHROPIC_API_KEY (AES-GCM decrypted client-side)
          codexKey   — BYOK: client's OPENAI_API_KEY (AES-GCM decrypted client-side)
        Server env vars always win over BYOK keys so fleet admins can lock
        the backend. A blank server key lets the client key through.
        SSE output uses the same {choices[0].delta.content} shape as
        /claudecode/stream. delta.type ('text'|'tool'|'error') lets the
        terminal UI colour-code different event kinds.
        """
        length = int(self.headers.get('content-length', 0) or 0)
        try:
            body = json.loads(self.rfile.read(length).decode('utf-8') or '{}')
        except json.JSONDecodeError:
            return self._send_json(400, {'error': 'bad json'})

        cli      = (body.get('cli') or 'claude').lower().strip()
        cwd      = (body.get('cwd') or '').strip()
        messages = body.get('messages') or []
        model    = (body.get('model') or '').strip()
        auth_method = (body.get('authMethod') or '').strip().lower()
        byok_claude = (body.get('claudeKey') or '').strip()
        byok_codex  = (body.get('codexKey') or '').strip()
        byok_gemini = (body.get('geminiKey') or '').strip()
        # Orchestrator session keys — used to tie multi-turn invocations to
        # the same project + CLI session so context (working dir, prior
        # output) can be threaded through. Logged for trace; future use:
        # pass to `claude --resume <session-id>` to recover stored CLI state.
        session_id  = (body.get('sessionId') or '').strip()[:80]
        project_id  = (body.get('projectId') or '').strip()[:80]

        if cli not in ('claude', 'codex', 'gemini'):
            return self._send_json(400, {'error': 'cli must be "claude", "codex", or "gemini"'})
        if not cwd:
            return self._send_json(400, {'error': 'cwd required'})
        cwd_path = pathlib.Path(_client_path(cwd)).resolve()
        if not cwd_path.is_dir():
            return self._send_json(400, {'error': f'directory not found: {cwd}'})

        history_lines = []
        for m in messages:
            role    = (m.get('role') or 'user').upper()
            content = (m.get('content') or '').strip()
            if content:
                history_lines.append(f'{role}: {content}')
        prompt = '\n\n'.join(history_lines)
        if not prompt:
            return self._send_json(400, {'error': 'no prompt'})

        import copy as _copy
        agent_env = _copy.deepcopy(os.environ)
        for _d in [r'C:\Program Files\Git\usr\bin',
                   r'C:\Program Files\Git\bin',
                   r'C:\Program Files\Git\mingw64\bin']:
            if os.path.isdir(_d) and _d not in agent_env.get('PATH', ''):
                agent_env['PATH'] = _d + os.pathsep + agent_env.get('PATH', '')

        if cli == 'claude':
            bin_ = self._claudecode_resolve()
            if not bin_:
                return self._send_json(503, {'error': 'claude CLI not found — install Claude Code or set CAFRESOHQ_CLAUDE_BIN'})
            if auth_method == 'subscription':
                # Strip any ANTHROPIC_API_KEY so the CLI falls through to its
                # own OAuth / subscription credentials (Pro/Max plan).
                agent_env.pop('ANTHROPIC_API_KEY', None)
            elif byok_claude and not agent_env.get('ANTHROPIC_API_KEY', '').strip():
                agent_env['ANTHROPIC_API_KEY'] = byok_claude
            cmd = [bin_, '--print',
                   '--output-format', 'stream-json',
                   '--verbose',
                   '--allowed-tools', 'Read,Glob,Grep,Bash,Edit,Write,WebFetch,WebSearch',
                   '--add-dir', str(cwd_path),
                   '--input-format', 'text']
            if model:
                cmd += ['--model', model]
        elif cli == 'gemini':
            bin_ = self._gemini_resolve()
            if not bin_:
                return self._send_json(503, {'error': 'gemini CLI not found — npm i -g @google/gemini-cli or set CAFRESOHQ_GEMINI_BIN'})
            # Auth mirrors the PTY path: 'subscription' = the gemini CLI's own
            # Google login (stored creds), otherwise inject the BYOK key into
            # GEMINI_API_KEY/GOOGLE_API_KEY (both lookups the CLI honours).
            if auth_method == 'subscription':
                agent_env.pop('GEMINI_API_KEY', None)
                agent_env.pop('GOOGLE_API_KEY', None)
            elif byok_gemini and not agent_env.get('GEMINI_API_KEY', '').strip():
                agent_env['GEMINI_API_KEY'] = byok_gemini
                agent_env.setdefault('GOOGLE_API_KEY', byok_gemini)
            # Non-interactive one-shot: `--prompt` runs once and exits, printing
            # the answer to stdout. `--yolo` auto-approves tool calls so the
            # agent can actually read/edit files in the project dir (parity with
            # claude --print's fixed allowed-tools); the project cwd scopes it.
            # The prompt goes in argv (not stdin) — gemini reads --prompt, so we
            # skip the stdin write below for it.
            cmd = [bin_, '--yolo', '--prompt', prompt]
            if model:
                cmd += ['--model', model]
        else:
            bin_ = self._codex_resolve()
            if not bin_:
                return self._send_json(503, {'error': 'codex CLI not found — npm i -g @openai/codex or set CAFRESOHQ_CODEX_BIN'})
            path_key = next((k for k in agent_env.keys() if k.lower() == 'path'), 'Path')
            path_value = agent_env.get(path_key, '')
            path_value = os.pathsep.join(
                p for p in path_value.split(os.pathsep)
                if p and r'\.codex\tmp\arg0' not in p.lower()
            )
            for _k in [k for k in list(agent_env.keys()) if k.lower() == 'path']:
                agent_env.pop(_k, None)
            agent_env['Path'] = path_value
            agent_env.pop('OPENAI_BASE_URL', None)

            if byok_codex and not agent_env.get('OPENAI_API_KEY', '').strip():
                agent_env['OPENAI_API_KEY'] = byok_codex
                cmd = [bin_, 'exec', '--json', '--skip-git-repo-check',
                       '--sandbox', 'workspace-write', '-C', str(cwd_path),
                       '-c', 'model_provider="openai"',
                       '-c', 'approval_policy="untrusted"']
                if model:
                    cmd += ['--model', model]
            else:
                # Codex → OpenAI directly (OPENAI_API_KEY); no base_url override.
                wire_model = model or 'gpt-4.1'
                cmd = [bin_, 'exec', '--json', '--skip-git-repo-check',
                       '--sandbox', 'workspace-write', '-C', str(cwd_path),
                       '-c', 'model_provider="openai"',
                       '-c', 'approval_policy="untrusted"']
                cmd += ['--model', wire_model]

        sys.stderr.write(
            f'[terminal] cli={cli} cwd={cwd_path} model={model or "(default)"}'
            f' auth={auth_method or "auto"} session={session_id[:8] or "-"}'
            f' project={project_id[:12] or "-"}'
            f' byok_claude={bool(byok_claude)} byok_codex={bool(byok_codex)}\n'
        )
        try:
            proc = subprocess.Popen(
                cmd,
                stdin=subprocess.PIPE,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True, bufsize=1, encoding='utf-8',
                cwd=str(cwd_path),
                env=agent_env,
            )
        except (FileNotFoundError, OSError) as e:
            return self._send_json(500, {'error': f'spawn {cli}: {e}'})

        # claude/codex read the prompt on stdin; gemini gets it in argv (--prompt)
        # so we just close its stdin to signal no interactive input.
        try:
            if cli != 'gemini':
                proc.stdin.write(prompt)
            proc.stdin.close()
        except Exception as e:
            try: proc.kill()
            except Exception: pass
            return self._send_json(500, {'error': f'stdin: {e}'})

        self.send_response(200)
        self.send_header('content-type', 'text/event-stream')
        self.send_header('cache-control', 'no-store')
        self.end_headers()

        def write_sse(obj):
            try:
                self.wfile.write(b'data: ' + json.dumps(obj).encode('utf-8') + b'\n\n')
                self.wfile.flush()
            except (BrokenPipeError, ConnectionResetError):
                return False
            return True

        def sse_delta(content, type_='text'):
            return write_sse({'choices': [{'index': 0, 'delta': {'content': content, 'type': type_}}]})

        stderr_buf = []
        def _drain_err():
            try:
                for line in proc.stderr:
                    stderr_buf.append(line)
            except Exception:
                pass
        threading.Thread(target=_drain_err, daemon=True).start()

        in_tokens = out_tokens = 0

        if cli == 'claude':
            try:
                for line in proc.stdout:
                    line = line.strip()
                    if not line: continue
                    try:
                        ev = json.loads(line)
                    except json.JSONDecodeError:
                        continue
                    t = ev.get('type')
                    if t == 'assistant':
                        msg = ev.get('message') or {}
                        for block in (msg.get('content') or []):
                            btype = block.get('type')
                            if btype == 'text':
                                text = block.get('text') or ''
                                if text:
                                    ok = sse_delta(text, 'text')
                                    if not ok: break
                            elif btype == 'tool_use':
                                name  = block.get('name') or 'tool'
                                inp   = block.get('input') or {}
                                inp_s = json.dumps(inp, ensure_ascii=False)
                                if len(inp_s) > 240: inp_s = inp_s[:240] + '…'
                                sse_delta(f'\n⚙ {name}: {inp_s}\n', 'tool')
                        u = msg.get('usage')
                        if u:
                            in_tokens  = u.get('input_tokens', in_tokens)
                            out_tokens = u.get('output_tokens', out_tokens)
                    elif t == 'result':
                        u = ev.get('usage') or {}
                        in_tokens  = u.get('input_tokens', in_tokens)
                        out_tokens = u.get('output_tokens', out_tokens)
                    elif t == 'error':
                        sse_delta(f'\n⚠ {ev.get("message") or ev}\n', 'error')
                if in_tokens or out_tokens:
                    write_sse({'choices': [], 'usage': {
                        'prompt_tokens': in_tokens,
                        'completion_tokens': out_tokens,
                        'total_tokens': in_tokens + out_tokens,
                    }})
            finally:
                try: proc.wait(timeout=2)
                except Exception:
                    try: proc.kill()
                    except Exception: pass
                if proc.returncode and proc.returncode != 0 and not (in_tokens or out_tokens):
                    err_text = ''.join(stderr_buf)[:600] or f'exit {proc.returncode}'
                    sse_delta(f'\n⚠ claude exited {proc.returncode}: {err_text}\n', 'error')

        elif cli == 'gemini':
            # Gemini CLI (`--prompt … --yolo`) prints its answer as plain text to
            # stdout (it may interleave a little ANSI for spinners). Stream it
            # straight through as 'text', stripping escape sequences so the chat
            # bubble stays clean. No token usage is reported by the CLI.
            _ansi_re = re.compile(r'\x1b\[[0-9;?]*[ -/]*[@-~]')
            text_emitted = False
            try:
                while True:
                    chunk = proc.stdout.read(256)
                    if not chunk:
                        break
                    clean = _ansi_re.sub('', chunk)
                    if clean:
                        text_emitted = True
                        if not sse_delta(clean, 'text'):
                            break
            finally:
                try: proc.wait(timeout=2)
                except Exception:
                    try: proc.kill()
                    except Exception: pass
                rc = proc.returncode
                stderr_text = ''.join(stderr_buf).strip()
                if rc and rc not in (0, None):
                    sse_delta(f'\n⚠ gemini exited {rc}: {stderr_text[:300] or "(no stderr)"}\n', 'error')
                elif not text_emitted:
                    hint = stderr_text[:300] or 'is the gemini CLI logged in? run it once in PTY mode to authenticate'
                    sse_delta(f'\n⚠ gemini returned no content — {hint}\n', 'error')
                try: self.wfile.write(b'data: [DONE]\n\n'); self.wfile.flush()
                except (BrokenPipeError, ConnectionResetError): pass

        else:  # codex
            def _content_text(value):
                if value is None: return ''
                if isinstance(value, str): return value
                if isinstance(value, list):
                    return ''.join(
                        (b if isinstance(b, str) else
                         b.get('text') or b.get('content') or b.get('output_text') or '')
                        for b in value
                    )
                if isinstance(value, dict):
                    return (value.get('text') or value.get('content')
                            or value.get('message') or '')
                return ''

            item_text_seen = {}
            events_seen = 0
            text_emitted = False
            try:
                for line in proc.stdout:
                    line = line.strip()
                    if not line: continue
                    try:
                        ev = json.loads(line)
                    except json.JSONDecodeError:
                        continue
                    events_seen += 1
                    t = ev.get('type', '')
                    if t in ('agent_message', 'message'):
                        text = _content_text(ev.get('content') or ev.get('message') or '')
                        if text:
                            text_emitted = True
                            sse_delta(text, 'text')
                    elif t in ('agent_message_delta', 'message_delta', 'response.output_text.delta'):
                        text = str(ev.get('delta') or ev.get('text') or ev.get('content') or '')
                        if text:
                            text_emitted = True
                            sse_delta(text, 'text')
                    elif t in ('item.updated', 'item.completed'):
                        item = ev.get('item') or {}
                        item_type = item.get('type', '')
                        if item_type == 'error':
                            msg = item.get('message') or item.get('text') or json.dumps(item)
                            if not ('disallowed by requirements' in msg
                                    or 'falling back to required value' in msg):
                                text_emitted = True
                                sse_delta(f'\n⚠ Codex: {msg}\n', 'error')
                        elif (item.get('role') == 'assistant'
                              or item_type in ('message', 'agent_message')):
                            text = (item.get('text')
                                    or _content_text(item.get('content') or item.get('message')))
                            item_id = item.get('id') or ev.get('item_id') or 'assistant'
                            prior = item_text_seen.get(item_id, '')
                            delta = (text[len(prior):]
                                     if text and text.startswith(prior)
                                     else (text if text != prior else ''))
                            item_text_seen[item_id] = text or prior
                            if delta:
                                text_emitted = True
                                sse_delta(delta, 'text')
                    elif t in ('turn.completed', 'response.completed'):
                        text = _content_text(
                            ev.get('last_agent_message') or ev.get('message') or '')
                        if text:
                            text_emitted = True
                            sse_delta(text, 'text')
                    elif t == 'function_call':
                        fname  = ev.get('name') or ev.get('function') or 'tool'
                        args   = ev.get('arguments') or ev.get('input') or {}
                        cmd_s  = ((args.get('cmd') or args.get('command') or json.dumps(args))
                                  if isinstance(args, dict) else str(args))
                        text_emitted = True
                        sse_delta(f'\n⚙ {fname}: {cmd_s}\n', 'tool')
                    elif t in ('error', 'turn.failed'):
                        msg = (ev.get('message')
                               or (ev.get('error') or {}).get('message')
                               or str(ev))
                        text_emitted = True
                        sse_delta(f'\n⚠ Codex error: {msg}\n', 'error')
            finally:
                try: proc.wait(timeout=2)
                except Exception:
                    try: proc.kill()
                    except Exception: pass
                rc = proc.returncode
                stderr_text = ''.join(stderr_buf).strip()
                if rc and rc not in (0, None):
                    sse_delta(f'\n⚠ Codex exited {rc}: {stderr_text[:300] or "(no stderr)"}\n', 'error')
                elif not text_emitted:
                    sse_delta(f'\n⚠ Codex returned no content (exit {rc})\n', 'error')
                try: self.wfile.write(b'data: [DONE]\n\n'); self.wfile.flush()
                except (BrokenPipeError, ConnectionResetError): pass

    # ---- Export endpoints (PowerPoint / Word / PDF) ----------------------
    def _vault_binary_path(self, rel: str, allowed_ext: tuple) -> pathlib.Path:
        """Resolve `rel` under the vault for binary files. Mirrors
        _vault_resolve but allows the caller's extension instead of forcing
        .md. Used by export/generate endpoints to drop files into the vault."""
        if not _vault_root:
            raise ValueError('vault directory not configured')
        root = pathlib.Path(_vault_root).resolve()
        if not root.is_dir():
            raise ValueError(f'vault directory does not exist: {root}')
        rel = (rel or '').lstrip('/').replace('\\', '/').strip()
        if not rel:
            raise ValueError('path required')
        ext = pathlib.Path(rel).suffix.lower()
        if ext not in allowed_ext:
            # If no extension was given, append the first allowed one.
            if not ext:
                rel = rel + allowed_ext[0]
                ext = allowed_ext[0]
            else:
                raise ValueError(f'extension must be one of {allowed_ext}, got {ext}')
        candidate = (root / rel).resolve()
        try: candidate.relative_to(root)
        except ValueError: raise ValueError('path escapes vault directory')
        candidate.parent.mkdir(parents=True, exist_ok=True)
        return candidate

    def _read_json_body(self):
        length = int(self.headers.get('content-length', 0) or 0)
        try:
            return json.loads(self.rfile.read(length).decode('utf-8') or '{}')
        except json.JSONDecodeError:
            return None

    def _export_pptx(self):
        """Render a markdown outline into a real .pptx and save to vault.
        Body: { path, content }. The outline uses `## Slide N: Title` headers
        and bullet lists; each `##` becomes a slide. Requires python-pptx."""
        body = self._read_json_body()
        if body is None: return self._send_json(400, {'error': 'bad json'})
        rel = (body.get('path') or '').strip()
        content = body.get('content') or ''
        if not rel or not content.strip():
            return self._send_json(400, {'error': 'path and content required'})
        try:
            out_path = self._vault_binary_path(rel, ('.pptx',))
        except ValueError as e:
            return self._send_json(400, {'error': str(e)})
        try:
            from pptx import Presentation  # noqa: PLC0415
            from pptx.util import Inches, Pt  # noqa: PLC0415
        except ImportError:
            return self._send_json(503, {'error': 'python-pptx not installed — run: pip install python-pptx'})

        prs = Presentation()
        # Title slide from first H1 if any.
        lines = content.splitlines()
        title_line = next((l for l in lines if l.strip().startswith('# ')), None)
        if title_line:
            layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title_line.lstrip('#').strip()
        # Each ## section becomes a slide.
        cur_title = None
        cur_bullets = []
        def flush():
            if cur_title is None and not cur_bullets: return
            layout = prs.slide_layouts[1]  # Title + Content
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = cur_title or 'Slide'
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for i, b in enumerate(cur_bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = b
                p.level = 0
        for ln in lines:
            s = ln.rstrip()
            if s.startswith('## '):
                flush()
                cur_title = s[3:].strip()
                # Strip leading "Slide N:" prefix the agents tend to write.
                cur_title = re.sub(r'^Slide\s*\d+\s*[:\-]\s*', '', cur_title, flags=re.IGNORECASE)
                cur_bullets = []
            elif re.match(r'^\s*[-*•]\s+', s):
                cur_bullets.append(re.sub(r'^\s*[-*•]\s+', '', s))
            elif s.strip() and cur_title is not None:
                # Plain paragraph under a slide → also a bullet.
                cur_bullets.append(s.strip())
        flush()

        try:
            prs.save(str(out_path))
        except Exception as e:
            return self._send_json(500, {'error': f'save failed: {e}'})
        rel_out = str(out_path.relative_to(pathlib.Path(_vault_root).resolve())).replace('\\', '/')
        return self._send_json(200, {'path': rel_out, 'slides': len(prs.slides)})

    def _export_docx(self):
        """Render markdown into a real .docx and save to vault. Body: {path, content}.
        Supports headings (# / ## / ###), bullets (- / *), and paragraphs.
        Requires python-docx."""
        body = self._read_json_body()
        if body is None: return self._send_json(400, {'error': 'bad json'})
        rel = (body.get('path') or '').strip()
        content = body.get('content') or ''
        if not rel or not content.strip():
            return self._send_json(400, {'error': 'path and content required'})
        try:
            out_path = self._vault_binary_path(rel, ('.docx',))
        except ValueError as e:
            return self._send_json(400, {'error': str(e)})
        try:
            from docx import Document  # noqa: PLC0415
        except ImportError:
            return self._send_json(503, {'error': 'python-docx not installed — run: pip install python-docx'})

        doc = Document()
        for ln in content.splitlines():
            s = ln.rstrip()
            if not s.strip():
                continue
            if s.startswith('### '):
                doc.add_heading(s[4:].strip(), level=3)
            elif s.startswith('## '):
                doc.add_heading(s[3:].strip(), level=2)
            elif s.startswith('# '):
                doc.add_heading(s[2:].strip(), level=1)
            elif re.match(r'^\s*[-*•]\s+', s):
                doc.add_paragraph(re.sub(r'^\s*[-*•]\s+', '', s), style='List Bullet')
            elif re.match(r'^\s*\d+\.\s+', s):
                doc.add_paragraph(re.sub(r'^\s*\d+\.\s+', '', s), style='List Number')
            else:
                doc.add_paragraph(s.strip())
        try:
            doc.save(str(out_path))
        except Exception as e:
            return self._send_json(500, {'error': f'save failed: {e}'})
        rel_out = str(out_path.relative_to(pathlib.Path(_vault_root).resolve())).replace('\\', '/')
        return self._send_json(200, {'path': rel_out})

    def _export_pdf(self):
        """Render markdown into a .pdf and save to vault. Body: {path, content}.
        Tries weasyprint (best) → reportlab (fallback). Returns 503 if neither."""
        body = self._read_json_body()
        if body is None: return self._send_json(400, {'error': 'bad json'})
        rel = (body.get('path') or '').strip()
        content = body.get('content') or ''
        if not rel or not content.strip():
            return self._send_json(400, {'error': 'path and content required'})
        try:
            out_path = self._vault_binary_path(rel, ('.pdf',))
        except ValueError as e:
            return self._send_json(400, {'error': str(e)})

        # Path A — weasyprint (real CSS, good typography).
        try:
            import markdown as _md  # noqa: PLC0415
            from weasyprint import HTML as _WeasyHTML  # noqa: PLC0415
            html_body = _md.markdown(content, extensions=['tables', 'fenced_code'])
            full_html = f'<html><head><meta charset="utf-8"><style>body{{font-family:Helvetica,Arial,sans-serif;line-height:1.5;padding:48px;}} h1,h2,h3{{color:#222}} code{{background:#f4f4f4;padding:2px 6px;border-radius:3px}} pre{{background:#f4f4f4;padding:12px;border-radius:6px;overflow-x:auto}}</style></head><body>{html_body}</body></html>'
            _WeasyHTML(string=full_html).write_pdf(str(out_path))
            rel_out = str(out_path.relative_to(pathlib.Path(_vault_root).resolve())).replace('\\', '/')
            return self._send_json(200, {'path': rel_out, 'renderer': 'weasyprint'})
        except ImportError:
            pass
        except Exception as e:
            sys.stderr.write(f'[export pdf] weasyprint failed: {e}\n')

        # Path B — reportlab (simpler, no CSS, ships with most Pythons via pip).
        try:
            from reportlab.lib.pagesizes import letter  # noqa: PLC0415
            from reportlab.lib.styles import getSampleStyleSheet  # noqa: PLC0415
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer  # noqa: PLC0415
        except ImportError:
            return self._send_json(503, {'error': 'install either weasyprint+markdown OR reportlab — run: pip install weasyprint markdown  (or)  pip install reportlab'})

        doc = SimpleDocTemplate(str(out_path), pagesize=letter, topMargin=48, bottomMargin=48, leftMargin=48, rightMargin=48)
        styles = getSampleStyleSheet()
        story = []
        for ln in content.splitlines():
            s = ln.rstrip()
            if not s.strip():
                story.append(Spacer(1, 6))
                continue
            if s.startswith('# '):
                story.append(Paragraph(s[2:].strip(), styles['Title']))
            elif s.startswith('## '):
                story.append(Paragraph(s[3:].strip(), styles['Heading2']))
            elif s.startswith('### '):
                story.append(Paragraph(s[4:].strip(), styles['Heading3']))
            elif re.match(r'^\s*[-*•]\s+', s):
                story.append(Paragraph('• ' + re.sub(r'^\s*[-*•]\s+', '', s), styles['BodyText']))
            else:
                story.append(Paragraph(s.strip(), styles['BodyText']))
        try:
            doc.build(story)
        except Exception as e:
            return self._send_json(500, {'error': f'pdf build: {e}'})
        rel_out = str(out_path.relative_to(pathlib.Path(_vault_root).resolve())).replace('\\', '/')
        return self._send_json(200, {'path': rel_out, 'renderer': 'reportlab'})

    # ---- Image / video generation (user-configured provider) -------------
    def _generate_image(self):
        """Generate an image via the user's configured provider, save to vault.
        Body: { path, prompt, provider, model, size?, apiKey? }
        - provider: 'openai' | 'google' | 'fal'
        - apiKey: forwarded from the browser (decrypted client-side); server
          env vars override if set."""
        body = self._read_json_body()
        if body is None: return self._send_json(400, {'error': 'bad json'})
        rel = (body.get('path') or '').strip()
        prompt = (body.get('prompt') or '').strip()
        provider = (body.get('provider') or '').strip().lower()
        model = (body.get('model') or '').strip()
        size = (body.get('size') or '1024x1024').strip()
        if not rel or not prompt or not provider:
            return self._send_json(400, {'error': 'path, prompt, and provider required'})
        try:
            out_path = self._vault_binary_path(rel, ('.png', '.jpg', '.jpeg', '.webp'))
        except ValueError as e:
            return self._send_json(400, {'error': str(e)})

        api_key = body.get('apiKey') or ''
        if provider == 'openai':
            api_key = os.environ.get('OPENAI_API_KEY') or api_key
            if not api_key: return self._send_json(400, {'error': 'OPENAI_API_KEY required'})
            payload = {'model': model or 'dall-e-3', 'prompt': prompt, 'size': size, 'n': 1, 'response_format': 'b64_json'}
            try:
                req = urllib.request.Request('https://api.openai.com/v1/images/generations',
                    data=json.dumps(payload).encode('utf-8'),
                    headers={'authorization': f'Bearer {api_key}', 'content-type': 'application/json'})
                with urllib.request.urlopen(req, timeout=120) as r:
                    j = json.loads(r.read().decode('utf-8'))
                b64 = j.get('data', [{}])[0].get('b64_json', '')
                if not b64: return self._send_json(502, {'error': 'no image returned'})
                out_path.write_bytes(base64.b64decode(b64))
            except urllib.error.HTTPError as e:
                return self._send_json(e.code, {'error': e.read().decode('utf-8', 'replace')[:600]})
            except Exception as e:
                return self._send_json(500, {'error': str(e)})
        elif provider == 'google':
            api_key = os.environ.get('GOOGLE_API_KEY') or api_key
            if not api_key: return self._send_json(400, {'error': 'GOOGLE_API_KEY required'})
            model_id = model or 'gemini-2.5-flash-image-preview'
            # Two flavors of Google image gen:
            #  - Imagen (imagen-3.0-*, imagen-4-*) → :predict endpoint, predictions[].bytesBase64Encoded
            #  - Gemini Flash Image / "nano banana" (gemini-*-image-*) → :generateContent,
            #    candidates[0].content.parts[*].inlineData.data
            # We pick the right endpoint based on the model name so the same
            # GOOGLE_API_KEY works for both — no extra setting required.
            is_gemini = 'gemini' in model_id.lower()
            try:
                if is_gemini:
                    url = f'https://generativelanguage.googleapis.com/v1beta/models/{model_id}:generateContent?key={api_key}'
                    payload = {
                        'contents': [{'parts': [{'text': prompt}]}],
                        'generationConfig': {'responseModalities': ['IMAGE']},
                    }
                    req = urllib.request.Request(url, data=json.dumps(payload).encode('utf-8'),
                        headers={'content-type': 'application/json'})
                    with urllib.request.urlopen(req, timeout=180) as r:
                        j = json.loads(r.read().decode('utf-8'))
                    b64 = ''
                    mime = ''
                    cands = j.get('candidates') or []
                    if cands:
                        for part in (cands[0].get('content', {}).get('parts') or []):
                            inline = part.get('inlineData') or part.get('inline_data') or {}
                            if inline.get('data'):
                                b64 = inline['data']
                                mime = (inline.get('mimeType') or inline.get('mime_type') or '').lower()
                                break
                    if not b64:
                        # Surface text fallback if model returned a refusal instead of image bytes
                        text_out = ''
                        for c in cands:
                            for p in (c.get('content', {}).get('parts') or []):
                                if p.get('text'): text_out += p['text']
                        return self._send_json(502, {'error': f'no image bytes in Gemini response{": " + text_out[:200] if text_out else ""}'})
                    out_path.write_bytes(base64.b64decode(b64))
                else:
                    url = f'https://generativelanguage.googleapis.com/v1beta/models/{model_id}:predict?key={api_key}'
                    payload = {'instances': [{'prompt': prompt}], 'parameters': {'sampleCount': 1}}
                    req = urllib.request.Request(url, data=json.dumps(payload).encode('utf-8'),
                        headers={'content-type': 'application/json'})
                    with urllib.request.urlopen(req, timeout=120) as r:
                        j = json.loads(r.read().decode('utf-8'))
                    preds = j.get('predictions', [])
                    if not preds: return self._send_json(502, {'error': 'no image returned'})
                    b64 = preds[0].get('bytesBase64Encoded') or preds[0].get('image', {}).get('bytesBase64Encoded') or ''
                    if not b64: return self._send_json(502, {'error': 'no image bytes in response'})
                    out_path.write_bytes(base64.b64decode(b64))
            except urllib.error.HTTPError as e:
                return self._send_json(e.code, {'error': e.read().decode('utf-8', 'replace')[:600]})
            except Exception as e:
                return self._send_json(500, {'error': str(e)})
        elif provider == 'fal':
            api_key = os.environ.get('FAL_KEY') or api_key
            if not api_key: return self._send_json(400, {'error': 'FAL_KEY required'})
            model_id = model or 'fal-ai/flux/schnell'
            try:
                req = urllib.request.Request(f'https://fal.run/{model_id}',
                    data=json.dumps({'prompt': prompt}).encode('utf-8'),
                    headers={'authorization': f'Key {api_key}', 'content-type': 'application/json'})
                with urllib.request.urlopen(req, timeout=180) as r:
                    j = json.loads(r.read().decode('utf-8'))
                img_url = ((j.get('images') or [{}])[0]).get('url')
                if not img_url: return self._send_json(502, {'error': 'no image url in fal response'})
                with urllib.request.urlopen(img_url, timeout=60) as img:
                    out_path.write_bytes(img.read())
            except urllib.error.HTTPError as e:
                return self._send_json(e.code, {'error': e.read().decode('utf-8', 'replace')[:600]})
            except Exception as e:
                return self._send_json(500, {'error': str(e)})
        elif provider == 'a1111':
            # Automatic1111 Stable Diffusion WebUI — REST API at /sdapi/v1/txt2img.
            # Run A1111 with --api (or --api --listen for LAN). No API key needed.
            base = (body.get('baseUrl') or 'http://127.0.0.1:7860').rstrip('/')
            try:
                w, h = (int(x) for x in size.split('x')[:2]) if 'x' in size else (1024, 1024)
            except Exception:
                w, h = 1024, 1024
            payload = {
                'prompt': prompt,
                'width': w, 'height': h,
                'steps': int(body.get('steps') or 25),
                'cfg_scale': float(body.get('cfgScale') or 7.5),
                'sampler_name': body.get('sampler') or 'DPM++ 2M Karras',
            }
            if model: payload['override_settings'] = {'sd_model_checkpoint': model}
            try:
                req = urllib.request.Request(f'{base}/sdapi/v1/txt2img',
                    data=json.dumps(payload).encode('utf-8'),
                    headers={'content-type': 'application/json'})
                with urllib.request.urlopen(req, timeout=300) as r:
                    j = json.loads(r.read().decode('utf-8'))
                imgs = j.get('images') or []
                if not imgs: return self._send_json(502, {'error': 'no images in A1111 response'})
                out_path.write_bytes(base64.b64decode(imgs[0]))
            except urllib.error.URLError as e:
                return self._send_json(502, {'error': f'A1111 not reachable at {base} — start it with --api flag: {e}'})
            except Exception as e:
                return self._send_json(500, {'error': f'A1111: {e}'})
        elif provider == 'comfyui':
            # ComfyUI — accepts a `workflow` JSON OR a simple `prompt` that we
            # wrap into a minimal txt2img graph. Polls /history for completion
            # then downloads the image bytes from /view.
            base = (body.get('baseUrl') or 'http://127.0.0.1:8188').rstrip('/')
            workflow = body.get('workflow')
            if not workflow:
                # Default minimal txt2img workflow — assumes a SD1.5/SDXL ckpt is
                # loaded under the name in `model` (or "model.safetensors").
                ckpt = model or 'model.safetensors'
                try:
                    w, h = (int(x) for x in size.split('x')[:2]) if 'x' in size else (1024, 1024)
                except Exception:
                    w, h = 1024, 1024
                workflow = {
                    "3": {"class_type": "KSampler", "inputs": {"seed": secrets.randbits(31),
                          "steps": int(body.get('steps') or 25), "cfg": float(body.get('cfgScale') or 7.5),
                          "sampler_name": "euler", "scheduler": "normal", "denoise": 1.0,
                          "model": ["4", 0], "positive": ["6", 0], "negative": ["7", 0], "latent_image": ["5", 0]}},
                    "4": {"class_type": "CheckpointLoaderSimple", "inputs": {"ckpt_name": ckpt}},
                    "5": {"class_type": "EmptyLatentImage", "inputs": {"width": w, "height": h, "batch_size": 1}},
                    "6": {"class_type": "CLIPTextEncode", "inputs": {"text": prompt, "clip": ["4", 1]}},
                    "7": {"class_type": "CLIPTextEncode", "inputs": {"text": body.get('negative') or '', "clip": ["4", 1]}},
                    "8": {"class_type": "VAEDecode", "inputs": {"samples": ["3", 0], "vae": ["4", 2]}},
                    "9": {"class_type": "SaveImage", "inputs": {"filename_prefix": "cafresohq", "images": ["8", 0]}},
                }
            client_id = uuid.uuid4().hex
            try:
                # Queue the prompt.
                req = urllib.request.Request(f'{base}/prompt',
                    data=json.dumps({'prompt': workflow, 'client_id': client_id}).encode('utf-8'),
                    headers={'content-type': 'application/json'})
                with urllib.request.urlopen(req, timeout=30) as r:
                    j = json.loads(r.read().decode('utf-8'))
                prompt_id = j.get('prompt_id')
                if not prompt_id: return self._send_json(502, {'error': 'no prompt_id in ComfyUI response'})
                # Poll history until done (max ~5 min).
                hist = None
                for _i in range(150):
                    time.sleep(2)
                    try:
                        with urllib.request.urlopen(f'{base}/history/{prompt_id}', timeout=10) as r:
                            h = json.loads(r.read().decode('utf-8'))
                        if h.get(prompt_id):
                            hist = h[prompt_id]
                            break
                    except Exception:
                        continue
                if not hist: return self._send_json(504, {'error': 'ComfyUI: prompt did not complete within 5 minutes'})
                # Find first image output.
                outputs = hist.get('outputs') or {}
                img_meta = None
                for _node_id, node_out in outputs.items():
                    if node_out.get('images'):
                        img_meta = node_out['images'][0]
                        break
                if not img_meta: return self._send_json(502, {'error': 'ComfyUI: no image in outputs'})
                # Download the image bytes.
                qs = urllib.parse.urlencode({k: v for k, v in img_meta.items() if k in ('filename','subfolder','type')})
                with urllib.request.urlopen(f'{base}/view?{qs}', timeout=60) as img:
                    out_path.write_bytes(img.read())
            except urllib.error.URLError as e:
                return self._send_json(502, {'error': f'ComfyUI not reachable at {base}: {e}'})
            except Exception as e:
                return self._send_json(500, {'error': f'ComfyUI: {e}'})
        else:
            return self._send_json(400, {'error': f'unsupported provider: {provider}'})
        rel_out = str(out_path.relative_to(pathlib.Path(_vault_root).resolve())).replace('\\', '/')
        return self._send_json(200, {'path': rel_out, 'provider': provider, 'model': model})

    def _generate_video(self):
        """Generate a video via the user's configured provider, save to vault.
        Body: { path, prompt, provider, model, duration?, apiKey? }
        Most video APIs are long-running — this endpoint kicks off the job,
        polls for completion, then writes the bytes. Times out after 10min."""
        body = self._read_json_body()
        if body is None: return self._send_json(400, {'error': 'bad json'})
        rel = (body.get('path') or '').strip()
        prompt = (body.get('prompt') or '').strip()
        provider = (body.get('provider') or '').strip().lower()
        model = (body.get('model') or '').strip()
        duration = int(body.get('duration') or 5)
        if not rel or not prompt or not provider:
            return self._send_json(400, {'error': 'path, prompt, and provider required'})
        try:
            out_path = self._vault_binary_path(rel, ('.mp4', '.mov', '.webm'))
        except ValueError as e:
            return self._send_json(400, {'error': str(e)})

        api_key = body.get('apiKey') or ''
        if provider == 'fal':
            api_key = os.environ.get('FAL_KEY') or api_key
            if not api_key: return self._send_json(400, {'error': 'FAL_KEY required'})
            model_id = model or 'fal-ai/bytedance/seedance/v1/lite/text-to-video'
            try:
                req = urllib.request.Request(f'https://fal.run/{model_id}',
                    data=json.dumps({'prompt': prompt, 'duration': duration}).encode('utf-8'),
                    headers={'authorization': f'Key {api_key}', 'content-type': 'application/json'})
                with urllib.request.urlopen(req, timeout=600) as r:
                    j = json.loads(r.read().decode('utf-8'))
                video_url = (j.get('video') or {}).get('url') or j.get('url') or ''
                if not video_url: return self._send_json(502, {'error': 'no video url in fal response'})
                with urllib.request.urlopen(video_url, timeout=300) as vid:
                    out_path.write_bytes(vid.read())
            except urllib.error.HTTPError as e:
                return self._send_json(e.code, {'error': e.read().decode('utf-8', 'replace')[:600]})
            except Exception as e:
                return self._send_json(500, {'error': str(e)})
        elif provider == 'openai':
            # Sora API is gated; we provide the call shape but most accounts
            # will get a 403. The error message guides the user.
            api_key = os.environ.get('OPENAI_API_KEY') or api_key
            if not api_key: return self._send_json(400, {'error': 'OPENAI_API_KEY required'})
            return self._send_json(501, {'error': 'OpenAI Sora video generation not yet wired (API still gated). Try provider=fal with a Seedance/Veo model instead.'})
        elif provider == 'google':
            return self._send_json(501, {'error': 'Google Veo direct API not yet wired. Try provider=fal with a Veo model instead.'})
        elif provider == 'comfyui':
            # ComfyUI for video — caller passes a `workflow` JSON describing
            # an AnimateDiff / SVD / Hunyuan / Mochi graph. We don't synthesize
            # a default one because video workflows are model-specific.
            base = (body.get('baseUrl') or 'http://127.0.0.1:8188').rstrip('/')
            workflow = body.get('workflow')
            if not workflow:
                return self._send_json(400, {
                    'error': 'ComfyUI video requires a `workflow` JSON (export from your AnimateDiff/SVD/Mochi/Hunyuan setup). Auto-default not provided because video workflows are model-specific. See docs.'
                })
            client_id = uuid.uuid4().hex
            try:
                req = urllib.request.Request(f'{base}/prompt',
                    data=json.dumps({'prompt': workflow, 'client_id': client_id}).encode('utf-8'),
                    headers={'content-type': 'application/json'})
                with urllib.request.urlopen(req, timeout=30) as r:
                    j = json.loads(r.read().decode('utf-8'))
                prompt_id = j.get('prompt_id')
                if not prompt_id: return self._send_json(502, {'error': 'no prompt_id in ComfyUI response'})
                hist = None
                for _i in range(300):  # video takes longer — up to ~10 min
                    time.sleep(2)
                    try:
                        with urllib.request.urlopen(f'{base}/history/{prompt_id}', timeout=10) as r:
                            h = json.loads(r.read().decode('utf-8'))
                        if h.get(prompt_id):
                            hist = h[prompt_id]
                            break
                    except Exception:
                        continue
                if not hist: return self._send_json(504, {'error': 'ComfyUI: prompt did not complete within 10 minutes'})
                outputs = hist.get('outputs') or {}
                # Look for video output — VHS_VideoCombine, SaveVideo, etc.
                vid_meta = None
                for _node_id, node_out in outputs.items():
                    for key in ('videos', 'gifs', 'images'):  # some nodes call the file 'images' even for mp4
                        if node_out.get(key):
                            vid_meta = node_out[key][0]
                            break
                    if vid_meta: break
                if not vid_meta: return self._send_json(502, {'error': 'ComfyUI: no video in outputs — does your workflow include VHS_VideoCombine or SaveVideo?'})
                qs = urllib.parse.urlencode({k: v for k, v in vid_meta.items() if k in ('filename','subfolder','type')})
                with urllib.request.urlopen(f'{base}/view?{qs}', timeout=300) as vid:
                    out_path.write_bytes(vid.read())
            except urllib.error.URLError as e:
                return self._send_json(502, {'error': f'ComfyUI not reachable at {base}: {e}'})
            except Exception as e:
                return self._send_json(500, {'error': f'ComfyUI: {e}'})
        else:
            return self._send_json(400, {'error': f'unsupported provider: {provider}'})
        rel_out = str(out_path.relative_to(pathlib.Path(_vault_root).resolve())).replace('\\', '/')
        return self._send_json(200, {'path': rel_out, 'provider': provider, 'model': model})

    # ---- External approvals (Claude Code hook bridge) --------------------
    def _approval_submit(self):
        """Hook script POSTs a tool-use request here. Returns {id} immediately
        so the script can long-poll /approvals/external/wait until the human
        decides. Body: {tool, input, cwd, agent, sessionId, summary?}."""
        _gc_approvals()
        length = int(self.headers.get('content-length', 0) or 0)
        try:
            body = json.loads(self.rfile.read(length).decode('utf-8') or '{}')
        except json.JSONDecodeError:
            return self._send_json(400, {'error': 'bad json'})
        entry = _new_approval(body)
        sys.stderr.write(f'[approvals] queued {entry["id"]} tool={entry["tool"]} '
                         f'agent={entry["agent"]} cwd={entry["cwd"]}\n')
        return self._send_json(200, {'id': entry['id']})

    def _approval_wait(self):
        """Hook script long-polls here. Returns when a decision is in, or after
        timeout (script should re-poll). Query: ?id=...&timeout=25."""
        qs = urllib.parse.parse_qs(urllib.parse.urlparse(self.path).query)
        aid = qs.get('id', [''])[0]
        try:
            wait_s = max(1, min(55, int(qs.get('timeout', ['25'])[0])))
        except ValueError:
            wait_s = 25
        with _approvals_lock:
            entry = _approvals_pending.get(aid)
        if not entry:
            return self._send_json(404, {'error': 'unknown approval id'})
        # Block this thread until decision or timeout. ThreadingMixIn means
        # other requests keep flowing while we wait.
        decided = entry['event'].wait(timeout=wait_s)
        if not decided:
            return self._send_json(200, {'pending': True, 'id': aid})
        # Decided — pop the entry so we don't leak.
        with _approvals_lock:
            _approvals_pending.pop(aid, None)
        return self._send_json(200, {
            'pending': False,
            'id': aid,
            'decision': entry['decision'] or 'deny',
            'reason': entry['reason'] or '',
        })

    def _approval_list(self):
        """UI polls this for pending external approvals to render in the tray."""
        _gc_approvals()
        with _approvals_lock:
            items = [{
                'id': e['id'], 'ts': e['ts'], 'tool': e['tool'],
                'input': e['input'], 'cwd': e['cwd'], 'agent': e['agent'],
                'sessionId': e['sessionId'], 'summary': e['summary'],
            } for e in _approvals_pending.values() if e['decision'] is None]
        items.sort(key=lambda x: x['ts'])
        return self._send_json(200, {'pending': items})

    def _approval_decide(self):
        """UI POSTs the human decision: {id, decision: 'allow'|'deny', reason?}.
        We flip the entry's decision and signal the waiting hook thread."""
        length = int(self.headers.get('content-length', 0) or 0)
        try:
            body = json.loads(self.rfile.read(length).decode('utf-8') or '{}')
        except json.JSONDecodeError:
            return self._send_json(400, {'error': 'bad json'})
        aid = (body.get('id') or '').strip()
        decision = (body.get('decision') or '').strip().lower()
        if decision not in ('allow', 'deny'):
            return self._send_json(400, {'error': 'decision must be allow|deny'})
        reason = (body.get('reason') or '').strip()[:400]
        with _approvals_lock:
            entry = _approvals_pending.get(aid)
            if not entry:
                return self._send_json(404, {'error': 'unknown approval id'})
            if entry['decision'] is not None:
                return self._send_json(409, {'error': 'already decided'})
            entry['decision'] = decision
            entry['reason'] = reason
        entry['event'].set()
        sys.stderr.write(f'[approvals] decided {aid} -> {decision}\n')
        return self._send_json(200, {'id': aid, 'decision': decision})

    def _hq_handler(self, method):
        """Read/write HQ state and memory files.
        GET  /hq/state/<name>    → read hq-state/<name>.json
        PUT  /hq/state/<name>    → write hq-state/<name>.json
        GET  /hq/memory/<name>   → read memory/<name>.json
        PUT  /hq/memory/<name>   → write memory/<name>.json
        """
        parts = self.path.split('/')  # ['', 'hq', 'state'|'memory', '<name>']
        if len(parts) < 4 or not parts[3]:
            return self._send_json(400, {'error': 'path must be /hq/state/<name> or /hq/memory/<name>'})

        scope, name = parts[2], parts[3]
        # Sanitize name — alphanumeric, hyphens, underscores only
        if not _re.match(r'^[\w\-]+$', name):
            return self._send_json(400, {'error': f'invalid name: {name!r}'})

        if scope == 'state':
            base = _hq_state_dir
        elif scope == 'memory':
            base = _hq_memory_dir
        else:
            return self._send_json(404, {'error': f'unknown scope: {scope}'})

        try:
            base.mkdir(parents=True, exist_ok=True)
        except Exception as e:
            return self._send_json(503, {'error': f'cannot create dir: {e}'})

        filepath = base / f'{name}.json'

        if method == 'GET':
            if not filepath.exists():
                return self._send_json(404, {'error': 'not found', 'key': name})
            try:
                data = filepath.read_text(encoding='utf-8')
                self.send_response(200)
                self.send_header('Content-Type', 'application/json')
                self.send_header('Access-Control-Allow-Origin', '*')
                self.end_headers()
                self.wfile.write(data.encode('utf-8'))
            except Exception as e:
                return self._send_json(500, {'error': str(e)})

        elif method == 'PUT':
            length = int(self.headers.get('content-length', 0) or 0)
            try:
                body = self.rfile.read(length)
                # Validate it's JSON
                json.loads(body.decode('utf-8'))
                filepath.write_bytes(body)
                # For agent roster, also write a human-readable markdown summary
                if scope == 'memory' and name == 'agents':
                    self._write_agents_md(base, json.loads(body.decode('utf-8')))
                return self._send_json(200, {'ok': True, 'path': str(filepath)})
            except json.JSONDecodeError:
                return self._send_json(400, {'error': 'body must be valid JSON'})
            except Exception as e:
                return self._send_json(500, {'error': str(e)})

    def _write_agents_md(self, base, agents):
        """Write a human-readable agent roster markdown file alongside the JSON."""
        try:
            from datetime import datetime, timezone
            lines = [
                '# CafresoHQ Agent Roster',
                f'_Updated: {datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")}_',
                '',
            ]
            for a in (agents or []):
                lines.append(f"## {a.get('name', '?')} — {a.get('role', '')}")
                lines.append(f"- Model: {a.get('model', 'default')}")
                tools = a.get('tools') or []
                if tools:
                    lines.append(f"- Tools: {', '.join(tools)}")
                lines.append(f"- Elevated: {'yes' if a.get('elevated') else 'no'}")
                lines.append('')
            (base / 'hq-agents.md').write_text('\n'.join(lines), encoding='utf-8')
        except Exception:
            pass  # markdown summary is best-effort

    def _send_json(self, code, payload):
        body = json.dumps(payload).encode('utf-8')
        self.send_response(code)
        self.send_header('content-type', 'application/json; charset=utf-8')
        self.send_header('content-length', str(len(body)))
        self.end_headers()
        try: self.wfile.write(body)
        except (BrokenPipeError, ConnectionResetError): pass

    def _vault(self, method):
        global _vault_root, _vault_backend, _vault_rest_url, _vault_rest_key
        url = urllib.parse.urlparse(self.path)
        path = url.path
        qs = urllib.parse.parse_qs(url.query)

        # ---------- Status ----------
        if path == '/vault/status' and method == 'GET':
            if _vault_backend == 'fs' and _vault_root:
                try:
                    pathlib.Path(_vault_root).mkdir(parents=True, exist_ok=True)
                except OSError:
                    pass
            fs_ok = bool(_vault_root) and pathlib.Path(_vault_root).is_dir()
            rest_ok = False
            rest_detail = ''
            if _vault_rest_url and _vault_rest_key:
                try:
                    s, _h, _b = _obsidian_request('GET', '/')
                    rest_ok = (s == 200)
                    rest_detail = '' if rest_ok else f'http {s}'
                except Exception as e:
                    rest_detail = str(e)[:120]
            configured = (_vault_backend == 'rest' and rest_ok) or (_vault_backend == 'fs' and fs_ok)
            return self._send_json(200, {
                'configured': configured,
                'backend': _vault_backend,
                'root': _vault_root,
                'defaultRoot': str(_default_vault_root),
                'restUrl': _vault_rest_url,
                'restKey': '••••' if _vault_rest_key else '',
                'fsExists': fs_ok,
                'restReachable': rest_ok,
                'restDetail': rest_detail,
                # Legacy field for older clients.
                'exists': configured,
            })

        # ---------- Discover local vaults ----------
        if path == '/vault/discover' and method == 'GET':
            return self._send_json(200, {
                'defaultRoot': str(_default_vault_root),
                'vaults': _discover_obsidian_vaults(),
            })

        # ---------- Opaque-blob vault (zero-knowledge / vetKeys) ----------
        # Stores client-encrypted ciphertext blobs identified by hex IDs.
        # The container has zero visibility into contents — it's an
        # untrusted blob store from the client's perspective. Used by the
        # SvelteKit /vault route with vetKeys-derived AES-GCM encryption.
        #
        #   PUT    /vault/blob/<id>   body=<base64 ciphertext>
        #   GET    /vault/blob/<id>   → base64 ciphertext (text/plain)
        #   DELETE /vault/blob/<id>
        #
        # IDs are accepted if they match [a-f0-9]{16,64} OR the literal
        # 'index' alias. No directory hierarchy — flat blob namespace.
        if path.startswith('/vault/blob/'):
            blob_id = path[len('/vault/blob/'):]
            # Strict ID validation: hex chars only (max 64 = SHA-256 length).
            # Rejects any path traversal, slashes, or shell metacharacters.
            if not re.match(r'^[A-Za-z0-9_-]{1,64}$', blob_id):
                return self._send_json(400, {'error': 'invalid blob id'})

            # OCI backend: blobs live under <user-prefix>/blobs/<id>.bin
            if _vault_backend == 'oci':
                if not (_oci_vault_namespace and _oci_vault_bucket):
                    return self._send_json(503, {'error': 'OCI vault not configured'})
                cli = _oci_object_client()
                key = _oci_obj_key('blobs/' + blob_id + '.bin')
                if method == 'PUT':
                    length = int(self.headers.get('content-length', 0) or 0)
                    if length == 0:
                        return self._send_json(400, {'error': 'empty blob'})
                    if length > 300 * 1024 * 1024:  # 300 MB cap per blob
                        return self._send_json(413, {'error': 'blob too large (max 300 MB)'})
                    body = self.rfile.read(length)
                    try:
                        cli.put_object(_oci_vault_namespace, _oci_vault_bucket, key, body,
                                       content_type='application/octet-stream')
                    except Exception as e:
                        return self._send_json(502, {'error': f'oci put: {e}'})
                    return self._send_json(200, {'id': blob_id, 'size': len(body)})
                if method == 'GET':
                    try:
                        resp = cli.get_object(_oci_vault_namespace, _oci_vault_bucket, key)
                        body = resp.data.content
                    except Exception as e:
                        err = str(e)
                        if '404' in err or 'NoSuchKey' in err or 'ObjectNotFound' in err:
                            return self._send_json(404, {'error': 'not found'})
                        return self._send_json(502, {'error': f'oci get: {e}'})
                    self.send_response(200)
                    self.send_header('content-type', 'application/octet-stream')
                    self.send_header('content-length', str(len(body)))
                    self.end_headers()
                    try: self.wfile.write(body)
                    except (BrokenPipeError, ConnectionResetError): pass
                    return
                if method == 'DELETE':
                    try:
                        cli.delete_object(_oci_vault_namespace, _oci_vault_bucket, key)
                    except Exception as e:
                        err = str(e)
                        if '404' in err or 'NoSuchKey' in err or 'ObjectNotFound' in err:
                            return self._send_json(404, {'error': 'not found'})
                        return self._send_json(502, {'error': f'oci delete: {e}'})
                    return self._send_json(200, {'id': blob_id, 'deleted': True})
                return self._send_json(405, {'error': 'method not allowed'})

            # FS backend: blobs live under <vault_root>/.blobs/<id>.bin
            if _vault_root:
                root = pathlib.Path(_vault_root).resolve()
                blob_dir = root / '.blobs'
                try: blob_dir.mkdir(parents=True, exist_ok=True)
                except OSError as e:
                    return self._send_json(500, {'error': f'cannot create blob dir: {e}'})
                blob_path = blob_dir / (blob_id + '.bin')
                # Defense in depth: confirm the resolved path is still under the dir.
                try:
                    blob_path.resolve().relative_to(blob_dir.resolve())
                except ValueError:
                    return self._send_json(400, {'error': 'invalid blob path'})

                if method == 'PUT':
                    length = int(self.headers.get('content-length', 0) or 0)
                    if length == 0:
                        return self._send_json(400, {'error': 'empty blob'})
                    if length > 25 * 1024 * 1024:
                        return self._send_json(413, {'error': 'blob too large (max 25 MB)'})
                    body = self.rfile.read(length)
                    try:
                        blob_path.write_bytes(body)
                    except OSError as e:
                        return self._send_json(500, {'error': str(e)})
                    return self._send_json(200, {'id': blob_id, 'size': len(body)})
                if method == 'GET':
                    if not blob_path.exists():
                        return self._send_json(404, {'error': 'not found'})
                    try:
                        body = blob_path.read_bytes()
                    except OSError as e:
                        return self._send_json(500, {'error': str(e)})
                    self.send_response(200)
                    self.send_header('content-type', 'application/octet-stream')
                    self.send_header('content-length', str(len(body)))
                    self.end_headers()
                    try: self.wfile.write(body)
                    except (BrokenPipeError, ConnectionResetError): pass
                    return
                if method == 'DELETE':
                    if not blob_path.exists():
                        return self._send_json(404, {'error': 'not found'})
                    try: blob_path.unlink()
                    except OSError as e:
                        return self._send_json(500, {'error': str(e)})
                    return self._send_json(200, {'id': blob_id, 'deleted': True})
                return self._send_json(405, {'error': 'method not allowed'})

            return self._send_json(503, {'error': 'no vault backend configured'})

        # ---------- Configure ----------
        if path == '/vault/configure' and method == 'POST':
            length = int(self.headers.get('content-length', 0) or 0)
            try:
                body = json.loads(self.rfile.read(length).decode('utf-8') or '{}')
            except json.JSONDecodeError:
                return self._send_json(400, {'error': 'bad json'})
            # Update only what was sent so partial PATCH-style works.
            if 'root' in body:
                root = _clean_vault_root(body.get('root'))
                if root:
                    try:
                        root = str(pathlib.Path(root).resolve())
                    except OSError as e:
                        return self._send_json(400, {'error': f'invalid directory path: {e}'})
                if root and pathlib.Path(root) == _default_vault_root.resolve():
                    try:
                        pathlib.Path(root).mkdir(parents=True, exist_ok=True)
                    except OSError as e:
                        return self._send_json(500, {'error': f'could not create default vault: {e}'})
                if root and not pathlib.Path(root).is_dir():
                    return self._send_json(400, {'error': f'not a directory: {root}'})
                _vault_root = root
            if 'backend' in body:
                bk = (body.get('backend') or 'fs').strip()
                if bk not in ('fs', 'rest'):
                    return self._send_json(400, {'error': f'bad backend: {bk}'})
                _vault_backend = bk
            if 'restUrl' in body:
                _vault_rest_url = (body.get('restUrl') or '').strip()
            if 'restKey' in body and body['restKey']:  # ignore empty string (preserves existing)
                _vault_rest_key = body['restKey'].strip()
            return self._send_json(200, {
                'backend': _vault_backend,
                'root': _vault_root,
                'defaultRoot': str(_default_vault_root),
                'restUrl': _vault_rest_url,
                'restKey': '••••' if _vault_rest_key else '',
            })

        # Routes from here require a configured backend.
        if _vault_backend == 'rest':
            if not (_vault_rest_url and _vault_rest_key):
                return self._send_json(503, {'error': 'Obsidian REST not configured (URL + API key)'})
        else:
            if not _vault_root:
                return self._send_json(503, {'error': 'vault not configured — POST /vault/configure {"root": "..."}'})

        # ---------- List ----------
        if path == '/vault/list' and method == 'GET':
            if _vault_backend == 'rest':
                try:
                    return self._send_json(200, {'files': _rest_list_all()[:500]})
                except Exception as e:
                    return self._send_json(502, {'error': f'obsidian: {e}'})
            if _vault_backend == 'oci':
                try:
                    cli = _oci_object_client()
                    prefix = (_oci_vault_prefix.rstrip('/') + '/') if _oci_vault_prefix else ''
                    resp = cli.list_objects(
                        _oci_vault_namespace, _oci_vault_bucket,
                        prefix=prefix, fields='name,size,timeModified', limit=1000)
                    files = []
                    for obj in resp.data.objects:
                        rel = obj.name[len(prefix):]  # strip user prefix
                        if not rel.endswith('.md') or rel.startswith('.'):
                            continue
                        mtime = 0
                        if obj.time_modified:
                            try: mtime = int(obj.time_modified.timestamp() * 1000)
                            except Exception: pass
                        files.append({
                            'path': rel,
                            'title': rel.rsplit('/', 1)[-1].removesuffix('.md'),
                            'mtime': mtime,
                            'size': obj.size or 0,
                        })
                    files.sort(key=lambda f: f['mtime'], reverse=True)
                    return self._send_json(200, {'files': files[:500]})
                except Exception as e:
                    return self._send_json(502, {'error': f'oci: {e}'})
            root = pathlib.Path(_vault_root).resolve()
            files = []
            for p in root.rglob('*.md'):
                try:
                    rel = str(p.relative_to(root)).replace('\\', '/')
                except ValueError:
                    continue
                if any(part.startswith('.') for part in p.relative_to(root).parts):
                    continue
                try:
                    st = p.stat()
                    files.append({'path': rel, 'title': p.stem,
                                  'mtime': int(st.st_mtime * 1000), 'size': st.st_size})
                except OSError:
                    continue
            files.sort(key=lambda f: f['mtime'], reverse=True)
            return self._send_json(200, {'files': files[:500]})

        # ---------- Read ----------
        if path == '/vault/note' and method == 'GET':
            rel = qs.get('path', [''])[0]
            if _vault_backend == 'rest':
                try:
                    s, hdrs, body = _obsidian_request(
                        'GET', '/vault/' + urllib.parse.quote(rel),
                        extra_headers={'Accept': 'text/markdown'})
                except Exception as e:
                    return self._send_json(502, {'error': f'obsidian: {e}'})
                if s == 404:
                    return self._send_json(404, {'error': 'not found'})
                if s != 200:
                    return self._send_json(s, {'error': body[:300].decode('utf-8', 'replace')})
                self.send_response(200)
                self.send_header('content-type', 'text/markdown; charset=utf-8')
                self.end_headers()
                try: self.wfile.write(body)
                except (BrokenPipeError, ConnectionResetError): pass
                return
            if _vault_backend == 'oci':
                try:
                    cli = _oci_object_client()
                    resp = cli.get_object(_oci_vault_namespace, _oci_vault_bucket, _oci_obj_key(rel))
                    content = resp.data.content
                    self.send_response(200)
                    self.send_header('content-type', 'text/markdown; charset=utf-8')
                    self.end_headers()
                    try: self.wfile.write(content)
                    except (BrokenPipeError, ConnectionResetError): pass
                    return
                except Exception as e:
                    err = str(e)
                    if '404' in err or 'NoSuchKey' in err or 'ObjectNotFound' in err:
                        return self._send_json(404, {'error': 'not found'})
                    return self._send_json(502, {'error': f'oci: {e}'})
            try:
                target = _vault_resolve(rel)
            except ValueError as e:
                return self._send_json(400, {'error': str(e)})
            if not target.exists():
                return self._send_json(404, {'error': 'not found'})
            try:
                text = target.read_text(encoding='utf-8')
            except Exception as e:
                return self._send_json(500, {'error': str(e)})
            self.send_response(200)
            self.send_header('content-type', 'text/markdown; charset=utf-8')
            self.end_headers()
            try: self.wfile.write(text.encode('utf-8'))
            except (BrokenPipeError, ConnectionResetError): pass
            return

        # ---------- Write/Append ----------
        if path == '/vault/note' and method == 'PUT':
            rel = qs.get('path', [''])[0]
            mode = qs.get('mode', ['write'])[0]
            length = int(self.headers.get('content-length', 0) or 0)
            body = self.rfile.read(length).decode('utf-8') if length else ''
            if _vault_backend == 'rest':
                # PUT replaces, POST appends.
                http_method = 'POST' if mode == 'append' else 'PUT'
                try:
                    s, _h, resp = _obsidian_request(
                        http_method, '/vault/' + urllib.parse.quote(rel),
                        body=body.encode('utf-8'), content_type='text/markdown')
                except Exception as e:
                    return self._send_json(502, {'error': f'obsidian: {e}'})
                if s not in (200, 204):
                    return self._send_json(s, {'error': resp[:300].decode('utf-8', 'replace')})
                return self._send_json(200, {'path': rel, 'mode': mode, 'backend': 'rest'})
            if _vault_backend == 'oci':
                try:
                    cli = _oci_object_client()
                    key = _oci_obj_key(rel)
                    content = body
                    if mode == 'append':
                        try:
                            existing = cli.get_object(
                                _oci_vault_namespace, _oci_vault_bucket, key).data.content
                            existing_text = existing.decode('utf-8', 'replace')
                            sep = '' if existing_text.endswith('\n') else '\n'
                            content = existing_text + sep + body
                        except Exception:
                            pass  # file doesn't exist yet — treat as write
                    encoded = content.encode('utf-8') if isinstance(content, str) else content
                    cli.put_object(_oci_vault_namespace, _oci_vault_bucket, key,
                                   put_object_body=encoded)
                    return self._send_json(200, {
                        'path': rel, 'mode': mode,
                        'size': len(encoded), 'backend': 'oci'})
                except Exception as e:
                    return self._send_json(502, {'error': f'oci: {e}'})
            try:
                target = _vault_resolve(rel)
            except ValueError as e:
                return self._send_json(400, {'error': str(e)})
            try:
                target.parent.mkdir(parents=True, exist_ok=True)
                if mode == 'append' and target.exists():
                    existing = target.read_text(encoding='utf-8')
                    sep = '' if existing.endswith('\n') else '\n'
                    target.write_text(existing + sep + body, encoding='utf-8')
                else:
                    target.write_text(body, encoding='utf-8')
            except Exception as e:
                return self._send_json(500, {'error': str(e)})
            return self._send_json(200, {'path': rel, 'mode': mode, 'size': target.stat().st_size})

        # ---------- Delete ----------
        if path == '/vault/note' and method == 'DELETE':
            rel = qs.get('path', [''])[0]
            if _vault_backend == 'rest':
                try:
                    s, _h, resp = _obsidian_request('DELETE', '/vault/' + urllib.parse.quote(rel))
                except Exception as e:
                    return self._send_json(502, {'error': f'obsidian: {e}'})
                if s not in (200, 204):
                    return self._send_json(s, {'error': resp[:300].decode('utf-8','replace')})
                return self._send_json(200, {'deleted': rel})
            if _vault_backend == 'oci':
                try:
                    cli = _oci_object_client()
                    cli.delete_object(_oci_vault_namespace, _oci_vault_bucket, _oci_obj_key(rel))
                except Exception:
                    pass  # 404 on delete is fine
                return self._send_json(200, {'deleted': rel, 'backend': 'oci'})
            try:
                target = _vault_resolve(rel)
            except ValueError as e:
                return self._send_json(400, {'error': str(e)})
            if target.exists():
                try: target.unlink()
                except Exception as e: return self._send_json(500, {'error': str(e)})
            return self._send_json(200, {'deleted': rel})

        # ---------- Rename / move ----------
        if path == '/vault/rename' and method == 'POST':
            length = int(self.headers.get('content-length', 0) or 0)
            try:
                req = json.loads(self.rfile.read(length) or b'{}')
            except Exception:
                return self._send_json(400, {'error': 'bad json'})
            src = str(req.get('from', '')).strip()
            dst = str(req.get('to', '')).strip()
            if not src or not dst:
                return self._send_json(400, {'error': 'need from + to'})
            if _vault_backend == 'rest':
                # Obsidian REST has no native move: copy then delete.
                try:
                    s, _h, body = _obsidian_request('GET', '/vault/' + urllib.parse.quote(src))
                    if s != 200:
                        return self._send_json(404, {'error': f'source not found ({s})'})
                    s2, _h2, resp = _obsidian_request(
                        'PUT', '/vault/' + urllib.parse.quote(dst),
                        body=body, content_type='text/markdown')
                    if s2 not in (200, 204):
                        return self._send_json(s2, {'error': resp[:300].decode('utf-8', 'replace')})
                    _obsidian_request('DELETE', '/vault/' + urllib.parse.quote(src))
                except Exception as e:
                    return self._send_json(502, {'error': f'obsidian: {e}'})
                return self._send_json(200, {'from': src, 'to': dst, 'backend': 'rest'})
            if _vault_backend == 'oci':
                try:
                    cli = _oci_object_client()
                    data = cli.get_object(_oci_vault_namespace, _oci_vault_bucket,
                                          _oci_obj_key(src)).data.content
                    cli.put_object(_oci_vault_namespace, _oci_vault_bucket,
                                   _oci_obj_key(dst), put_object_body=data)
                    cli.delete_object(_oci_vault_namespace, _oci_vault_bucket, _oci_obj_key(src))
                except Exception as e:
                    return self._send_json(502, {'error': f'oci: {e}'})
                return self._send_json(200, {'from': src, 'to': dst, 'backend': 'oci'})
            try:
                s_path = _vault_resolve(src)
                d_path = _vault_resolve(dst)
            except ValueError as e:
                return self._send_json(400, {'error': str(e)})
            if not s_path.exists():
                return self._send_json(404, {'error': 'source not found'})
            if d_path.exists():
                return self._send_json(409, {'error': 'target already exists'})
            try:
                d_path.parent.mkdir(parents=True, exist_ok=True)
                os.replace(str(s_path), str(d_path))
            except Exception as e:
                return self._send_json(500, {'error': str(e)})
            return self._send_json(200, {'from': src, 'to': dst})

        # ---------- Upload (multipart) ----------
        # The browser-side file picker / drag-drop lands here. Filenames are
        # flattened + sanitized; ?dir= picks a target folder inside the vault.
        if path == '/vault/upload' and method == 'POST':
            ctype = self.headers.get('content-type', '')
            if 'multipart/form-data' not in ctype:
                return self._send_json(400, {'error': 'expected multipart/form-data'})
            length = int(self.headers.get('content-length', 0) or 0)
            if length <= 0:
                return self._send_json(400, {'error': 'empty upload'})
            if length > 50 * 1024 * 1024:
                return self._send_json(413, {'error': 'upload too large (50 MB max)'})
            raw = self.rfile.read(length)
            import email.parser as _ep
            import re as _re
            msg = _ep.BytesParser().parsebytes(
                b'Content-Type: ' + ctype.encode('latin-1', 'replace') + b'\r\n\r\n' + raw)
            if not msg.is_multipart():
                return self._send_json(400, {'error': 'bad multipart body'})
            folder = qs.get('dir', [''])[0].strip().strip('/')
            saved, errors = [], []
            for part in msg.get_payload():
                fname = part.get_filename()
                if not fname:
                    continue
                fname = fname.replace('\\', '/').split('/')[-1]
                fname = _re.sub(r'[^\w .()\[\]\-]+', '_', fname).strip()
                if not fname or fname.startswith('.'):
                    continue
                data = part.get_payload(decode=True) or b''
                rel = (folder + '/' if folder else '') + fname
                try:
                    if _vault_backend == 'rest':
                        s, _h, resp = _obsidian_request(
                            'PUT', '/vault/' + urllib.parse.quote(rel),
                            body=data, content_type='application/octet-stream')
                        if s not in (200, 204):
                            raise RuntimeError(resp[:200].decode('utf-8', 'replace'))
                    elif _vault_backend == 'oci':
                        cli = _oci_object_client()
                        cli.put_object(_oci_vault_namespace, _oci_vault_bucket,
                                       _oci_obj_key(rel), put_object_body=data)
                    else:
                        target = _vault_resolve(rel)
                        target.parent.mkdir(parents=True, exist_ok=True)
                        target.write_bytes(data)
                    saved.append({'path': rel, 'size': len(data)})
                except Exception as e:
                    errors.append({'path': rel, 'error': str(e)})
            if not saved and errors:
                return self._send_json(500, {'error': errors[0]['error'], 'failed': errors})
            return self._send_json(200, {'uploaded': saved, 'count': len(saved),
                                         **({'failed': errors} if errors else {})})

        # ---------- Search ----------
        if path == '/vault/search' and method == 'GET':
            query = qs.get('q', [''])[0].strip()
            limit = int(qs.get('limit', ['10'])[0])
            if not query:
                return self._send_json(400, {'error': 'missing q'})
            if _vault_backend == 'rest':
                try:
                    return self._send_json(200, {'hits': _rest_search(query, limit)})
                except Exception as e:
                    return self._send_json(502, {'error': f'obsidian: {e}'})
            root = pathlib.Path(_vault_root).resolve()
            ql = query.lower()
            hits = []
            for p in root.rglob('*.md'):
                try:
                    rel = str(p.relative_to(root)).replace('\\', '/')
                except ValueError:
                    continue
                if any(part.startswith('.') for part in p.relative_to(root).parts):
                    continue
                try:
                    text = p.read_text(encoding='utf-8', errors='replace')
                except OSError:
                    continue
                tl = text.lower()
                title_score = 3 if ql in p.stem.lower() else 0
                count = tl.count(ql)
                if not (title_score or count):
                    continue
                idx = tl.find(ql)
                snippet = ''
                if idx >= 0:
                    s = max(0, idx - 60)
                    e = min(len(text), idx + len(query) + 60)
                    snippet = ('…' if s > 0 else '') + text[s:e].replace('\n', ' ').strip() + ('…' if e < len(text) else '')
                hits.append({'path': rel, 'title': p.stem,
                             'score': title_score + count, 'snippet': snippet})
            hits.sort(key=lambda h: h['score'], reverse=True)
            return self._send_json(200, {'hits': hits[:limit], 'total': len(hits)})

        # ---------- Graph (nodes + wikilink edges) ----------
        if path == '/vault/graph' and method == 'GET':
            try:
                graph = _build_graph_rest() if _vault_backend == 'rest' else _build_graph_fs_cached()
                return self._send_json(200, graph)
            except Exception as e:
                return self._send_json(502, {'error': str(e)})

        # ---------- Open in Obsidian (REST only) ----------
        if path == '/vault/open' and method == 'POST':
            if _vault_backend != 'rest':
                return self._send_json(400, {'error': 'open-in-Obsidian requires REST backend'})
            length = int(self.headers.get('content-length', 0) or 0)
            try:
                body = json.loads(self.rfile.read(length).decode('utf-8') or '{}')
            except json.JSONDecodeError:
                return self._send_json(400, {'error': 'bad json'})
            rel = (body.get('path') or '').strip()
            if not rel:
                return self._send_json(400, {'error': 'missing path'})
            try:
                s, _h, resp = _obsidian_request('POST', '/open/' + urllib.parse.quote(rel))
            except Exception as e:
                return self._send_json(502, {'error': f'obsidian: {e}'})
            return self._send_json(s if s != 200 else 200,
                                   {'opened': rel} if s == 200 else {'error': resp[:200].decode('utf-8','replace')})

        return self._send_json(404, {'error': f'unknown vault route: {path}'})

    def _brave_search(self):
        """Forward /brave/search?q=... to Brave Web Search.

        Auth: caller sends X-Brave-Key, we send X-Subscription-Token.
        Falls back to BRAVE_API_KEY env var if header is absent.
        """
        if not self.path.startswith('/brave/search'):
            return self._send_json(404, {'error': 'unknown brave endpoint'})
        key = self.headers.get('X-Brave-Key') or os.environ.get('BRAVE_API_KEY', '')
        if not key:
            return self._send_json(401, {'error': 'no Brave key (set X-Brave-Key header or BRAVE_API_KEY env)'})
        # Preserve all query params after /brave/search
        _, _, qs = self.path.partition('?')
        upstream_path = '/res/v1/web/search' + ('?' + qs if qs else '')
        ctx = ssl.create_default_context()
        conn = http.client.HTTPSConnection('api.search.brave.com', 443, timeout=30, context=ctx)
        try:
            conn.request('GET', upstream_path, headers={
                'Accept': 'application/json',
                'Accept-Encoding': 'identity',
                'X-Subscription-Token': key,
                'User-Agent': 'CafresoHQ/1.0',
            })
            resp = conn.getresponse()
            self.send_response(resp.status)
            for k, v in resp.getheaders():
                if k.lower() in HOP_HEADERS or k.lower() == 'content-encoding':
                    continue
                self.send_header(k, v)
            self.end_headers()
            while True:
                chunk = resp.read(2048)
                if not chunk:
                    break
                try:
                    self.wfile.write(chunk)
                except (BrokenPipeError, ConnectionResetError):
                    break
        except Exception as e:
            try: self._send_json(502, {'error': f'brave: {e}'})
            except Exception: pass
        finally:
            conn.close()

    def _proxy(self, method, prefix, target):
        UPSTREAM_HOST, UPSTREAM_PORT = target
        upstream_path = self.path[len(prefix) - 1:]  # keep leading /
        length = int(self.headers.get('content-length', 0) or 0)
        body = self.rfile.read(length) if length else None

        # Strip Accept-Encoding so the upstream doesn't gzip/compress, which
        # would add buffering and break streaming (SSE, etc.).
        _drop = HOP_HEADERS | {'accept-encoding'}
        headers = {k: v for k, v in self.headers.items()
                   if k.lower() not in _drop}
        headers['host'] = f'{UPSTREAM_HOST}:{UPSTREAM_PORT}'
        headers['Accept-Encoding'] = 'identity'

        conn = http.client.HTTPConnection(UPSTREAM_HOST, UPSTREAM_PORT,
                                          timeout=600)
        try:
            conn.request(method, upstream_path, body=body, headers=headers)
            resp = conn.getresponse()
            self.send_response(resp.status)
            for k, v in resp.getheaders():
                if k.lower() in HOP_HEADERS or k.lower() == 'content-encoding':
                    continue
                self.send_header(k, v)
            self.send_header('Connection', 'close')
            self.end_headers()
            # Use the underlying raw socket (resp.fp.read / resp.fp.read1)
            # to avoid BufferedReader.read(n) blocking until it accumulates
            # exactly n bytes — which kills SSE streaming.  HTTPResponse.read
            # wraps a BufferedReader; read1 returns as soon as ANY data is
            # available in the buffer (at most n bytes), which is exactly the
            # behaviour we need for real-time proxy streaming.
            raw_fp = resp  # HTTPResponse itself
            while True:
                try:
                    # read1 = "one underlying read, return whatever we got"
                    chunk = raw_fp.fp.read1(8192) if hasattr(raw_fp.fp, 'read1') else raw_fp.read(1024)
                except Exception:
                    break
                if not chunk:
                    break
                try:
                    self.wfile.write(chunk)
                    self.wfile.flush()
                except (BrokenPipeError, ConnectionResetError):
                    break
        except Exception as e:
            try:
                self._send_json(502, {'error': f'upstream: {e}'})
            except Exception:
                pass
        finally:
            conn.close()

    # ── Hermes capability mode (lite ↔ full) ─────────────────────────────────
    # Hermes injects its full toolset into the system prompt (~14-17k tokens),
    # which exceeds free hosted tiers' request-size caps (e.g. Groq free ~5-6k →
    # HTTP 413). 'lite' forces tool_search deferral + trims the preamble so the
    # prompt fits the free tier; 'full' restores the rich prompt for BYOK/paid
    # keys that can afford it. The HQ Settings toggle drives this.
    def _hermes_capability_file(self):
        import os as _os
        home = _os.environ.get('HERMES_HOME', '').strip() or _os.path.expanduser('~/.hermes')
        return _os.path.join(home, 'capability_mode')

    # ── Hermes model quick-switch ─────────────────────────────────────────────
    # GET  /hermes/model            → {model, presets}
    # POST /hermes/model {model}    → rewrite config.yaml model.default + restart
    # Presets are curated OpenRouter free open-weights ids verified to accept
    # Hermes' large prompt. The UI offers these as one-click switches.
    _HERMES_MODEL_PRESETS = [
        {'id': 'openai/gpt-oss-120b:free',                  'label': 'GPT-OSS 120B (default)'},
        {'id': 'nvidia/nemotron-3-super-120b-a12b:free',    'label': 'Nemotron 3 Super 120B'},
        {'id': 'nousresearch/hermes-3-llama-3.1-405b:free', 'label': 'Hermes 3 405B (Nous)'},
        {'id': 'meta-llama/llama-3.3-70b-instruct:free',    'label': 'Llama 3.3 70B'},
        {'id': 'qwen/qwen3-next-80b-a3b-instruct:free',     'label': 'Qwen3-Next 80B'},
    ]

    def _hermes_config_path(self):
        import os as _os
        home = _os.environ.get('HERMES_HOME', '').strip() or _os.path.expanduser('~/.hermes')
        return _os.path.join(home, 'config.yaml')

    def _hermes_get_model(self):
        import re as _re
        model = ''
        try:
            with open(self._hermes_config_path(), 'r', encoding='utf-8') as f:
                m = _re.search(r'^\s*default:\s*(.+)\s*$', f.read(), _re.MULTILINE)
                if m:
                    model = m.group(1).strip()
        except Exception:
            pass
        return self._send_json(200, {'model': model, 'presets': self._HERMES_MODEL_PRESETS})

    def _hermes_set_model(self):
        """POST {model} → rewrite config.yaml model.default + restart gateway.
        Only the model id changes; provider/base_url (OpenRouter) are preserved."""
        length = int(self.headers.get('content-length', 0) or 0)
        try:
            req = json.loads(self.rfile.read(length) or b'{}')
        except Exception:
            return self._send_json(400, {'error': 'bad json'})
        model = str(req.get('model', '')).strip()
        # Allow presets OR any plausible model id. The vendor slash is OPTIONAL:
        # cloud ids look like "vendor/model[:tag]" but local ones often don't —
        # Ollama's are bare ("llama3.3:70b"), and requiring the slash 400'd every
        # local backend.
        import re as _re
        valid = any(p['id'] == model for p in self._HERMES_MODEL_PRESETS) or \
            bool(_re.match(r'^[\w.\-]+(/[\w.\-]+)*(:[\w.\-]+)?$', model))
        if not model or not valid:
            return self._send_json(400, {'error': 'invalid model id'})

        cfg_path = self._hermes_config_path()
        try:
            with open(cfg_path, 'r', encoding='utf-8') as f:
                cfg = f.read()
        except Exception as e:
            return self._send_json(500, {'error': f'read config: {e}'})

        new_cfg, n = _re.subn(r'(^\s*default:\s*).+$',
                              lambda m: m.group(1) + model, cfg,
                              count=1, flags=_re.MULTILINE)
        if n == 0:
            return self._send_json(500, {'error': 'no model.default line in config'})
        try:
            with open(cfg_path, 'w', encoding='utf-8') as f:
                f.write(new_cfg)
        except Exception as e:
            return self._send_json(500, {'error': f'write config: {e}'})

        restarted = False
        try:
            subprocess.Popen(['hermes', 'gateway', 'restart'],
                             stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            restarted = True
        except Exception as e:
            sys.stderr.write(f'[hermes] model restart failed: {e}\n')
        return self._send_json(200, {'model': model, 'restarted': restarted,
                                     'note': 'gateway reloading; allow ~10s'})

    def _hermes_get_capability(self):
        try:
            with open(self._hermes_capability_file(), 'r', encoding='utf-8') as f:
                mode = (f.read().strip() or 'lite')
        except Exception:
            mode = 'lite'
        return self._send_json(200, {'mode': mode if mode in ('lite', 'full') else 'lite'})

    def _hermes_set_capability(self):
        """POST {mode: 'lite'|'full'} → rewrite config.yaml's capability block and
        restart the hermes gateway so the new system-prompt size takes effect."""
        length = int(self.headers.get('content-length', 0) or 0)
        try:
            req = json.loads(self.rfile.read(length) or b'{}')
        except Exception:
            return self._send_json(400, {'error': 'bad json'})
        mode = str(req.get('mode', '')).strip().lower()
        if mode not in ('lite', 'full'):
            return self._send_json(400, {'error': "mode must be 'lite' or 'full'"})

        import os as _os
        home = _os.environ.get('HERMES_HOME', '').strip() or _os.path.expanduser('~/.hermes')
        cfg_path = _os.path.join(home, 'config.yaml')
        try:
            with open(cfg_path, 'r', encoding='utf-8') as f:
                cfg = f.read()
        except Exception as e:
            return self._send_json(500, {'error': f'read config: {e}'})

        # Replace the toolsets/agent/tools tail (everything from the first
        # 'toolsets:' line) with the requested capability block. The block layout
        # matches hermes-bootstrap.py:_capability_block so the two stay in sync.
        if mode == 'full':
            block = ('toolsets:\n  - hermes-cli\n'
                     'agent:\n  environment_probe: true\n  task_completion_guidance: true\n'
                     'tools:\n  tool_search:\n    enabled: auto\n')
        else:
            block = ('toolsets:\n  - hermes-cli\n'
                     'agent:\n  environment_probe: false\n  task_completion_guidance: false\n'
                     'tools:\n  tool_search:\n    enabled: true\n    threshold_pct: 0\n')

        idx = cfg.find('\ntoolsets:')
        new_cfg = (cfg[:idx + 1] if idx >= 0 else cfg.rstrip() + '\n') + block
        try:
            with open(cfg_path, 'w', encoding='utf-8') as f:
                f.write(new_cfg)
            with open(self._hermes_capability_file(), 'w', encoding='utf-8') as f:
                f.write(mode)
        except Exception as e:
            return self._send_json(500, {'error': f'write config: {e}'})

        # Restart the gateway so it reloads config.yaml. Best-effort; the proxy
        # keeps serving until the new gateway binds. `hermes gateway restart`
        # replaces the running singleton.
        restarted = False
        try:
            subprocess.Popen(['hermes', 'gateway', 'restart'],
                             stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            restarted = True
        except Exception as e:
            sys.stderr.write(f'[hermes] capability restart failed: {e}\n')
        return self._send_json(200, {'mode': mode, 'restarted': restarted,
                                     'note': 'gateway reloading; allow ~10s'})

    # ── Backend providers the user can pick in HQ Settings ──────────────────
    # Each logical provider maps to the Hermes model block + the env var that
    # carries its API key. OpenRouter is the zero-config default; **Gemini
    # (direct) is the most RELIABLE free tier** (≈15 RPM / 1500 RPD on Flash vs
    # OpenRouter :free's 20 RPM / 50 RPD — the documented reliability pain);
    # Groq is fast + free. Gemini + Groq ride their OpenAI-compatible endpoints
    # (custom_providers) so they use the exact chat_completions path OpenRouter
    # uses — no unverified native adapter, no silent config keys.
    # 'local' backends are the operator's own hardware: no key, reached at a
    # base_url they supply. They write `provider: lmstudio` either way — there is
    # no `ollama` provider in hermes-agent (bootstrap's model block is explicit
    # that Ollama rides the lmstudio block as a generic OpenAI-compatible
    # endpoint), and since hermes is a pinned third-party package we can't add
    # one. 'ollama' is a UI label, not a config value.
    _HERMES_PROVIDERS = {
        'openrouter': {'env': 'OPENROUTER_API_KEY', 'model': 'openai/gpt-oss-120b:free',
                       're': r'^sk-or-[A-Za-z0-9_\-]{8,}$', 'label': 'OpenRouter'},
        'gemini':     {'env': 'GOOGLE_API_KEY',     'model': 'gemini-2.5-flash',
                       're': r'^AIza[A-Za-z0-9_\-]{30,}$', 'label': 'Google Gemini'},
        'groq':       {'env': 'GROQ_API_KEY',       'model': 'llama-3.3-70b-versatile',
                       're': r'^gsk_[A-Za-z0-9]{20,}$', 'label': 'Groq'},
        'lmstudio':   {'env': '', 'model': 'local-model', 're': None, 'local': True,
                       'label': 'LM Studio (local)', 'default_url': 'http://localhost:1234/v1'},
        'ollama':     {'env': '', 'model': 'llama3.1', 're': None, 'local': True,
                       'label': 'Ollama (local)', 'default_url': 'http://localhost:11434/v1'},
    }

    @staticmethod
    def _hermes_model_block(provider, model, base_url=None):
        """Return the config.yaml model block (+ custom_providers) for a provider."""
        if provider in ('lmstudio', 'ollama'):
            # Mirrors oci-fleet/hermes-bootstrap.py's lmstudio block exactly:
            # no key_env, no custom_providers. Both UI labels write `lmstudio`.
            return (f'model:\n  default: {model}\n  provider: lmstudio\n'
                    f'  base_url: {base_url}\n')
        if provider == 'gemini':
            return (f'model:\n  default: {model}\n  provider: google-openai\n'
                    '  base_url: https://generativelanguage.googleapis.com/v1beta/openai\n'
                    'custom_providers:\n'
                    '  - name: google-openai\n'
                    '    base_url: https://generativelanguage.googleapis.com/v1beta/openai\n'
                    '    key_env: GOOGLE_API_KEY\n'
                    '    api_mode: chat_completions\n')
        if provider == 'groq':
            return (f'model:\n  default: {model}\n  provider: groq\n'
                    '  base_url: https://api.groq.com/openai/v1\n'
                    'custom_providers:\n'
                    '  - name: groq\n'
                    '    base_url: https://api.groq.com/openai/v1\n'
                    '    key_env: GROQ_API_KEY\n'
                    '    api_mode: chat_completions\n')
        # openrouter — Hermes' native default provider (no custom_providers)
        return (f'model:\n  default: {model}\n  provider: openrouter\n'
                '  base_url: https://openrouter.ai/api/v1\n')

    def _hermes_local_models(self, query):
        """GET /hermes/local-models?base_url=… → {models:[{id, state, loaded}], detail}

        Which models the operator's local backend actually has, and crucially
        WHICH ARE LOADED. Asking a local backend for an unloaded model fails
        every call; showing loaded state in the picker makes that unpickable
        instead of a mystery.

        The browser can't fetch this itself: the /lmstudio/ proxy is hardcoded to
        one host, so it can't serve a URL the operator just typed. Runs the same
        allowlist as the write path — discovery must not become the SSRF hole the
        writer isn't. Always 200: a backend that's merely offline must not block
        the operator from saving a config for it.
        """
        base_url = urllib.parse.parse_qs(query or '').get('base_url', [''])[0].strip()
        ok, err = _validate_local_base_url(base_url)
        if not ok:
            return self._send_json(400, {'error': err})
        root = base_url.rstrip('/')
        # LM Studio's native REST API carries load state; the OpenAI-compatible
        # /models (which Ollama also answers) does not.
        native = re.sub(r'/v\d+$', '', root) + '/api/v0/models'
        for url, has_state in ((native, True), (root + '/models', False)):
            try:
                with urllib.request.urlopen(url, timeout=3) as r:
                    data = json.loads(r.read().decode('utf-8'))
            except Exception:
                continue
            out = []
            for m in (data.get('data') or []):
                mid = m.get('id')
                if not mid:
                    continue
                if has_state:
                    out.append({'id': mid, 'state': m.get('state'),
                                'loaded': m.get('state') == 'loaded',
                                'type': m.get('type')})
                else:
                    out.append({'id': mid, 'state': None, 'loaded': None,
                                'type': m.get('type')})
            if out:
                # Loaded first — those are the ones that answer immediately.
                out.sort(key=lambda x: (x['loaded'] is not True, x['id']))
                return self._send_json(200, {'models': out, 'source': url})
        return self._send_json(200, {'models': [], 'detail': 'no model list at %s' % root})

    def _hermes_get_provider(self):
        """GET /hermes/provider → {provider, model, configured}. The UI reads this
        on load to decide whether to re-push the user's saved key: a container
        recreate wipes the ephemeral ~/.hermes, so the key must be re-applied from
        the browser-side settings copy (that's the 'keys vanish on recreate' fix)."""
        import os as _os, re as _re
        import night_runner as _nr
        home = _hermes_home()
        cfg = _nr.read_model_config(home)          # scoped to the model: block
        provider = cfg['raw_provider'] or 'openrouter'
        model, base_url = cfg['model'], cfg['base_url']
        logical = {'google-openai': 'gemini'}.get(provider, provider)
        spec = self._HERMES_PROVIDERS.get(logical)
        configured = False
        if spec and spec.get('local'):
            # A local backend has no key — having a base_url IS being configured.
            configured = bool(base_url)
        elif spec:
            if _os.environ.get(spec['env'], '').strip():
                configured = True
            else:
                try:
                    with open(_os.path.join(home, '.env'), 'r', encoding='utf-8') as f:
                        configured = bool(_re.search(
                            r'(?m)^%s\s*=\s*\S' % _re.escape(spec['env']), f.read()))
                except Exception:
                    pass
        return self._send_json(200, {'provider': logical, 'model': model,
                                     'base_url': base_url, 'configured': configured})

    def _hermes_set_provider(self, force_provider=None):
        """POST /hermes/provider {provider, key, model?} → write the provider's key
        into ~/.hermes/.env, REWRITE config.yaml's model block to that provider, and
        restart the gateway. Generalizes the old OpenRouter-only path so users can
        switch to a more reliable free backend (Gemini direct / Groq) with their own
        free key. /hermes/openrouter-key routes here with force_provider='openrouter'.

        Unlike first-boot bootstrap (skip-if-exists), this REWRITES the model block
        so switching providers actually changes the live backend. The capability
        tail (toolsets/agent/tools) is preserved. Key lives only in the container
        .env (0600), never persisted server-side beyond that file."""
        length = int(self.headers.get('content-length', 0) or 0)
        try:
            req = json.loads(self.rfile.read(length) or b'{}')
        except Exception:
            return self._send_json(400, {'error': 'bad json'})
        provider = (force_provider or str(req.get('provider', 'openrouter'))).strip().lower()
        spec = self._HERMES_PROVIDERS.get(provider)
        if not spec:
            return self._send_json(400, {'error': f'unknown provider: {provider}'})
        key = str(req.get('key', '')).strip()
        import os as _os, re as _re
        local = bool(spec.get('local'))
        base_url = str(req.get('base_url', '')).strip() or spec.get('default_url')
        if local:
            # No key to validate — the gate here is the URL, which the container
            # will go on to POST to. See _validate_local_base_url: allowlist.
            ok, err = _validate_local_base_url(base_url)
            if not ok:
                return self._send_json(400, {'error': err})
        elif not _re.match(spec['re'], key):
            return self._send_json(400, {'error': f"invalid {spec['label']} key"})
        model = str(req.get('model', '')).strip() or spec['model']

        home = _os.environ.get('HERMES_HOME', '').strip() or _os.path.expanduser('~/.hermes')
        env_path = _os.path.join(home, '.env')
        cfg_path = _os.path.join(home, 'config.yaml')
        try:
            _os.makedirs(home, exist_ok=True)
            # 1. write the key into .env (replace any prior line for THIS env var).
            #    A local backend has no key and no env var — skip entirely.
            if not local:
                lines = []
                if _os.path.exists(env_path):
                    with open(env_path, 'r', encoding='utf-8') as f:
                        lines = [l for l in f.read().splitlines()
                                 if not l.startswith(spec['env'] + '=')]
                lines.append(f"{spec['env']}={key}")
                with open(env_path, 'w', encoding='utf-8') as f:
                    f.write('\n'.join(lines) + '\n')
                try:
                    _os.chmod(env_path, 0o600)
                except Exception:
                    pass
        except Exception as e:
            return self._send_json(500, {'error': f'write .env: {e}'})

        # 2. rewrite config.yaml's model block (everything before `approvals:`),
        #    preserving the capability tail so the lite/full toggle survives.
        try:
            cfg = ''
            if _os.path.exists(cfg_path):
                with open(cfg_path, 'r', encoding='utf-8') as f:
                    cfg = f.read()
            header = ('# CafresoHQ — Hermes config (provider set via HQ Settings).\n'
                      '# capability_mode controls system-prompt size (lite=free-tier-safe).\n')
            block = self._hermes_model_block(provider, model, base_url)
            m = _re.search(r'^approvals:', cfg, _re.MULTILINE)
            if m:
                new_cfg = header + block + cfg[m.start():]
            else:
                # fresh/unknown config — write a complete minimal one (lite caps)
                new_cfg = (header + block +
                           'approvals:\n  mode: manual\n'
                           'toolsets:\n  - hermes-cli\n'
                           'agent:\n  environment_probe: false\n  task_completion_guidance: false\n'
                           'tools:\n  tool_search:\n    enabled: true\n    threshold_pct: 0\n')
            with open(cfg_path, 'w', encoding='utf-8') as f:
                f.write(new_cfg)
        except Exception as e:
            return self._send_json(500, {'error': f'write config: {e}'})

        # The user just brought their OWN key (or their own hardware) — end the
        # shared-trial cap. Their key draws on their own account, so it's never
        # metered here.
        _trial_deactivate()

        # 3. export into THIS process env so the restarted gateway (which inherits
        #    serve.py's env via the bootstrap) sees the key immediately.
        if not local:
            _os.environ[spec['env']] = key
        restarted = False
        try:
            subprocess.Popen(['hermes', 'gateway', 'restart'],
                             stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            restarted = True
        except Exception as e:
            sys.stderr.write(f'[hermes] provider restart failed: {e}\n')
        return self._send_json(200, {'ok': True, 'provider': provider, 'model': model,
                                     'restarted': restarted,
                                     'note': 'gateway reloading; allow ~10s'})

    def _hermes_trial_status(self):
        """GET /hermes/trial-status → whether this HQ is on the shared trial
        brain and how much of today's per-principal allowance is left. Drives the
        onboarding upsell ('you're on the free shared brain — add your own key
        for unlimited'). Returns active:false once the user brings their own key."""
        st = _trial_state()
        principal = (self.headers.get('X-User-Principal') or '').strip() or 'local'
        used, cap = _trial_usage(principal)
        return self._send_json(200, {
            'active': st.get('active', False),
            'provider': st.get('provider', ''),
            'used': used,
            'cap': cap,
            'remaining': max(0, cap - used),
        })

    # ── Hermes config import/export ──────────────────────────────────────────
    # Lets users carry a Hermes agent setup between HQs (or in from a local
    # ~/.hermes) in one click. config.yaml holds NO secrets — keys live in
    # .env, which is never exported and never accepted on import.
    def _hermes_export_config(self):
        import os as _os
        home = _os.environ.get('HERMES_HOME', '').strip() or _os.path.expanduser('~/.hermes')
        cfg = ''
        try:
            with open(_os.path.join(home, 'config.yaml'), 'r', encoding='utf-8') as f:
                cfg = f.read()
        except Exception:
            pass
        mode = 'lite'
        try:
            with open(self._hermes_capability_file(), 'r', encoding='utf-8') as f:
                mode = (f.read().strip() or 'lite')
        except Exception:
            pass
        return self._send_json(200, {
            'version': 1,
            'kind': 'cafresohq-hermes-config',
            'capability': mode if mode in ('lite', 'full') else 'lite',
            'config_yaml': cfg,
            'note': 'keys are NOT included — set them in Settings → Connections',
        })

    def _hermes_import_config(self):
        """POST {config_yaml, capability?} → replace ~/.hermes/config.yaml and
        restart the gateway. Refuses key material (keys belong in .env via
        Settings) and obviously-broken payloads. Accepts a raw Hermes
        config.yaml or our export envelope's config_yaml field."""
        length = int(self.headers.get('content-length', 0) or 0)
        if length > 64 * 1024:
            return self._send_json(413, {'error': 'config too large (64 KB max)'})
        try:
            req = json.loads(self.rfile.read(length) or b'{}')
        except Exception:
            return self._send_json(400, {'error': 'bad json'})
        cfg = str(req.get('config_yaml', '') or '')
        if not cfg.strip():
            return self._send_json(400, {'error': 'config_yaml is empty'})
        if 'model:' not in cfg:
            return self._send_json(400, {'error': "doesn't look like a Hermes config (no model: block)"})
        import re as _re
        if _re.search(r'(?im)^\s*[\w-]*(api[_-]?key|secret|token|password)\s*:\s*\S', cfg):
            return self._send_json(400, {
                'error': 'config contains key material — remove it; API keys are set in Settings → Connections'})
        import os as _os
        home = _os.environ.get('HERMES_HOME', '').strip() or _os.path.expanduser('~/.hermes')
        cfg_path = _os.path.join(home, 'config.yaml')
        try:
            _os.makedirs(home, exist_ok=True)
            # Keep one rollback copy in case the imported config breaks the gateway.
            if _os.path.exists(cfg_path):
                try:
                    with open(cfg_path, 'r', encoding='utf-8') as f:
                        prev = f.read()
                    with open(cfg_path + '.bak', 'w', encoding='utf-8') as f:
                        f.write(prev)
                except Exception:
                    pass
            with open(cfg_path, 'w', encoding='utf-8') as f:
                f.write(cfg)
        except Exception as e:
            return self._send_json(500, {'error': f'write config: {e}'})
        cap = str(req.get('capability', '')).strip().lower()
        if cap in ('lite', 'full'):
            try:
                with open(self._hermes_capability_file(), 'w', encoding='utf-8') as f:
                    f.write(cap)
            except Exception:
                pass
        restarted = False
        try:
            subprocess.Popen(['hermes', 'gateway', 'restart'],
                             stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            restarted = True
        except Exception as e:
            sys.stderr.write(f'[hermes] import restart failed: {e}\n')
        return self._send_json(200, {'ok': True, 'restarted': restarted,
                                     'rollback': cfg_path + '.bak',
                                     'note': 'gateway reloading; allow ~10s'})

    def _hermes_proxy(self, method):
        """Proxy /hermes/* → the local `hermes gateway` OpenAI-compatible API
        server (HERMES_HOST:HERMES_PORT, default 127.0.0.1:8642).

        /hermes/v1/chat/completions → http://127.0.0.1:8642/v1/chat/completions

        Auth: API_SERVER_KEY env var is injected as 'Authorization: Bearer …'
        server-side and takes priority over any header the browser sends, so
        the key never lives in the browser. Falls back to the client's
        Authorization header if the env var is absent (dev convenience).
        Accept-Encoding is forced to 'identity' so SSE arrives uncompressed
        and streams straight through. A bounded tail of the response is scanned
        for the OpenAI `usage` object and recorded per-principal (metering).
        """
        upstream_path = self.path[len('/hermes'):]  # keep leading /
        if not upstream_path.startswith('/'):
            upstream_path = '/' + upstream_path

        length = int(self.headers.get('content-length', 0) or 0)
        body = self.rfile.read(length) if length else None

        # Trial-brain cap: when this HQ is running on the shared trial key, meter
        # chat completions per-principal per-day so one user can't drain the
        # shared pool. A user's OWN key clears the trial (below) and is never
        # metered. Body is already consumed, so a 429 here is a clean response.
        if method == 'POST' and upstream_path.startswith('/v1/chat/completions') \
                and _trial_state().get('active'):
            _principal = (self.headers.get('X-User-Principal') or '').strip() or 'local'
            _ok, _used, _cap = _trial_check_and_bump(_principal)
            if not _ok:
                return self._send_json(429, {
                    'error': 'trial_limit', 'used': _used, 'cap': _cap,
                    'message': (f"You've used your {_cap} free messages for today. "
                                "Add your own free key in Settings — it's quick and "
                                "gives you unlimited use."),
                })

        # Also strip browser-context headers. The in-container Hermes gateway's
        # aiohttp server REJECTS any request carrying an Origin it doesn't know
        # with 403 (CORS/CSRF) — and our cross-origin canister UI always sends
        # Origin: https://hq-ui.cafreso.com — so forwarding it broke every Hermes
        # call (chat + /models). The gateway authenticates via the Bearer we
        # inject below, not the browser cookie/origin, so drop all three.
        _drop = HOP_HEADERS | {'host', 'authorization', 'accept-encoding',
                               'origin', 'referer', 'cookie'}
        headers = {k: v for k, v in self.headers.items()
                   if k.lower() not in _drop}
        headers['Host'] = f'{HERMES_HOST}:{HERMES_PORT}'
        headers['Accept-Encoding'] = 'identity'
        # Force the upstream (aiohttp) to close after responding. We relay the
        # body with a raw read1() loop (so SSE streams straight through) which
        # bypasses http.client's Content-Length framing — on a keep-alive
        # response read1() would block waiting for an EOF the gateway never
        # sends, hanging the request (and, since we drop Content-Length as a
        # hop header, the browser's fetch never resolves either). Connection:
        # close makes the gateway send EOF after the body, so the loop ends.
        headers['Connection'] = 'close'

        env_key = os.environ.get('API_SERVER_KEY', '').strip()
        if not env_key:
            # Fallback: read API_SERVER_KEY from ~/.hermes/.env so the proxy
            # works even when the env var wasn't injected at container start.
            try:
                import re as _re
                _hermes_home = os.environ.get('HERMES_HOME', '').strip() or os.path.expanduser('~/.hermes')
                with open(os.path.join(_hermes_home, '.env'), 'r', encoding='utf-8') as _hf:
                    _km = _re.search(r'^API_SERVER_KEY\s*=\s*([^\r\n]+)', _hf.read(), _re.MULTILINE)
                    if _km:
                        env_key = _km.group(1).strip().strip('"\'')
            except Exception:
                pass
        if env_key:
            headers['Authorization'] = 'Bearer ' + env_key
        else:
            client_auth = self.headers.get('Authorization') or self.headers.get('authorization')
            if client_auth:
                headers['Authorization'] = client_auth

        principal = (self.headers.get('X-User-Principal') or '').strip() or 'local'
        # The gateway is briefly unavailable right after a restart (key / model /
        # capability change replaces the running singleton). Retry the upstream
        # CONNECT a few times so the user sees a short delay instead of a transient
        # 502 during that ~10-15s window. Safe to retry: nothing is written to the
        # client until getresponse() succeeds, and the request body is buffered in
        # memory. Only connection-level failures (gateway not listening / dropped)
        # are retried — a real HTTP error comes back as a status, not an exception.
        conn = None
        resp = None
        last_err = None
        for _attempt in range(10):
            conn = http.client.HTTPConnection(HERMES_HOST, HERMES_PORT, timeout=600)
            try:
                conn.request(method, upstream_path, body=body, headers=headers)
                resp = conn.getresponse()
                break
            except (ConnectionRefusedError, ConnectionResetError,
                    http.client.RemoteDisconnected) as e:
                last_err = e
                try:
                    conn.close()
                except Exception:
                    pass
                conn = None
                time.sleep(1.5)
            except Exception as e:
                last_err = e
                break
        if resp is None:
            try:
                if conn:
                    conn.close()
            except Exception:
                pass
            try:
                self._send_json(502, {'error': f'hermes: {last_err or "gateway unavailable"}',
                                      'hint': 'the agent gateway is restarting or down — retry in ~15s'})
            except Exception:
                pass
            return
        try:
            self.send_response(resp.status)
            for k, v in resp.getheaders():
                if k.lower() in HOP_HEADERS or k.lower() == 'content-encoding':
                    continue
                self.send_header(k, v)
            self.send_header('Connection', 'close')
            self.end_headers()
            tail = b''  # rolling buffer (last 16KB) for usage extraction
            while True:
                try:
                    chunk = resp.fp.read1(8192) if hasattr(resp.fp, 'read1') else resp.read(1024)
                except Exception:
                    break
                if not chunk:
                    break
                try:
                    self.wfile.write(chunk)
                    self.wfile.flush()
                except (BrokenPipeError, ConnectionResetError):
                    break
                tail = (tail + chunk)[-16384:]
            self._record_hermes_usage(principal, upstream_path, tail)
        except Exception as e:
            try:
                self._send_json(502, {'error': f'hermes: {e}'})
            except Exception:
                pass
        finally:
            conn.close()

    # Deterministic, writable usage-log path: prefer an explicit override, then
    # the persistent Hermes data dir (mounted volume in the fleet), then /data,
    # then the home dir. Avoids the HOME-ambiguity that silently broke writes.
    def _hermes_usage_log_path(self):
        override = os.environ.get('HERMES_USAGE_LOG', '').strip()
        if override:
            return override
        for base in (os.environ.get('HERMES_HOME', '').strip(), '/data',
                     os.path.expanduser('~')):
            if base and os.path.isdir(base):
                return os.path.join(base, 'hermes-usage.log')
        return os.path.join(os.path.expanduser('~'), 'hermes-usage.log')

    def _record_hermes_usage(self, principal, path, tail_bytes):
        """Best-effort per-principal token accounting. Scans the response tail
        for the LAST OpenAI `usage` object (covers both non-stream JSON bodies
        and the terminal SSE chunk), then appends a JSONL record to the usage
        log. Never raises — metering must not break the proxy. Phase 5 replaces
        this flat log with the vault-backed quota store."""
        try:
            text = tail_bytes.decode('utf-8', 'replace')
            sys.stderr.write(f'[hermes] _record called tail={len(text)}B path={path}\n')
            usage = None
            # Fast path: non-streaming responses are a single JSON object whose
            # top-level `usage` we can read directly.
            try:
                obj = json.loads(text)
                if isinstance(obj, dict) and isinstance(obj.get('usage'), dict):
                    usage = obj['usage']
            except Exception:
                pass
            # SSE / partial path: find the LAST `"usage"` key and brace-match
            # the object that FOLLOWS it (the usage object's own braces).
            if not isinstance(usage, dict):
                idx = text.rfind('"usage"')
                brace = text.find('{', idx) if idx >= 0 else -1
                # Guard against `"usage": null` (intermediate stream chunks).
                if brace >= 0 and 'null' not in text[idx:brace]:
                    depth = 0
                    for j in range(brace, len(text)):
                        c = text[j]
                        if c == '{':
                            depth += 1
                        elif c == '}':
                            depth -= 1
                            if depth == 0:
                                try:
                                    cand = json.loads(text[brace:j + 1])
                                    if isinstance(cand, dict):
                                        usage = cand
                                except Exception:
                                    usage = None
                                break
            if not isinstance(usage, dict):
                sys.stderr.write('[hermes] _record: no usage object found in tail\n')
                return
            rec = {
                'ts': time.time(),
                'principal': principal,
                'path': path,
                'prompt_tokens': usage.get('prompt_tokens'),
                'completion_tokens': usage.get('completion_tokens'),
                'total_tokens': usage.get('total_tokens'),
            }
            log_path = self._hermes_usage_log_path()
            with open(log_path, 'a', encoding='utf-8') as f:
                f.write(json.dumps(rec, separators=(',', ':')) + '\n')
            sys.stderr.write(f'[hermes] usage principal={principal} total={rec["total_tokens"]} -> {log_path}\n')
        except Exception as e:
            sys.stderr.write(f'[hermes] _record error: {e!r}\n')

    def log_message(self, fmt, *args):
        sys.stderr.write(f'{self.address_string()} - {fmt % args}\n')


class ThreadedServer(socketserver.ThreadingMixIn, socketserver.TCPServer):
    allow_reuse_address = True
    daemon_threads = True


def _local_ip():
    """Best-effort LAN IP (not loopback) for sharing with mobile devices."""
    try:
        s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
        s.connect(('8.8.8.8', 80))
        ip = s.getsockname()[0]
        s.close()
        return ip
    except Exception:
        return None


def _ensure_local_tls(state_dir, lan_ip='', allow_selfsigned=True):
    """Provision a localhost TLS cert+key so a LOCAL run can serve HTTPS.

    A browser-TRUSTED https://localhost is what lets the SvelteKit shell at
    https://ai.cafreso.com embed this app in an <iframe> in every browser
    (Safari has no localhost mixed-content exemption). Returns
    (cert_path, key_path, trusted) or (None, None, False).

    Trust tiers, best first:
      1. mkcert — issues a cert signed by a CA installed in the OS/browser trust
         store, so the embed works with NO warning. Used if `mkcert` is on PATH.
      2. self-signed (openssl, then the `cryptography` lib) — only when
         allow_selfsigned=True (CAFRESOHQ_TLS_AUTO=1). HTTPS works after the user
         manually trusts the cert; NOT used by default because a TLS-only server
         with an untrusted cert is unreachable (worse than plain HTTP).

    Certs are cached under <state_dir>/tls and reused on later starts. A sibling
    `.mkcert` marker records whether the cached pair is browser-trusted.
    """
    try:
        tls_dir = pathlib.Path(state_dir) / 'tls'
        tls_dir.mkdir(parents=True, exist_ok=True)
    except OSError as e:
        sys.stderr.write(f'[tls] cannot create cert dir: {e}\n')
        return None, None, False

    cert = tls_dir / 'localhost.pem'
    key  = tls_dir / 'localhost-key.pem'
    marker = tls_dir / '.mkcert'

    # Reuse a previously generated pair — but never hand a cached SELF-SIGNED
    # pair to a caller that only wants trusted certs (it would flip the server
    # to TLS-only with a cert the browser rejects).
    if cert.is_file() and key.is_file() and cert.stat().st_size and key.stat().st_size:
        if marker.exists() or allow_selfsigned:
            return str(cert), str(key), marker.exists()

    hosts = ['localhost', '127.0.0.1', '::1']
    if lan_ip and lan_ip not in hosts:
        hosts.append(lan_ip)

    # ── Tier 1: mkcert (trusted) ──────────────────────────────────────────────
    mkcert = shutil.which('mkcert')
    if mkcert:
        try:
            # -install is idempotent; it may prompt for admin the very first time
            # (to add the local CA). Non-fatal if it fails — the cert still works
            # for direct use, just not the trusted embed.
            subprocess.run([mkcert, '-install'], capture_output=True, text=True, timeout=60)
            r = subprocess.run(
                [mkcert, '-cert-file', str(cert), '-key-file', str(key), *hosts],
                capture_output=True, text=True, timeout=60,
            )
            if r.returncode == 0 and cert.is_file() and key.is_file():
                marker.write_text('mkcert\n', encoding='utf-8')
                sys.stderr.write('[tls] generated a browser-trusted cert via mkcert\n')
                return str(cert), str(key), True
            sys.stderr.write(f'[tls] mkcert failed: {(r.stderr or r.stdout)[:200]}\n')
        except (subprocess.SubprocessError, OSError) as e:
            sys.stderr.write(f'[tls] mkcert error: {e}\n')

    if not allow_selfsigned:
        # Trusted-only mode (the auto default): no mkcert → stay on HTTP rather
        # than degrade to an unreachable self-signed TLS-only server.
        return None, None, False

    # ── Tier 2a: openssl self-signed ──────────────────────────────────────────
    openssl = shutil.which('openssl')
    if openssl:
        try:
            san = 'subjectAltName=' + ','.join(
                (f'IP:{h}' if (h.replace('.', '').isdigit() or ':' in h) else f'DNS:{h}')
                for h in hosts)
            r = subprocess.run(
                [openssl, 'req', '-x509', '-newkey', 'rsa:2048', '-nodes',
                 '-keyout', str(key), '-out', str(cert), '-days', '825',
                 '-subj', '/CN=localhost', '-addext', san],
                capture_output=True, text=True, timeout=60,
            )
            if r.returncode == 0 and cert.is_file() and key.is_file():
                sys.stderr.write('[tls] generated a self-signed cert via openssl '
                                 '(browser will warn until trusted)\n')
                return str(cert), str(key), False
            sys.stderr.write(f'[tls] openssl failed: {(r.stderr or r.stdout)[:200]}\n')
        except (subprocess.SubprocessError, OSError) as e:
            sys.stderr.write(f'[tls] openssl error: {e}\n')

    # ── Tier 2b: pure-Python self-signed via cryptography ─────────────────────
    try:
        from cryptography import x509
        from cryptography.x509.oid import NameOID
        from cryptography.hazmat.primitives import hashes, serialization
        from cryptography.hazmat.primitives.asymmetric import rsa
        import datetime as _dt, ipaddress as _ip

        k = rsa.generate_private_key(public_exponent=65537, key_size=2048)
        san_list = []
        for h in hosts:
            try:
                san_list.append(x509.IPAddress(_ip.ip_address(h)))
            except ValueError:
                san_list.append(x509.DNSName(h))
        name = x509.Name([x509.NameAttribute(NameOID.COMMON_NAME, 'localhost')])
        # epoch-based dates — Date.now()/utcnow are fine here (real wall clock at
        # startup), but use a fixed past start to dodge clock-skew rejections.
        not_before = _dt.datetime(2020, 1, 1)
        not_after  = _dt.datetime(2035, 1, 1)
        cert_obj = (
            x509.CertificateBuilder()
            .subject_name(name).issuer_name(name)
            .public_key(k.public_key())
            .serial_number(x509.random_serial_number())
            .not_valid_before(not_before).not_valid_after(not_after)
            .add_extension(x509.SubjectAlternativeName(san_list), critical=False)
            .sign(k, hashes.SHA256())
        )
        cert.write_bytes(cert_obj.public_bytes(serialization.Encoding.PEM))
        key.write_bytes(k.private_bytes(
            serialization.Encoding.PEM,
            serialization.PrivateFormat.TraditionalOpenSSL,
            serialization.NoEncryption()))
        sys.stderr.write('[tls] generated a self-signed cert via cryptography '
                         '(browser will warn until trusted)\n')
        return str(cert), str(key), False
    except ImportError:
        pass
    except Exception as e:
        sys.stderr.write(f'[tls] cryptography generation failed: {e}\n')

    sys.stderr.write('[tls] no cert tool available (mkcert/openssl/cryptography) '
                     '— staying on HTTP\n')
    return None, None, False


if __name__ == '__main__':
    # When packaged as a PyInstaller exe, static files live next to the exe.
    if getattr(sys, 'frozen', False):
        os.chdir(os.path.dirname(sys.executable))
    else:
        os.chdir(os.path.dirname(os.path.abspath(__file__)))

    # Is this a genuine local desktop run (vs. an OCI Fleet container)? Computed
    # here because both the TLS-auto and bind-host decisions below depend on it.
    _is_local_run = _RUNTIME_ENV == 'local' and _fleet_mode not in ('oci-fleet', 'fleet')

    # TLS — explicit cert/key (mkcert or any trusted pair) win. Otherwise, for a
    # local deployment (anything NOT behind the OCI Fleet gateway — native run
    # OR `docker run -e CAFRESOHQ_FLEET_MODE=local`), auto-provision a localhost
    # cert so the app serves HTTPS and embeds in https://ai.cafreso.com.
    #
    # CRITICAL default: auto only upgrades to HTTPS when the cert is
    # BROWSER-TRUSTED (mkcert present). A self-signed cert here would be a
    # regression, not an upgrade — the server becomes TLS-only, plain
    # http://localhost:8787 stops answering, the browser refuses the untrusted
    # cert, and the app turns unreachable (esp. in Docker, where mkcert can't
    # install a CA into the HOST browser). And http://localhost is already a
    # secure context in Chrome/Edge/Firefox, so self-signed buys nothing there.
    # CAFRESOHQ_TLS_AUTO=1 forces HTTPS incl. the self-signed fallback (power
    # users who will trust the cert manually); =0 disables auto entirely.
    _local_deploy = _fleet_mode not in ('oci-fleet', 'fleet')
    _tls_cert = os.environ.get('CAFRESOHQ_TLS_CERT', '').strip()
    _tls_key  = os.environ.get('CAFRESOHQ_TLS_KEY',  '').strip()
    _tls_trusted = bool(_tls_cert and _tls_key)   # operator-supplied → assume managed
    if not (_tls_cert and _tls_key):
        _auto = os.environ.get('CAFRESOHQ_TLS_AUTO', '').strip().lower()
        _tls_forced   = _auto in ('1', 'true', 'yes', 'on')
        _tls_disabled = _auto in ('0', 'false', 'no', 'off')
        if not _tls_disabled and (_tls_forced or _local_deploy):
            _gc, _gk, _trusted = _ensure_local_tls(
                _hq_state_dir, _local_ip() or '', allow_selfsigned=_tls_forced)
            if _gc and (_trusted or _tls_forced):
                _tls_cert, _tls_key, _tls_trusted = _gc, _gk, _trusted
    _tls_on   = bool(_tls_cert and _tls_key and
                     pathlib.Path(_tls_cert).is_file() and
                     pathlib.Path(_tls_key).is_file())
    _scheme   = 'https' if _tls_on else 'http'
    # Night-shift runner self-calls loop back over the live scheme/port.
    _night_base_url[0] = '%s://127.0.0.1:%d' % (_scheme, PORT)

    # Bind host: containers/fleet listen on all interfaces (the gateway reaches
    # them by private IP); genuine LOCAL runs default to loopback so an accidental
    # port map doesn't expose the (RCE-capable) terminal. NOTE: OCI Container
    # Instances detect as runtime_env='local' (none of the docker markers exist),
    # so gate the loopback default on the FLEET MODE too — a fleet/container
    # deploy (CAFRESOHQ_FLEET_MODE=oci-fleet) MUST bind all interfaces or the
    # gateway gets 502. Override with CAFRESOHQ_BIND.
    _bind_host = os.environ.get('CAFRESOHQ_BIND', '').strip()
    if not _bind_host:
        _bind_host = '127.0.0.1' if _is_local_run else ''
    if _is_local_run and _bind_host not in ('127.0.0.1', 'localhost') \
            and not CAFRESOHQ_API_KEY:
        print('  ⚠  bound to a non-loopback interface with NO CAFRESOHQ_API_KEY — '
              'the terminal PTY is exposed. Set CAFRESOHQ_API_KEY to lock it down.')
    with ThreadedServer((_bind_host, PORT), Handler) as httpd:
        if _tls_on:
            ctx = ssl.SSLContext(ssl.PROTOCOL_TLS_SERVER)
            ctx.load_cert_chain(_tls_cert, _tls_key)
            httpd.socket = ctx.wrap_socket(httpd.socket, server_side=True)
            print(f'  🔒 TLS enabled  cert={_tls_cert}'
                  + ('  (browser-trusted)' if _tls_trusted else '  (self-signed)'))

        lan = _local_ip()
        print(f'CafresoHQ -> {_scheme}://localhost:{PORT}/hq.html')
        if lan:
            print(f'  📱 Mobile / LAN  -> {_scheme}://{lan}:{PORT}/hq.html')
        if _tls_on and _local_deploy:
            if _tls_trusted:
                print('  ✅ Trusted HTTPS — this HQ can load embedded inside '
                      'https://ai.cafreso.com')
            else:
                print('  ⚠  Self-signed HTTPS. For the seamless embed in '
                      'ai.cafreso.com, install mkcert (https://github.com/FiloSottile/mkcert) '
                      'and restart — or open https://localhost:%d/hq.html once and '
                      'trust the cert.' % PORT)
        elif not _tls_on:
            print('  💡 Set CAFRESOHQ_TLS_CERT + CAFRESOHQ_TLS_KEY (or install mkcert) '
                  'for HTTPS — enables the ai.cafreso.com embed + iOS service worker')
        for prefix, (h, p) in ROUTES.items():
            print(f'  proxy {prefix}* -> {h}:{p}/*')
        if _cafresohq_allowed_dirs:
            print(f'  /cafresohq/stream  ELEVATED · tools={",".join(_cafresohq_allowed_tools)}'
                  f' · dirs={_cafresohq_allowed_dirs}')
        else:
            print('  /cafresohq/stream  DISABLED (set CAFRESOHQ_ALLOWED_DIRS to enable)')
        _codex_found = (_codex_bin if _codex_bin and pathlib.Path(_codex_bin).is_file() else '') \
            or shutil.which(_codex_bin or 'codex') or shutil.which('codex.cmd')
        if _codex_found and _cafresohq_allowed_dirs:
            print(f'  /codex/stream     CODEX · dirs={_cafresohq_allowed_dirs}')
        else:
            print('  /codex/stream     DISABLED (codex not found or CAFRESOHQ_ALLOWED_DIRS not set)')
        print('  /approvals/external  Claude Code PreToolUse hook bridge')
        httpd.serve_forever()
